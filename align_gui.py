"""
PowerPoint Image Grid Generator GUI

This application provides a graphical interface for creating PowerPoint presentations
with images arranged in a grid layout.

Features:
- Interactive REAL-TIME PREVIEW of the slide layout.
- Layout Mode (Grid vs Flow).
- **NEW**: Visual Crop Editor.
  - Open images directly from selected folders.
  - Draw crop regions with mouse drag.
  - Auto-detection of image aspect ratio on folder add.
  - **UPDATED**: Crop Editor now always creates a NEW region instead of editing selected one.
- Flow Alignment (Left/Center/Right).
- Fixed image size specification.
- Detailed crop settings (Alignment, Offset, Custom Gaps).
- **UPDATED**: Full editing capability for Crop Regions (Name, X, Y, W, H, Color) in the list.
- **UPDATED**: Separate controls for Crop Border (on source) and Zoom Border (on crop) thickness.
- Border shape selection (Rectangle / Rounded) for Zoom Border.
- **FIXED**: Layout calculation now accounts for border width to prevent overlap.
- Precise Gap Control.
- Save/Load configuration to/from YAML files.
"""

import os
import re
import sys
import shutil
import tempfile
import threading
from pathlib import Path
from typing import Optional, List, Tuple, Dict, Union
from dataclasses import dataclass, field

import tkinter as tk
from tkinter import ttk, filedialog, messagebox, colorchooser

import yaml
from PIL import Image, ImageTk
from pptx import Presentation
from pptx.util import Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


# --- Backend Logic ---

CM_TO_EMU = 360000
PT_TO_EMU = 12700


def cm_to_emu(cm: float) -> int:
    return int(cm * CM_TO_EMU)


def pt_to_emu(pt: float) -> int:
    return int(pt * PT_TO_EMU)


def pt_to_cm(pt: float) -> float:
    """Convert points to centimeters."""
    # 1 inch = 2.54 cm = 72 points
    return pt * (2.54 / 72.0)


@dataclass
class CropRegion:
    x: int
    y: int
    width: int
    height: int
    color: tuple[int, int, int] = (255, 0, 0)
    name: str = ""
    # Per-crop positioning
    align: str = "auto"  # 'auto', 'start', 'center', 'end'
    offset: float = 0.0  # cm, offset from alignment anchor
    gap: Optional[float] = None  # cm, overrides global main-crop gap if set


@dataclass
class GapConfig:
    value: float = 0.5
    mode: str = "cm"  # 'cm' or 'scale'

    def to_cm(self, ref_size: float) -> float:
        if self.mode == "scale":
            return ref_size * self.value
        return self.value


@dataclass
class CropDisplayConfig:
    position: str = "right"
    main_crop_gap: GapConfig = field(default_factory=lambda: GapConfig(0.15, "cm"))
    crop_crop_gap: GapConfig = field(default_factory=lambda: GapConfig(0.15, "cm"))
    crop_bottom_gap: GapConfig = field(default_factory=lambda: GapConfig(0.0, "cm"))
    size: Optional[float] = None
    scale: Optional[float] = None


@dataclass
class GridConfig:
    slide_width: float = 33.867
    slide_height: float = 19.05
    rows: int = 2
    cols: int = 3
    margin_left: float = 1.0
    margin_top: float = 1.0
    margin_right: float = 1.0
    margin_bottom: float = 1.0

    # Gap settings
    gap_h: GapConfig = field(default_factory=lambda: GapConfig(0.5, "cm"))
    gap_v: GapConfig = field(default_factory=lambda: GapConfig(0.5, "cm"))

    # Layout Logic
    layout_mode: str = "grid"  # 'grid' (aligned) or 'flow' (compact)
    flow_align: str = "left"  # 'left', 'center', 'right' (Only for flow mode)

    # Image sizing
    size_mode: str = "fit"  # 'fit' or 'fixed'
    fit_mode: str = "fit"  # 'fit', 'width', 'height'
    image_width: Optional[float] = None
    image_height: Optional[float] = None
    image_scale: float = 1.0

    arrangement: str = "row"
    crop_regions: List[CropRegion] = field(default_factory=list)
    crop_rows: Optional[List[int]] = None
    crop_cols: Optional[List[int]] = None
    crop_display: CropDisplayConfig = field(default_factory=CropDisplayConfig)

    # Border settings
    show_crop_border: bool = True
    crop_border_width: float = 1.5
    show_zoom_border: bool = True
    zoom_border_width: float = 1.5
    zoom_border_shape: str = "rectangle"  # 'rectangle' or 'rounded'

    folders: List[str] = field(default_factory=list)
    output: str = "output.pptx"


def parse_color(color_value) -> tuple[int, int, int]:
    if isinstance(color_value, (list, tuple)):
        return tuple(color_value[:3])
    elif isinstance(color_value, str):
        if color_value.startswith("#"):
            hex_color = color_value.lstrip("#")
            return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))
        else:
            parts = color_value.split(",")
            return tuple(int(p.strip()) for p in parts)
    return (255, 0, 0)


def parse_gap(data) -> GapConfig:
    if isinstance(data, (int, float)):
        return GapConfig(float(data), "cm")
    elif isinstance(data, dict):
        return GapConfig(float(data.get("value", 0.5)), data.get("mode", "cm"))
    return GapConfig(0.5, "cm")


def load_config(config_path: str) -> GridConfig:
    with open(config_path, "r", encoding="utf-8") as f:
        data = yaml.safe_load(f)

    config = GridConfig()

    if "slide" in data:
        slide = data["slide"]
        config.slide_width = slide.get("width", config.slide_width)
        config.slide_height = slide.get("height", config.slide_height)

    if "grid" in data:
        grid = data["grid"]
        config.rows = grid.get("rows", config.rows)
        config.cols = grid.get("cols", config.cols)
        config.arrangement = grid.get("arrangement", config.arrangement)
        config.layout_mode = grid.get("layout_mode", config.layout_mode)
        config.flow_align = grid.get("flow_align", config.flow_align)

    if "margin" in data:
        margin = data["margin"]
        if isinstance(margin, (int, float)):
            config.margin_left = margin
            config.margin_top = margin
            config.margin_right = margin
            config.margin_bottom = margin
        else:
            config.margin_left = margin.get("left", config.margin_left)
            config.margin_top = margin.get("top", config.margin_top)
            config.margin_right = margin.get("right", config.margin_right)
            config.margin_bottom = margin.get("bottom", config.margin_bottom)

    if "gap" in data:
        gap = data["gap"]
        if isinstance(gap, (int, float)):
            config.gap_h = GapConfig(float(gap), "cm")
            config.gap_v = GapConfig(float(gap), "cm")
        else:
            if "horizontal" in gap:
                config.gap_h = parse_gap(gap["horizontal"])
            if "vertical" in gap:
                config.gap_v = parse_gap(gap["vertical"])

    if "image" in data:
        img = data["image"]
        config.size_mode = img.get("size_mode", config.size_mode)
        config.fit_mode = img.get("fit_mode", config.fit_mode)
        config.image_scale = img.get("scale", config.image_scale)
        config.image_width = img.get("width")
        config.image_height = img.get("height")

    if "crop" in data:
        crop = data["crop"]
        if "regions" in crop:
            for i, r in enumerate(crop["regions"]):
                region = CropRegion(
                    x=r.get("x", 0),
                    y=r.get("y", 0),
                    width=r.get("width", 100),
                    height=r.get("height", 100),
                    color=parse_color(r.get("color", "#FF0000")),
                    name=r.get("name", f"crop_{i + 1}"),
                    align=r.get("align", "auto"),
                    offset=r.get("offset", 0.0),
                    gap=r.get("gap", None),
                )
                config.crop_regions.append(region)
        elif "region" in crop:
            r = crop["region"]
            region = CropRegion(
                x=r.get("x", 0),
                y=r.get("y", 0),
                width=r.get("width", 100),
                height=r.get("height", 100),
                color=parse_color(r.get("color", "#FF0000")),
                name="crop_1",
                align=r.get("align", "auto"),
                offset=r.get("offset", 0.0),
                gap=r.get("gap", None),
            )
            config.crop_regions.append(region)

        config.crop_rows = crop.get("rows")
        config.crop_cols = crop.get("cols")

        if "display" in crop:
            disp = crop["display"]
            config.crop_display.position = disp.get("position", "right")
            config.crop_display.size = disp.get("size")
            config.crop_display.scale = disp.get("scale")

            legacy_gap = disp.get("gap")
            if legacy_gap is not None:
                if isinstance(legacy_gap, (int, float)):
                    config.crop_display.main_crop_gap = GapConfig(
                        float(legacy_gap), "cm"
                    )
                    config.crop_display.crop_crop_gap = GapConfig(
                        float(legacy_gap), "cm"
                    )

            if "main_crop_gap" in disp:
                config.crop_display.main_crop_gap = parse_gap(disp["main_crop_gap"])
            if "crop_crop_gap" in disp:
                config.crop_display.crop_crop_gap = parse_gap(disp["crop_crop_gap"])
            if "crop_bottom_gap" in disp:
                config.crop_display.crop_bottom_gap = parse_gap(disp["crop_bottom_gap"])

    if "border" in data:
        border = data["border"]
        if "crop" in border:
            cb = border["crop"]
            config.show_crop_border = cb.get("show", config.show_crop_border)
            config.crop_border_width = cb.get("width", config.crop_border_width)
        if "zoom" in border:
            zb = border["zoom"]
            config.show_zoom_border = zb.get("show", config.show_zoom_border)
            config.zoom_border_width = zb.get("width", config.zoom_border_width)
            config.zoom_border_shape = zb.get("shape", config.zoom_border_shape)

    if "folders" in data:
        config.folders = data["folders"]
    config.output = data.get("output", config.output)

    return config


def extract_number_from_filename(filename: str) -> int:
    numbers = re.findall(r"\d+", filename)
    return int(numbers[0]) if numbers else 0


def get_sorted_images(folder_path: str) -> List[str]:
    supported_extensions = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp"}
    folder = Path(folder_path)
    if not folder.exists():
        return []

    image_files = [
        f
        for f in folder.iterdir()
        if f.is_file() and f.suffix.lower() in supported_extensions
    ]
    image_files.sort(key=lambda f: extract_number_from_filename(f.stem))
    return [str(f) for f in image_files]


def crop_image(image_path: str, region: CropRegion, output_path: str) -> str:
    with Image.open(image_path) as img:
        box = (region.x, region.y, region.x + region.width, region.y + region.height)
        cropped = img.crop(box)
        cropped.save(output_path)
    return output_path


def get_image_dimensions(image_path: str) -> Tuple[int, int]:
    with Image.open(image_path) as img:
        return img.size


def should_apply_crop(row: int, col: int, config: GridConfig) -> bool:
    if not config.crop_regions:
        return False
    if config.crop_rows is not None and row not in config.crop_rows:
        return False
    if config.crop_cols is not None and col not in config.crop_cols:
        return False
    return True


def calculate_size_fit_static(
    img_w: int, img_h: int, max_width: float, max_height: float, fit_mode: str = "fit"
) -> Tuple[float, float]:
    if max_height <= 0 or max_width <= 0:
        return 0, 0
    if img_h == 0:
        return 0, 0
    aspect_ratio = img_w / img_h

    if fit_mode == "width":
        width = max_width
        height = width / aspect_ratio
    elif fit_mode == "height":
        height = max_height
        width = height * aspect_ratio
    else:  # 'fit'
        cell_aspect = max_width / max_height
        if aspect_ratio > cell_aspect:
            width = max_width
            height = width / aspect_ratio
        else:
            height = max_height
            width = height * aspect_ratio

    return width, height


def calculate_image_size_fit(
    image_path: str, max_width: float, max_height: float, fit_mode: str = "fit"
) -> Tuple[float, float]:
    try:
        img_width_px, img_height_px = get_image_dimensions(image_path)
    except Exception:
        return 0, 0
    return calculate_size_fit_static(
        img_width_px, img_height_px, max_width, max_height, fit_mode
    )


def add_border_shape(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    border_color: tuple[int, int, int],
    border_width: float,
    shape_type: str = "rectangle",
):
    ms_shape_type = MSO_SHAPE.RECTANGLE
    if shape_type == "rounded":
        ms_shape_type = MSO_SHAPE.ROUNDED_RECTANGLE

    shape = slide.shapes.add_shape(
        ms_shape_type,
        cm_to_emu(left),
        cm_to_emu(top),
        cm_to_emu(width),
        cm_to_emu(height),
    )
    shape.fill.background()
    shape.line.color.rgb = RGBColor(*border_color)
    shape.line.width = pt_to_emu(border_width)
    shape.shadow.inherit = False


def add_crop_borders_to_image(
    slide,
    image_left: float,
    image_top: float,
    image_width: float,
    image_height: float,
    original_image_path: str,
    crop_regions: List[CropRegion],
    border_width: float,
):
    try:
        orig_width_px, orig_height_px = get_image_dimensions(original_image_path)
    except Exception:
        return

    scale_x = image_width / orig_width_px
    scale_y = image_height / orig_height_px

    for region in crop_regions:
        border_left = image_left + region.x * scale_x
        border_top = image_top + region.y * scale_y
        border_w = region.width * scale_x
        border_h = region.height * scale_y
        add_border_shape(
            slide,
            border_left,
            border_top,
            border_w,
            border_h,
            region.color,
            border_width,
            "rectangle",
        )  # Original image crop marks are always rects


@dataclass
class LayoutMetrics:
    main_width: float
    main_height: float
    col_widths: List[float]
    row_heights: List[float]
    crop_size: float
    crop_main_gap: float
    crop_crop_gap: float
    crop_bottom_gap: float


def calculate_grid_metrics(config: GridConfig) -> LayoutMetrics:
    total_grid_w = config.slide_width - config.margin_left - config.margin_right
    total_grid_h = config.slide_height - config.margin_top - config.margin_bottom

    est_cols = max(1, config.cols)
    est_rows = max(1, config.rows)
    est_main_w = total_grid_w / est_cols
    est_main_h = total_grid_h / est_rows

    gap_h_val = config.gap_h.to_cm(est_main_w)
    gap_v_val = config.gap_v.to_cm(est_main_h)

    avail_w_for_cells = total_grid_w - (gap_h_val * (config.cols - 1))
    avail_h_for_cells = total_grid_h - (gap_v_val * (config.rows - 1))

    has_crops = len(config.crop_regions) > 0
    expanded_cols = set()
    expanded_rows = set()
    if has_crops:
        expanded_cols = (
            set(range(config.cols))
            if config.crop_cols is None
            else set(config.crop_cols)
        )
        expanded_rows = (
            set(range(config.rows))
            if config.crop_rows is None
            else set(config.crop_rows)
        )
    num_exp_cols = len([c for c in range(config.cols) if c in expanded_cols])
    num_exp_rows = len([r for r in range(config.rows) if r in expanded_rows])

    disp = config.crop_display

    gap_mc = disp.main_crop_gap.to_cm(est_main_w)
    gap_cc = disp.crop_crop_gap.to_cm(est_main_w)
    gap_cb = disp.crop_bottom_gap.to_cm(
        est_main_w if disp.position == "right" else est_main_h
    )

    # Border offset compensation
    # If zoom border is shown, we need extra space to avoid overlap because border is drawn on the line.
    # Usually half width extends out. We add full width for safety/spacing on both sides relative to gap.
    border_offset = 0.0
    if config.show_zoom_border:
        border_offset = pt_to_cm(config.zoom_border_width)
        # Add a bit of space to gap_mc and gap_cc effectively
        gap_mc += border_offset
        gap_cc += border_offset

    main_w = 0.0
    main_h = 0.0
    crop_box_size = 0.0

    num_crops = max(1, len(config.crop_regions))

    is_fixed_size = (
        config.size_mode == "fixed"
        and config.image_width
        and config.image_width > 0
        and config.image_height
        and config.image_height > 0
    )

    if is_fixed_size:
        main_w = config.image_width
        main_h = config.image_height

        if disp.size is not None:
            crop_box_size = disp.size
        elif disp.scale is not None:
            if disp.position == "right":
                crop_box_size = main_w * disp.scale
            else:
                crop_box_size = main_h * disp.scale
        else:
            if disp.position == "right":
                crop_box_size = main_w * 0.33
            else:
                crop_box_size = main_h * 0.33

    else:
        if disp.position == "right":
            main_h = avail_h_for_cells / config.rows
            crop_k = 0.0
            crop_c = 0.0
            if disp.size is not None:
                crop_c = disp.size
            elif disp.scale is not None:
                crop_k = disp.scale
            else:
                crop_k = 0.33

            # Note: gap_mc includes border offset now
            denom = config.cols + num_exp_cols * crop_k
            numer = avail_w_for_cells - num_exp_cols * (gap_mc + crop_c)
            main_w = numer / denom
            crop_box_size = main_w * crop_k + crop_c

        else:  # bottom
            main_w = avail_w_for_cells / config.cols
            crop_k = 0.0
            crop_c = 0.0
            if disp.size is not None:
                crop_c = disp.size
            elif disp.scale is not None:
                crop_k = disp.scale
            else:
                crop_k = 0.33

            denom = config.rows + num_exp_rows * crop_k
            numer = avail_h_for_cells - num_exp_rows * (gap_mc + crop_c)
            main_h = numer / denom
            crop_box_size = main_h * crop_k + crop_c

    col_widths = []
    row_heights = []

    if disp.position == "right":
        extra = crop_box_size + gap_mc
        for c in range(config.cols):
            if c in expanded_cols:
                col_widths.append(main_w + extra)
            else:
                col_widths.append(main_w)
        # Note: Vertical space for crops (if they exceed main_h) isn't fully calculated here for Grid mode
        # Simple approximation for now
        row_heights = [main_h] * config.rows
    else:
        # Bottom placement expands row height
        extra = crop_box_size + gap_mc + gap_cb
        col_widths = [main_w] * config.cols
        for r in range(config.rows):
            if r in expanded_rows:
                row_heights.append(main_h + extra)
            else:
                row_heights.append(main_h)

    return LayoutMetrics(
        main_w, main_h, col_widths, row_heights, crop_box_size, gap_mc, gap_cc, gap_cb
    )


def create_grid_presentation(config: GridConfig) -> str:
    prs = Presentation()
    prs.slide_width = cm_to_emu(config.slide_width)
    prs.slide_height = cm_to_emu(config.slide_height)
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    if config.arrangement == "row":
        image_grid = [get_sorted_images(folder) for folder in config.folders]
    else:
        columns = [get_sorted_images(folder) for folder in config.folders]
        max_len = max(len(col) for col in columns) if columns else 0
        image_grid = [
            [col[i] if i < len(col) else None for col in columns]
            for i in range(max_len)
        ]

    temp_dir = tempfile.mkdtemp()
    num_crops = len(config.crop_regions)
    metrics = calculate_grid_metrics(config)

    # Calculate Border Offset for placement
    border_offset_cm = 0.0
    if config.show_zoom_border:
        border_offset_cm = pt_to_cm(config.zoom_border_width)

    try:
        current_y = config.margin_top

        for row_idx, row_images in enumerate(image_grid):
            if row_idx >= config.rows:
                break

            this_gap_v = config.gap_v.to_cm(metrics.main_height)
            current_x = config.margin_left
            current_row_height = (
                metrics.row_heights[row_idx]
                if row_idx < len(metrics.row_heights)
                else metrics.main_height
            )

            # --- Pre-calculate row total width for alignment if in flow mode ---
            row_total_content_width = 0.0
            valid_items = 0

            if config.layout_mode == "flow":
                for col_idx, image_path in enumerate(row_images):
                    if col_idx >= config.cols:
                        break
                    if image_path is None:
                        continue

                    has_crops = should_apply_crop(row_idx, col_idx, config)
                    orig_w, orig_h = calculate_image_size_fit(
                        image_path,
                        metrics.main_width,
                        metrics.main_height,
                        config.fit_mode,
                    )

                    item_width = orig_w
                    if has_crops and config.crop_regions:
                        num_crops = len(config.crop_regions)
                        disp = config.crop_display
                        actual_gap_mc = disp.main_crop_gap.to_cm(
                            orig_w if disp.position == "right" else orig_h
                        )
                        if config.show_zoom_border:
                            actual_gap_mc += border_offset_cm  # compensate gap

                        if disp.position == "right":
                            max_crop_ext = 0
                            for ci, r_crop in enumerate(config.crop_regions):
                                crop_w = 0
                                if disp.size is not None:
                                    if disp.position == "right":
                                        crop_w, _ = calculate_size_fit_static(
                                            100, 100, disp.size, 9999, "width"
                                        )
                                    else:
                                        crop_w, _ = calculate_size_fit_static(
                                            100, 100, 9999, disp.size, "height"
                                        )
                                elif disp.scale is not None:
                                    if disp.position == "right":
                                        tw = orig_w * disp.scale
                                        crop_w, _ = calculate_size_fit_static(
                                            100, 100, tw, 9999, "width"
                                        )
                                    else:
                                        th = orig_h * disp.scale
                                        crop_w, _ = calculate_size_fit_static(
                                            100, 100, 9999, th, "height"
                                        )
                                else:
                                    crop_w = metrics.crop_size

                                this_gap = (
                                    r_crop.gap
                                    if r_crop.gap is not None
                                    else actual_gap_mc
                                )
                                current_ext = this_gap + crop_w
                                if current_ext > max_crop_ext:
                                    max_crop_ext = current_ext
                            item_width += max_crop_ext
                        else:
                            # bottom
                            # Check if crop stack is wider than main
                            actual_gap_cc = disp.crop_crop_gap.to_cm(
                                orig_h
                            )  # bottom uses height ref approx
                            if config.show_zoom_border:
                                actual_gap_cc += border_offset_cm

                            # Estimate widths
                            c_w_sum = 0
                            for ci, r_crop in enumerate(config.crop_regions):
                                c_w = 0
                                if disp.size is not None:
                                    c_w, _ = calculate_size_fit_static(
                                        100, 100, 9999, disp.size, "height"
                                    )
                                elif disp.scale is not None:
                                    th = orig_h * disp.scale
                                    c_w, _ = calculate_size_fit_static(
                                        100, 100, 9999, th, "height"
                                    )
                                else:
                                    sw = (
                                        orig_w - actual_gap_cc * (num_crops - 1)
                                    ) / num_crops
                                    c_w, _ = calculate_size_fit_static(
                                        100, 100, sw, metrics.crop_size, "fit"
                                    )
                                c_w_sum += c_w

                            if num_crops > 1:
                                c_w_sum += actual_gap_cc * (num_crops - 1)
                            if c_w_sum > item_width:
                                item_width = c_w_sum

                    row_total_content_width += item_width
                    valid_items += 1

                if valid_items > 1:
                    gap_val = config.gap_h.to_cm(metrics.main_width)  # approx
                    row_content_width += (valid_items - 1) * gap_val

                avail_w = config.slide_width - config.margin_left - config.margin_right
                if config.flow_align == "center":
                    current_x = (
                        config.margin_left + (avail_w - row_total_content_width) / 2
                    )
                elif config.flow_align == "right":
                    current_x = config.margin_left + (avail_w - row_total_content_width)
                else:
                    current_x = config.margin_left

            # --- End Pre-calculation ---

            for col_idx, image_path in enumerate(row_images):
                if col_idx >= config.cols:
                    break

                this_gap_h = config.gap_h.to_cm(metrics.main_width)

                if image_path is None:
                    if config.layout_mode == "grid":
                        w = (
                            metrics.col_widths[col_idx]
                            if col_idx < len(metrics.col_widths)
                            else metrics.main_width
                        )
                        current_x += w + this_gap_h
                    continue

                has_crops = should_apply_crop(row_idx, col_idx, config)
                orig_w, orig_h = calculate_image_size_fit(
                    image_path, metrics.main_width, metrics.main_height, config.fit_mode
                )

                # Global dynamic gaps
                global_gap_mc = config.crop_display.main_crop_gap.to_cm(
                    orig_w if config.crop_display.position == "right" else orig_h
                )
                global_gap_cc = config.crop_display.crop_crop_gap.to_cm(
                    orig_w if config.crop_display.position == "right" else orig_h
                )
                global_gap_cb = config.crop_display.crop_bottom_gap.to_cm(
                    orig_w if config.crop_display.position == "right" else orig_h
                )

                # Compensate for border
                if config.show_zoom_border:
                    global_gap_mc += border_offset_cm
                    global_gap_cc += border_offset_cm

                if config.layout_mode == "flow":
                    final_main_left = current_x
                    final_main_top = current_y + (metrics.main_height - orig_h) / 2
                else:
                    final_main_left = current_x + (metrics.main_width - orig_w) / 2
                    final_main_top = current_y + (metrics.main_height - orig_h) / 2

                pic = slide.shapes.add_picture(
                    image_path,
                    cm_to_emu(final_main_left),
                    cm_to_emu(final_main_top),
                    cm_to_emu(orig_w),
                    cm_to_emu(orig_h),
                )
                pic.shadow.inherit = False

                if has_crops and config.show_crop_border:
                    add_crop_borders_to_image(
                        slide,
                        final_main_left,
                        final_main_top,
                        orig_w,
                        orig_h,
                        image_path,
                        config.crop_regions,
                        config.crop_border_width,
                    )

                content_right_edge = final_main_left + orig_w
                content_bottom_edge = final_main_top + orig_h

                if has_crops and num_crops > 0:
                    disp = config.crop_display

                    for crop_idx, region in enumerate(config.crop_regions):
                        crop_filename = f"crop_{row_idx}_{col_idx}_{crop_idx}.png"
                        crop_path = os.path.join(temp_dir, crop_filename)
                        try:
                            crop_image(image_path, region, crop_path)
                        except Exception:
                            continue

                        # Resolve Size
                        if disp.size is not None:
                            if disp.position == "right":
                                cw, ch = calculate_image_size_fit(
                                    crop_path, disp.size, 9999, "width"
                                )
                            else:
                                cw, ch = calculate_image_size_fit(
                                    crop_path, 9999, disp.size, "height"
                                )
                        elif disp.scale is not None:
                            if disp.position == "right":
                                target_w = orig_w * disp.scale
                                cw, ch = calculate_image_size_fit(
                                    crop_path, target_w, 9999, "width"
                                )
                            else:
                                target_h = orig_h * disp.scale
                                cw, ch = calculate_image_size_fit(
                                    crop_path, 9999, target_h, "height"
                                )
                        else:
                            # Fallback fit
                            if disp.position == "right":
                                single_h = (
                                    orig_h - global_gap_cc * (num_crops - 1)
                                ) / num_crops
                                cw, ch = calculate_image_size_fit(
                                    crop_path, metrics.crop_size, single_h, "fit"
                                )
                            else:
                                single_w = (
                                    orig_w - global_gap_cc * (num_crops - 1)
                                ) / num_crops
                                cw, ch = calculate_image_size_fit(
                                    crop_path, single_w, metrics.crop_size, "fit"
                                )

                        # Resolve Gap
                        this_gap_mc = (
                            region.gap if region.gap is not None else global_gap_mc
                        )
                        if region.gap is not None and config.show_zoom_border:
                            this_gap_mc += border_offset_cm

                        # Resolve Alignment Position
                        if disp.position == "right":
                            c_left = final_main_left + orig_w + this_gap_mc

                            # Y-Alignment
                            if region.align == "start":
                                c_top = final_main_top + region.offset
                            elif region.align == "center":
                                c_top = (
                                    final_main_top + (orig_h - ch) / 2 + region.offset
                                )
                            elif region.align == "end":
                                c_top = final_main_top + orig_h - ch + region.offset
                            else:  # auto
                                if disp.scale is not None or disp.size is not None:
                                    # Stack logic
                                    c_top = final_main_top + crop_idx * (
                                        ch + global_gap_cc
                                    )
                                else:
                                    # Fit logic (Pin ends)
                                    if crop_idx == 0:
                                        c_top = final_main_top
                                    elif num_crops > 1 and crop_idx == num_crops - 1:
                                        c_top = (final_main_top + orig_h) - ch
                                    else:
                                        single_h = (
                                            orig_h - global_gap_cc * (num_crops - 1)
                                        ) / num_crops
                                        slot_top = final_main_top + crop_idx * (
                                            single_h + global_gap_cc
                                        )
                                        c_top = slot_top + (single_h - ch) / 2

                            content_right_edge = max(content_right_edge, c_left + cw)
                            content_bottom_edge = max(content_bottom_edge, c_top + ch)

                        else:  # bottom
                            c_top = final_main_top + orig_h + this_gap_mc

                            # X-Alignment
                            if region.align == "start":
                                c_left = final_main_left + region.offset
                            elif region.align == "center":
                                c_left = (
                                    final_main_left + (orig_w - cw) / 2 + region.offset
                                )
                            elif region.align == "end":
                                c_left = final_main_left + orig_w - cw + region.offset
                            else:  # auto
                                if disp.scale is not None or disp.size is not None:
                                    c_left = final_main_left + crop_idx * (
                                        cw + global_gap_cc
                                    )
                                else:
                                    if crop_idx == 0:
                                        c_left = final_main_left
                                    elif num_crops > 1 and crop_idx == num_crops - 1:
                                        c_left = (final_main_left + orig_w) - cw
                                    else:
                                        single_w = (
                                            orig_w - global_gap_cc * (num_crops - 1)
                                        ) / num_crops
                                        slot_left = final_main_left + crop_idx * (
                                            single_w + global_gap_cc
                                        )
                                        c_left = slot_left + (single_w - cw) / 2

                            content_bottom_edge = max(content_bottom_edge, c_top + ch)
                            content_right_edge = max(content_right_edge, c_left + cw)

                        pic_crop = slide.shapes.add_picture(
                            crop_path,
                            cm_to_emu(c_left),
                            cm_to_emu(c_top),
                            cm_to_emu(cw),
                            cm_to_emu(ch),
                        )
                        pic_crop.shadow.inherit = False

                        if config.show_zoom_border:
                            add_border_shape(
                                slide,
                                c_left,
                                c_top,
                                cw,
                                ch,
                                region.color,
                                config.zoom_border_width,
                                config.zoom_border_shape,
                            )

                # Apply Crop-Bottom Gap to content extent if crops exist
                if has_crops and num_crops > 0:
                    content_bottom_edge += global_gap_cb

                if config.layout_mode == "flow":
                    width_used = content_right_edge - current_x
                    current_x += width_used + this_gap_h
                else:
                    w = (
                        metrics.col_widths[col_idx]
                        if col_idx < len(metrics.col_widths)
                        else metrics.main_width
                    )
                    current_x += w + this_gap_h

            current_y += current_row_height + this_gap_v

        prs.save(config.output)
    except Exception as e:
        raise e
    finally:
        shutil.rmtree(temp_dir)
    return config.output


# --- GUI Application ---


class CropEditor(tk.Toplevel):
    def __init__(self, parent, image_path, callback):
        super().__init__(parent)
        self.title("Crop Editor")
        self.geometry("900x700")
        self.callback = callback

        self.image_path = image_path
        self.orig_img = Image.open(image_path)
        self.orig_w, self.orig_h = self.orig_img.size

        # UI
        self.canvas_frame = ttk.Frame(self)
        self.canvas_frame.pack(fill=tk.BOTH, expand=True)

        self.canvas = tk.Canvas(self.canvas_frame, bg="gray")
        self.canvas.pack(fill=tk.BOTH, expand=True)

        self.bottom_frame = ttk.Frame(self, padding=5)
        self.bottom_frame.pack(fill=tk.X)

        ttk.Label(self.bottom_frame, text="Name:").pack(side=tk.LEFT)
        self.var_name = tk.StringVar(value="Region")
        ttk.Entry(self.bottom_frame, textvariable=self.var_name, width=10).pack(
            side=tk.LEFT, padx=5
        )

        ttk.Button(self.bottom_frame, text="保存 (Save)", command=self.on_save).pack(
            side=tk.RIGHT, padx=5
        )
        ttk.Button(self.bottom_frame, text="キャンセル", command=self.destroy).pack(
            side=tk.RIGHT
        )

        # State
        self.rect_id = None
        self.start_x = None
        self.start_y = None
        self.cur_rect = None  # (x, y, w, h) in original coords

        # Resize image to fit canvas
        self.display_scale = 1.0
        self.tk_img = None

        self.bind("<Configure>", self.on_resize_window)
        self.canvas.bind("<ButtonPress-1>", self.on_press)
        self.canvas.bind("<B1-Motion>", self.on_drag)
        self.canvas.bind("<ButtonRelease-1>", self.on_release)

        # Initial draw
        self.after(100, self.redraw_image)

    def on_resize_window(self, event):
        if event.widget == self:
            self.redraw_image()

    def redraw_image(self):
        c_w = self.canvas.winfo_width()
        c_h = self.canvas.winfo_height()
        if c_w < 50 or c_h < 50:
            return

        scale_w = c_w / self.orig_w
        scale_h = c_h / self.orig_h
        self.display_scale = min(scale_w, scale_h) * 0.9  # margin

        new_w = int(self.orig_w * self.display_scale)
        new_h = int(self.orig_h * self.display_scale)

        resized = self.orig_img.resize((new_w, new_h), Image.Resampling.LANCZOS)
        self.tk_img = ImageTk.PhotoImage(resized)

        self.canvas.delete("all")
        # Center image
        self.off_x = (c_w - new_w) // 2
        self.off_y = (c_h - new_h) // 2

        self.canvas.create_image(
            self.off_x, self.off_y, anchor=tk.NW, image=self.tk_img
        )

        # Redraw existing rect if present
        if self.cur_rect:
            x, y, w, h = self.cur_rect
            sx = x * self.display_scale + self.off_x
            sy = y * self.display_scale + self.off_y
            ex = (x + w) * self.display_scale + self.off_x
            ey = (y + h) * self.display_scale + self.off_y
            self.rect_id = self.canvas.create_rectangle(
                sx, sy, ex, ey, outline="red", width=2
            )

    def on_press(self, event):
        self.start_x = event.x
        self.start_y = event.y
        if self.rect_id:
            self.canvas.delete(self.rect_id)
        self.rect_id = self.canvas.create_rectangle(
            self.start_x,
            self.start_y,
            self.start_x,
            self.start_y,
            outline="red",
            width=2,
        )

    def on_drag(self, event):
        self.canvas.coords(self.rect_id, self.start_x, self.start_y, event.x, event.y)

    def on_release(self, event):
        end_x, end_y = event.x, event.y

        # Normalize to image coords
        x1 = min(self.start_x, end_x) - self.off_x
        y1 = min(self.start_y, end_y) - self.off_y
        x2 = max(self.start_x, end_x) - self.off_x
        y2 = max(self.start_y, end_y) - self.off_y

        # Convert to original scale
        ox1 = max(0, int(x1 / self.display_scale))
        oy1 = max(0, int(y1 / self.display_scale))
        ox2 = min(self.orig_w, int(x2 / self.display_scale))
        oy2 = min(self.orig_h, int(y2 / self.display_scale))

        w = ox2 - ox1
        h = oy2 - oy1

        if w > 0 and h > 0:
            self.cur_rect = (ox1, oy1, w, h)

    def on_save(self):
        if self.cur_rect:
            self.callback(
                self.cur_rect[0],
                self.cur_rect[1],
                self.cur_rect[2],
                self.cur_rect[3],
                self.var_name.get(),
            )
        self.destroy()


class ImageGridApp:
    def __init__(self, root):
        self.root = root
        self.root.title("PowerPoint Grid Generator GUI")
        self.root.geometry("1400x950")

        # Variables
        self.output_path = tk.StringVar(value="output.pptx")
        self.rows = tk.IntVar(value=3)
        self.cols = tk.IntVar(value=3)
        self.arrangement = tk.StringVar(value="row")
        self.layout_mode = tk.StringVar(value="flow")
        self.flow_align = tk.StringVar(value="left")

        self.slide_w = tk.DoubleVar(value=33.867)
        self.slide_h = tk.DoubleVar(value=19.05)
        self.margin_l = tk.DoubleVar(value=1.0)
        self.margin_t = tk.DoubleVar(value=1.0)
        self.margin_r = tk.DoubleVar(value=1.0)
        self.margin_b = tk.DoubleVar(value=1.0)

        self.gap_h_val = tk.DoubleVar(value=0.5)
        self.gap_h_mode = tk.StringVar(value="cm")
        self.gap_v_val = tk.DoubleVar(value=0.5)
        self.gap_v_mode = tk.StringVar(value="cm")
        self.gap_mc_val = tk.DoubleVar(value=0.15)
        self.gap_mc_mode = tk.StringVar(value="cm")
        self.gap_cc_val = tk.DoubleVar(value=0.15)
        self.gap_cc_mode = tk.StringVar(value="cm")
        self.gap_cb_val = tk.DoubleVar(value=0.0)
        self.gap_cb_mode = tk.StringVar(value="cm")

        self.image_size_mode = tk.StringVar(value="fit")
        self.image_fit_mode = tk.StringVar(value="fit")
        self.image_w = tk.DoubleVar(value=10.0)
        self.image_h = tk.DoubleVar(value=7.5)
        self.crop_pos = tk.StringVar(value="right")

        self.crop_size_mode = tk.StringVar(value="scale")
        self.crop_size_val = tk.DoubleVar(value=0.0)
        self.crop_scale_val = tk.DoubleVar(value=0.4)

        self.show_crop_border = tk.BooleanVar(value=True)
        self.crop_border_w = tk.DoubleVar(value=1.5)
        self.zoom_border_shape = tk.StringVar(value="rectangle")  # New
        self.show_zoom_border = tk.BooleanVar(value=True)
        self.zoom_border_w = tk.DoubleVar(value=1.5)

        self.dummy_ratio_w = tk.DoubleVar(value=1.0)
        self.dummy_ratio_h = tk.DoubleVar(value=1.0)

        self.folders = []
        self.crop_regions = []
        self.crop_rows_filter = tk.StringVar(value="")
        self.crop_cols_filter = tk.StringVar(value="")

        # Selected region editing vars
        self.sel_idx = None
        self.r_name = tk.StringVar()
        self.r_x = tk.IntVar()
        self.r_y = tk.IntVar()
        self.r_w = tk.IntVar()
        self.r_h = tk.IntVar()
        self.r_color = (255, 0, 0)  # internal storage
        self.r_align = tk.StringVar(value="auto")
        self.r_offset = tk.DoubleVar(value=0.0)
        self.r_gap = tk.StringVar(value="")

        self.create_widgets()
        self.add_preview_tracers()
        self.update_preview()

    def create_widgets(self):
        self.paned = ttk.PanedWindow(self.root, orient=tk.HORIZONTAL)
        self.paned.pack(fill=tk.BOTH, expand=True)
        self.left_frame = ttk.Frame(self.paned)
        self.paned.add(self.left_frame, weight=1)
        self.right_frame = ttk.Frame(self.paned)
        self.paned.add(self.right_frame, weight=2)

        top_frame = ttk.Frame(self.left_frame, padding=5)
        top_frame.pack(fill=tk.X)
        ttk.Button(top_frame, text="設定読込", command=self.load_config_gui).pack(
            side=tk.LEFT, padx=2
        )
        ttk.Button(top_frame, text="設定保存", command=self.save_config).pack(
            side=tk.LEFT, padx=2
        )
        ttk.Button(
            top_frame, text="PPTX生成", command=self.generate, style="Accent.TButton"
        ).pack(side=tk.RIGHT, padx=2)

        self.notebook = ttk.Notebook(self.left_frame)
        self.notebook.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
        self.tab_basic = ttk.Frame(self.notebook, padding=10)
        self.tab_layout = ttk.Frame(self.notebook, padding=10)
        self.tab_crop = ttk.Frame(self.notebook, padding=10)
        self.tab_style = ttk.Frame(self.notebook, padding=10)
        self.notebook.add(self.tab_basic, text="基本・フォルダ")
        self.notebook.add(self.tab_layout, text="レイアウト")
        self.notebook.add(self.tab_crop, text="クロップ設定")
        self.notebook.add(self.tab_style, text="装飾")
        self.setup_tab_basic()
        self.setup_tab_layout()
        self.setup_tab_crop()
        self.setup_tab_style()

        f_prev_set = ttk.Frame(self.right_frame, padding=5)
        f_prev_set.pack(fill=tk.X)
        ttk.Label(f_prev_set, text="ダミー画像比率 (W:H)").pack(side=tk.LEFT, padx=5)
        ttk.Entry(f_prev_set, textvariable=self.dummy_ratio_w, width=5).pack(
            side=tk.LEFT
        )
        ttk.Label(f_prev_set, text=":").pack(side=tk.LEFT)
        ttk.Entry(f_prev_set, textvariable=self.dummy_ratio_h, width=5).pack(
            side=tk.LEFT
        )

        ttk.Label(
            self.right_frame,
            text="リアルタイムプレビュー",
            font=("Helvetica", 12, "bold"),
        ).pack(pady=5)
        canvas_frame = ttk.Frame(self.right_frame)
        canvas_frame.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
        self.preview_canvas = tk.Canvas(canvas_frame, bg="#e0e0e0")
        self.preview_canvas.pack(fill=tk.BOTH, expand=True)
        self.preview_canvas.bind("<Configure>", lambda e: self.update_preview())

    def add_preview_tracers(self):
        vars_to_trace = [
            self.rows,
            self.cols,
            self.arrangement,
            self.layout_mode,
            self.flow_align,
            self.slide_w,
            self.slide_h,
            self.margin_l,
            self.margin_t,
            self.margin_r,
            self.margin_b,
            self.gap_h_val,
            self.gap_h_mode,
            self.gap_v_val,
            self.gap_v_mode,
            self.image_size_mode,
            self.image_fit_mode,
            self.image_w,
            self.image_h,
            self.crop_pos,
            self.gap_mc_val,
            self.gap_mc_mode,
            self.gap_cc_val,
            self.gap_cc_mode,
            self.gap_cb_val,
            self.gap_cb_mode,
            self.crop_size_mode,
            self.crop_size_val,
            self.crop_scale_val,
            self.crop_rows_filter,
            self.crop_cols_filter,
            self.dummy_ratio_w,
            self.dummy_ratio_h,
            # Add crop region edit variables to trigger preview update
            self.r_name,
            self.r_x,
            self.r_y,
            self.r_w,
            self.r_h,
            self.r_align,
            self.r_offset,
            self.r_gap,
            self.zoom_border_shape,
            self.show_zoom_border,
            self.zoom_border_w,
            self.show_crop_border,
            self.crop_border_w,
        ]
        for v in vars_to_trace:
            v.trace_add("write", lambda *args: self.schedule_preview())

    def schedule_preview(self):
        if hasattr(self, "_after_id"):
            self.root.after_cancel(self._after_id)
        self._after_id = self.root.after(100, self.update_preview)

    def setup_tab_basic(self):
        f_out = ttk.LabelFrame(self.tab_basic, text="出力ファイル", padding=5)
        f_out.pack(fill=tk.X, pady=5)
        ttk.Entry(f_out, textvariable=self.output_path).pack(
            side=tk.LEFT, fill=tk.X, expand=True, padx=5
        )
        ttk.Button(f_out, text="参照", command=self.browse_output).pack(side=tk.RIGHT)

        f_folders = ttk.LabelFrame(self.tab_basic, text="入力画像フォルダ", padding=5)
        f_folders.pack(fill=tk.BOTH, expand=True, pady=5)
        btn_frame = ttk.Frame(f_folders)
        btn_frame.pack(fill=tk.X, pady=2)
        ttk.Button(btn_frame, text="追加", command=self.add_folder).pack(
            side=tk.LEFT, padx=2
        )
        ttk.Button(btn_frame, text="削除", command=self.remove_folder).pack(
            side=tk.LEFT, padx=2
        )
        ttk.Button(btn_frame, text="クリア", command=self.clear_folders).pack(
            side=tk.LEFT, padx=2
        )
        self.folder_listbox = tk.Listbox(f_folders, selectmode=tk.EXTENDED)
        self.folder_listbox.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)

        f_grid = ttk.LabelFrame(self.tab_basic, text="グリッド構成", padding=5)
        f_grid.pack(fill=tk.X, pady=5)
        ttk.Label(f_grid, text="Rows:").grid(row=0, column=0, padx=5)
        ttk.Entry(f_grid, textvariable=self.rows, width=5).grid(row=0, column=1)
        ttk.Label(f_grid, text="Cols:").grid(row=0, column=2, padx=5)
        ttk.Entry(f_grid, textvariable=self.cols, width=5).grid(row=0, column=3)
        ttk.Label(f_grid, text="並び:").grid(row=1, column=0, padx=5)
        ttk.Radiobutton(
            f_grid, text="行(Row)", variable=self.arrangement, value="row"
        ).grid(row=1, column=1)
        ttk.Radiobutton(
            f_grid, text="列(Col)", variable=self.arrangement, value="col"
        ).grid(row=1, column=2)

    def setup_tab_layout(self):
        f_mode = ttk.LabelFrame(self.tab_layout, text="レイアウトモード", padding=5)
        f_mode.pack(fill=tk.X, pady=5)
        ttk.Radiobutton(
            f_mode, text="Flow (詰める)", variable=self.layout_mode, value="flow"
        ).pack(side=tk.LEFT, padx=5)
        ttk.Radiobutton(
            f_mode, text="Grid (整列)", variable=self.layout_mode, value="grid"
        ).pack(side=tk.LEFT, padx=5)

        # New Flow Align
        ttk.Label(f_mode, text="| Flow Align:").pack(side=tk.LEFT, padx=5)
        ttk.Combobox(
            f_mode,
            textvariable=self.flow_align,
            values=["left", "center", "right"],
            width=8,
        ).pack(side=tk.LEFT)

        f_slide = ttk.LabelFrame(self.tab_layout, text="スライドサイズ", padding=5)
        f_slide.pack(fill=tk.X, pady=5)
        ttk.Label(f_slide, text="W:").pack(side=tk.LEFT)
        ttk.Entry(f_slide, textvariable=self.slide_w, width=6).pack(
            side=tk.LEFT, padx=5
        )
        ttk.Label(f_slide, text="H:").pack(side=tk.LEFT)
        ttk.Entry(f_slide, textvariable=self.slide_h, width=6).pack(
            side=tk.LEFT, padx=5
        )

        f_mg = ttk.LabelFrame(self.tab_layout, text="余白(cm)", padding=5)
        f_mg.pack(fill=tk.X, pady=5)
        ttk.Label(f_mg, text="L:").pack(side=tk.LEFT)
        ttk.Entry(f_mg, textvariable=self.margin_l, width=4).pack(side=tk.LEFT)
        ttk.Label(f_mg, text="T:").pack(side=tk.LEFT)
        ttk.Entry(f_mg, textvariable=self.margin_t, width=4).pack(side=tk.LEFT)
        ttk.Label(f_mg, text="R:").pack(side=tk.LEFT)
        ttk.Entry(f_mg, textvariable=self.margin_r, width=4).pack(side=tk.LEFT)
        ttk.Label(f_mg, text="B:").pack(side=tk.LEFT)
        ttk.Entry(f_mg, textvariable=self.margin_b, width=4).pack(side=tk.LEFT)

        f_gap = ttk.LabelFrame(self.tab_layout, text="グリッド間隔", padding=5)
        f_gap.pack(fill=tk.X, pady=5)
        r = 0
        ttk.Label(f_gap, text="横(H):").grid(row=r, column=0, sticky=tk.E)
        ttk.Entry(f_gap, textvariable=self.gap_h_val, width=5).grid(row=r, column=1)
        ttk.Radiobutton(f_gap, text="cm", variable=self.gap_h_mode, value="cm").grid(
            row=r, column=2
        )
        ttk.Radiobutton(
            f_gap, text="Scale", variable=self.gap_h_mode, value="scale"
        ).grid(row=r, column=3)
        r = 1
        ttk.Label(f_gap, text="縦(V):").grid(row=r, column=0, sticky=tk.E)
        ttk.Entry(f_gap, textvariable=self.gap_v_val, width=5).grid(row=r, column=1)
        ttk.Radiobutton(f_gap, text="cm", variable=self.gap_v_mode, value="cm").grid(
            row=r, column=2
        )
        ttk.Radiobutton(
            f_gap, text="Scale", variable=self.gap_v_mode, value="scale"
        ).grid(row=r, column=3)

        f_img = ttk.LabelFrame(self.tab_layout, text="画像サイズ設定", padding=5)
        f_img.pack(fill=tk.X, pady=5)
        f_size_m = ttk.Frame(f_img)
        f_size_m.pack(fill=tk.X)
        ttk.Radiobutton(
            f_size_m, text="Fit Mode", variable=self.image_size_mode, value="fit"
        ).pack(side=tk.LEFT)
        ttk.Radiobutton(
            f_size_m, text="Fixed Size", variable=self.image_size_mode, value="fixed"
        ).pack(side=tk.LEFT)
        f_fit = ttk.Frame(f_img)
        f_fit.pack(fill=tk.X, pady=2)
        ttk.Label(f_fit, text="基準(Fit):").pack(side=tk.LEFT)
        ttk.Radiobutton(
            f_fit, text="Fit", variable=self.image_fit_mode, value="fit"
        ).pack(side=tk.LEFT)
        ttk.Radiobutton(
            f_fit, text="Width", variable=self.image_fit_mode, value="width"
        ).pack(side=tk.LEFT)
        ttk.Radiobutton(
            f_fit, text="Height", variable=self.image_fit_mode, value="height"
        ).pack(side=tk.LEFT)
        f_fixed = ttk.Frame(f_img)
        f_fixed.pack(fill=tk.X, pady=2)
        ttk.Label(f_fixed, text="W:").pack(side=tk.LEFT)
        ttk.Entry(f_fixed, textvariable=self.image_w, width=5).pack(side=tk.LEFT)
        ttk.Label(f_fixed, text="H:").pack(side=tk.LEFT)
        ttk.Entry(f_fixed, textvariable=self.image_h, width=5).pack(side=tk.LEFT)

    def setup_tab_crop(self):
        # 1. List
        f_reg = ttk.LabelFrame(self.tab_crop, text="領域リスト", padding=5)
        f_reg.pack(fill=tk.X, pady=5)
        r_btn_frame = ttk.Frame(f_reg)
        r_btn_frame.pack(fill=tk.X)
        ttk.Button(
            r_btn_frame, text="画像から指定 (Editor)", command=self.open_crop_editor
        ).pack(side=tk.LEFT, padx=2)
        ttk.Button(
            r_btn_frame, text="追加 (数値)", command=self.add_region_dialog
        ).pack(side=tk.LEFT, padx=2)
        ttk.Button(r_btn_frame, text="削除", command=self.remove_region).pack(
            side=tk.LEFT, padx=2
        )
        self.region_tree = ttk.Treeview(
            f_reg, columns=("name", "xywh", "align"), show="headings", height=4
        )
        self.region_tree.heading("name", text="Name")
        self.region_tree.heading("xywh", text="Coord")
        self.region_tree.heading("align", text="Align")
        self.region_tree.pack(fill=tk.BOTH, expand=True)
        self.region_tree.bind("<<TreeviewSelect>>", self.on_region_select)

        # 2. Selected Detail
        f_detail = ttk.LabelFrame(
            self.tab_crop, text="選択した領域の詳細設定", padding=5
        )
        f_detail.pack(fill=tk.X, pady=5)

        # Row 1: Name, Color
        f_r1 = ttk.Frame(f_detail)
        f_r1.pack(fill=tk.X)
        ttk.Label(f_r1, text="Name:").pack(side=tk.LEFT)
        ttk.Entry(f_r1, textvariable=self.r_name, width=10).pack(side=tk.LEFT, padx=5)
        self.btn_r_color = tk.Button(
            f_r1, text="Color", width=5, command=lambda: self.pick_region_color()
        )
        self.btn_r_color.pack(side=tk.LEFT, padx=5)

        # Row 2: Coords
        f_r2 = ttk.Frame(f_detail)
        f_r2.pack(fill=tk.X)
        ttk.Label(f_r2, text="X:").pack(side=tk.LEFT)
        ttk.Entry(f_r2, textvariable=self.r_x, width=5).pack(side=tk.LEFT)
        ttk.Label(f_r2, text="Y:").pack(side=tk.LEFT)
        ttk.Entry(f_r2, textvariable=self.r_y, width=5).pack(side=tk.LEFT)
        ttk.Label(f_r2, text="W:").pack(side=tk.LEFT)
        ttk.Entry(f_r2, textvariable=self.r_w, width=5).pack(side=tk.LEFT)
        ttk.Label(f_r2, text="H:").pack(side=tk.LEFT)
        ttk.Entry(f_r2, textvariable=self.r_h, width=5).pack(side=tk.LEFT)

        # Row 3: Position props
        f_r3 = ttk.Frame(f_detail)
        f_r3.pack(fill=tk.X)
        ttk.Label(f_r3, text="Align:").pack(side=tk.LEFT)
        ttk.Combobox(
            f_r3,
            textvariable=self.r_align,
            values=["auto", "start", "center", "end"],
            width=7,
        ).pack(side=tk.LEFT)
        ttk.Label(f_r3, text="Offset:").pack(side=tk.LEFT)
        ttk.Entry(f_r3, textvariable=self.r_offset, width=5).pack(side=tk.LEFT)
        ttk.Label(f_r3, text="Gap:").pack(side=tk.LEFT)
        ttk.Entry(f_r3, textvariable=self.r_gap, width=5).pack(side=tk.LEFT)

        # Row 4: Action
        f_r4 = ttk.Frame(f_detail)
        f_r4.pack(fill=tk.X, pady=5)
        ttk.Button(f_r4, text="更新 (Update)", command=self.update_region_detail).pack(
            anchor=tk.E
        )

        # 3. Global Settings
        f_glob = ttk.LabelFrame(self.tab_crop, text="全体・配置設定", padding=5)
        f_glob.pack(fill=tk.X, pady=5)

        ttk.Label(f_glob, text="配置方向:").grid(row=0, column=0, sticky=tk.E)
        f_pos = ttk.Frame(f_glob)
        f_pos.grid(row=0, column=1, columnspan=3, sticky=tk.W)
        ttk.Radiobutton(
            f_pos, text="Right", variable=self.crop_pos, value="right"
        ).pack(side=tk.LEFT)
        ttk.Radiobutton(
            f_pos, text="Bottom", variable=self.crop_pos, value="bottom"
        ).pack(side=tk.LEFT)

        ttk.Label(f_glob, text="Main-Crop:").grid(row=1, column=0, sticky=tk.E)
        ttk.Entry(f_glob, textvariable=self.gap_mc_val, width=5).grid(row=1, column=1)
        ttk.Radiobutton(f_glob, text="cm", variable=self.gap_mc_mode, value="cm").grid(
            row=1, column=2
        )
        ttk.Radiobutton(
            f_glob, text="Scale", variable=self.gap_mc_mode, value="scale"
        ).grid(row=1, column=3)

        ttk.Label(f_glob, text="Crop-Crop:").grid(row=2, column=0, sticky=tk.E)
        ttk.Entry(f_glob, textvariable=self.gap_cc_val, width=5).grid(row=2, column=1)
        ttk.Radiobutton(f_glob, text="cm", variable=self.gap_cc_mode, value="cm").grid(
            row=2, column=2
        )
        ttk.Radiobutton(
            f_glob, text="Scale", variable=self.gap_cc_mode, value="scale"
        ).grid(row=2, column=3)

        ttk.Label(f_glob, text="Crop-Bottom:").grid(row=3, column=0, sticky=tk.E)
        ttk.Entry(f_glob, textvariable=self.gap_cb_val, width=5).grid(row=3, column=1)
        ttk.Radiobutton(f_glob, text="cm", variable=self.gap_cb_mode, value="cm").grid(
            row=3, column=2
        )
        ttk.Radiobutton(
            f_glob, text="Scale", variable=self.gap_cb_mode, value="scale"
        ).grid(row=3, column=3)

        ttk.Label(f_glob, text="サイズ:").grid(row=4, column=0, sticky=tk.E)
        f_sz = ttk.Frame(f_glob)
        f_sz.grid(row=4, column=1, columnspan=3, sticky=tk.W)
        ttk.Radiobutton(
            f_sz, text="倍率(Scale)", variable=self.crop_size_mode, value="scale"
        ).pack(side=tk.LEFT)
        ttk.Entry(f_sz, textvariable=self.crop_scale_val, width=5).pack(side=tk.LEFT)
        ttk.Radiobutton(
            f_sz, text="固定(cm)", variable=self.crop_size_mode, value="size"
        ).pack(side=tk.LEFT, padx=(10, 0))
        ttk.Entry(f_sz, textvariable=self.crop_size_val, width=5).pack(side=tk.LEFT)

        f_filter = ttk.LabelFrame(self.tab_crop, text="適用対象 (空欄=全て)", padding=5)
        f_filter.pack(fill=tk.X, pady=5)
        ttk.Label(f_filter, text="行:").pack(side=tk.LEFT)
        ttk.Entry(f_filter, textvariable=self.crop_rows_filter, width=10).pack(
            side=tk.LEFT, padx=5
        )
        ttk.Label(f_filter, text="列:").pack(side=tk.LEFT)
        ttk.Entry(f_filter, textvariable=self.crop_cols_filter, width=10).pack(
            side=tk.LEFT, padx=5
        )

    def pick_region_color(self):
        c = colorchooser.askcolor(color=self.r_color)
        if c[0]:
            self.r_color = tuple(map(int, c[0]))
            self.btn_r_color.config(bg=c[1])

    def on_region_select(self, event):
        sel = self.region_tree.selection()
        if not sel:
            return
        self.sel_idx = self.region_tree.index(sel[0])
        r = self.crop_regions[self.sel_idx]
        self.r_name.set(r.name)
        self.r_x.set(r.x)
        self.r_y.set(r.y)
        self.r_w.set(r.width)
        self.r_h.set(r.height)
        self.r_color = r.color
        self.btn_r_color.config(bg=f"#{r.color[0]:02x}{r.color[1]:02x}{r.color[2]:02x}")
        self.r_align.set(r.align)
        self.r_offset.set(r.offset)
        self.r_gap.set(str(r.gap) if r.gap is not None else "")

    def update_region_detail(self):
        if self.sel_idx is None or self.sel_idx >= len(self.crop_regions):
            return
        r = self.crop_regions[self.sel_idx]
        r.name = self.r_name.get()
        try:
            r.x = self.r_x.get()
        except:
            pass
        try:
            r.y = self.r_y.get()
        except:
            pass
        try:
            r.width = self.r_w.get()
        except:
            pass
        try:
            r.height = self.r_h.get()
        except:
            pass
        r.color = self.r_color
        r.align = self.r_align.get()
        try:
            r.offset = self.r_offset.get()
        except Exception:
            pass
        try:
            g_str = self.r_gap.get().strip()
            r.gap = float(g_str) if g_str else None
        except Exception:
            r.gap = None
        self.update_region_list()
        self.schedule_preview()

    def setup_tab_style(self):
        f_style = ttk.LabelFrame(self.tab_style, text="枠線設定", padding=5)
        f_style.pack(fill=tk.X, pady=5)

        # --- 1. Source Image Border ---
        f_src = ttk.LabelFrame(f_style, text="元画像上の枠線 (Source Image)", padding=5)
        f_src.pack(fill=tk.X, pady=5)
        f_src_row = ttk.Frame(f_src)
        f_src_row.pack(fill=tk.X)
        ttk.Checkbutton(
            f_src_row, text="表示する", variable=self.show_crop_border
        ).pack(side=tk.LEFT)
        ttk.Label(f_src_row, text="太さ (pt):").pack(side=tk.LEFT, padx=(20, 5))
        ttk.Entry(f_src_row, textvariable=self.crop_border_w, width=5).pack(
            side=tk.LEFT
        )

        # --- 2. Cropped Image Border ---
        f_zoom = ttk.LabelFrame(
            f_style, text="クロップ画像の枠線 (Cropped Image)", padding=5
        )
        f_zoom.pack(fill=tk.X, pady=5)
        f_zoom_row1 = ttk.Frame(f_zoom)
        f_zoom_row1.pack(fill=tk.X)
        ttk.Checkbutton(
            f_zoom_row1, text="表示する", variable=self.show_zoom_border
        ).pack(side=tk.LEFT)
        ttk.Label(f_zoom_row1, text="太さ (pt):").pack(side=tk.LEFT, padx=(20, 5))
        ttk.Entry(f_zoom_row1, textvariable=self.zoom_border_w, width=5).pack(
            side=tk.LEFT
        )

        f_zoom_row2 = ttk.Frame(f_zoom)
        f_zoom_row2.pack(fill=tk.X, pady=5)
        ttk.Label(f_zoom_row2, text="形状:").pack(side=tk.LEFT)
        ttk.Radiobutton(
            f_zoom_row2,
            text="角 (Rectangle)",
            variable=self.zoom_border_shape,
            value="rectangle",
        ).pack(side=tk.LEFT, padx=5)
        ttk.Radiobutton(
            f_zoom_row2,
            text="丸 (Rounded)",
            variable=self.zoom_border_shape,
            value="rounded",
        ).pack(side=tk.LEFT, padx=5)

    def browse_output(self):
        path = filedialog.asksaveasfilename(
            defaultextension=".pptx", filetypes=[("PPTX", "*.pptx")]
        )
        if path:
            self.output_path.set(path)

    def add_folder(self):
        p = filedialog.askdirectory()
        if p:
            self.folders.append(p)
            self.folder_listbox.insert(tk.END, p)
            self.schedule_preview()

            # Auto-detect aspect ratio from first image
            try:
                images = get_sorted_images(p)
                if images:
                    with Image.open(images[0]) as img:
                        w, h = img.size
                        self.dummy_ratio_w.set(1.0)
                        self.dummy_ratio_h.set(h / w)
            except:
                pass

    def remove_folder(self):
        s = self.folder_listbox.curselection()
        if s:
            self.folders.pop(s[0])
            self.folder_listbox.delete(s[0])
            self.schedule_preview()

    def clear_folders(self):
        self.folders = []
        self.folder_listbox.delete(0, tk.END)
        self.schedule_preview()

    def open_crop_editor(self):
        # Find first available image
        target_img = None
        for f in self.folders:
            imgs = get_sorted_images(f)
            if imgs:
                target_img = imgs[0]
                break

        if not target_img:
            # Ask user for a file
            target_img = filedialog.askopenfilename(
                filetypes=[("Images", "*.png;*.jpg;*.jpeg;*.bmp;*.webp")]
            )

        if not target_img:
            return

        # Always add new region regardless of selection
        CropEditor(self.root, target_img, self.add_region_from_editor)

    def add_region_from_editor(self, x, y, w, h, name):
        r = CropRegion(x, y, w, h, (255, 0, 0), name)
        self.crop_regions.append(r)
        self.update_region_list()
        self.schedule_preview()

    def update_region_from_editor(self, idx, x, y, w, h, name):
        r = self.crop_regions[idx]
        r.x, r.y, r.width, r.height, r.name = x, y, w, h, name
        self.update_region_list()
        self.schedule_preview()

    def add_region_dialog(self):
        d = tk.Toplevel(self.root)
        d.title("Add Region")
        tk.Label(d, text="Name").grid(row=0, column=0)
        v_name = tk.Entry(d)
        v_name.grid(row=0, column=1)
        v_name.insert(0, f"R{len(self.crop_regions) + 1}")

        tk.Label(d, text="x").grid(row=1, column=0)
        vx = tk.Entry(d)
        vx.grid(row=1, column=1)

        tk.Label(d, text="y").grid(row=2, column=0)
        vy = tk.Entry(d)
        vy.grid(row=2, column=1)

        tk.Label(d, text="w").grid(row=3, column=0)
        vw = tk.Entry(d)
        vw.grid(row=3, column=1)

        tk.Label(d, text="h").grid(row=4, column=0)
        vh = tk.Entry(d)
        vh.grid(row=4, column=1)

        c_var = [(255, 0, 0)]
        b_col = tk.Button(
            d,
            text="Color",
            bg="#FF0000",
            fg="white",
            command=lambda: self.pick_color(b_col, c_var),
        )
        b_col.grid(row=5, column=0, columnspan=2)

        def add():
            try:
                self.crop_regions.append(
                    CropRegion(
                        int(vx.get()),
                        int(vy.get()),
                        int(vw.get()),
                        int(vh.get()),
                        c_var[0],
                        v_name.get(),
                    )
                )
                self.update_region_list()
                self.schedule_preview()
                d.destroy()
            except Exception:
                pass

        tk.Button(d, text="OK", command=add).grid(row=6, column=0, columnspan=2)

    def pick_color(self, btn, store):
        c = colorchooser.askcolor(color=store[0])
        if c[0]:
            store[0] = tuple(map(int, c[0]))
            btn.config(bg=c[1])

    def remove_region(self):
        s = self.region_tree.selection()
        if s:
            self.crop_regions.pop(self.region_tree.index(s[0]))
            self.update_region_list()
            self.schedule_preview()

    def update_region_list(self):
        self.region_tree.delete(*self.region_tree.get_children())
        for r in self.crop_regions:
            self.region_tree.insert(
                "",
                tk.END,
                values=(r.name, f"{r.x},{r.y},{r.width},{r.height}", r.align),
            )

    def get_current_config(self) -> Optional[GridConfig]:
        c = GridConfig()
        c.layout_mode = self.layout_mode.get()
        c.flow_align = self.flow_align.get()
        c.slide_width = self.slide_w.get()
        c.slide_height = self.slide_h.get()
        c.rows = self.rows.get()
        c.cols = self.cols.get()
        c.margin_left = self.margin_l.get()
        c.margin_top = self.margin_t.get()
        c.margin_right = self.margin_r.get()
        c.margin_bottom = self.margin_b.get()
        c.gap_h = GapConfig(self.gap_h_val.get(), self.gap_h_mode.get())
        c.gap_v = GapConfig(self.gap_v_val.get(), self.gap_v_mode.get())
        c.size_mode = self.image_size_mode.get()
        c.fit_mode = self.image_fit_mode.get()
        c.image_width = self.image_w.get()
        c.image_height = self.image_h.get()
        c.folders = self.folders if self.folders else ["dummy"] * (c.rows * c.cols)
        c.crop_regions = self.crop_regions

        rs = self.crop_rows_filter.get().strip()
        if rs:
            c.crop_rows = [int(x.strip()) for x in rs.split(",") if x.strip()]
        cs = self.crop_cols_filter.get().strip()
        if cs:
            c.crop_cols = [int(x.strip()) for x in cs.split(",") if x.strip()]

        c.crop_display.position = self.crop_pos.get()
        c.crop_display.main_crop_gap = GapConfig(
            self.gap_mc_val.get(), self.gap_mc_mode.get()
        )
        c.crop_display.crop_crop_gap = GapConfig(
            self.gap_cc_val.get(), self.gap_cc_mode.get()
        )
        c.crop_display.crop_bottom_gap = GapConfig(
            self.gap_cb_val.get(), self.gap_cb_mode.get()
        )

        if self.crop_size_mode.get() == "size":
            if self.crop_size_val.get() > 0:
                c.crop_display.size = self.crop_size_val.get()
        else:
            if self.crop_scale_val.get() > 0:
                c.crop_display.scale = self.crop_scale_val.get()

        c.zoom_border_shape = self.zoom_border_shape.get()
        c.show_crop_border = self.show_crop_border.get()
        c.crop_border_width = self.crop_border_w.get()
        c.show_zoom_border = self.show_zoom_border.get()
        c.zoom_border_width = self.zoom_border_w.get()

        if c.rows == 0:
            c.rows = 1
        if c.cols == 0:
            c.cols = 1
        return c

    def update_preview(self):
        config = self.get_current_config()
        canvas = self.preview_canvas
        canvas.delete("all")
        c_w, c_h = canvas.winfo_width(), canvas.winfo_height()
        if c_w < 10:
            return

        scale = min((c_w - 20) / config.slide_width, (c_h - 20) / config.slide_height)
        sx = (c_w - config.slide_width * scale) / 2
        sy = (c_h - config.slide_height * scale) / 2

        def tx(v):
            return sx + v * scale

        def ty(v):
            return sy + v * scale

        canvas.create_rectangle(
            tx(0),
            ty(0),
            tx(config.slide_width),
            ty(config.slide_height),
            fill="white",
            outline="black",
        )
        metrics = calculate_grid_metrics(config)
        cy = config.margin_top

        d_rw = self.dummy_ratio_w.get()
        d_rh = self.dummy_ratio_h.get()
        if d_rw <= 0:
            d_rw = 1.0
        if d_rh <= 0:
            d_rh = 1.0
        dummy_w_px = 400
        dummy_h_px = 400 * (d_rh / d_rw)

        # Border offset for preview
        border_offset_cm = 0.0
        if config.show_zoom_border:
            border_offset_cm = pt_to_cm(config.zoom_border_width)

        for r in range(config.rows):
            current_row_h = (
                metrics.row_heights[r]
                if r < len(metrics.row_heights)
                else metrics.main_height
            )
            cx = config.margin_left

            # --- Flow Mode Alignment Logic (Preview) ---
            row_content_width = 0.0
            valid_items = 0
            if config.layout_mode == "flow":
                # Pre-calc row width
                sim_cx = 0
                for c in range(config.cols):
                    img_w_cm, img_h_cm = calculate_size_fit_static(
                        dummy_w_px,
                        dummy_h_px,
                        metrics.main_width,
                        metrics.main_height,
                        config.fit_mode,
                    )
                    has_crops = should_apply_crop(r, c, config)

                    item_width = img_w_cm
                    if has_crops and config.crop_regions:
                        num_crops = len(config.crop_regions)
                        disp = config.crop_display
                        actual_gap_mc = disp.main_crop_gap.to_cm(
                            img_w_cm if disp.position == "right" else img_h_cm
                        )
                        if config.show_zoom_border:
                            actual_gap_mc += border_offset_cm

                        if disp.position == "right":
                            # Add crops width
                            c_w_dummy = 0
                            if disp.scale is not None or disp.size is not None:
                                if disp.size:
                                    c_w_dummy = disp.size
                                else:
                                    c_w_dummy = img_w_cm * disp.scale
                            else:
                                single_h = (img_h_cm) / num_crops  # approx
                                c_w_dummy, _ = calculate_size_fit_static(
                                    100, 100, metrics.crop_size, single_h, "fit"
                                )

                            max_crop_ext = 0
                            for ci, r_crop in enumerate(config.crop_regions):
                                this_gap = (
                                    r_crop.gap
                                    if r_crop.gap is not None
                                    else actual_gap_mc
                                )
                                current_ext = this_gap + c_w_dummy
                                if current_ext > max_crop_ext:
                                    max_crop_ext = current_ext
                            item_width += max_crop_ext
                        else:
                            # bottom
                            # Check if crop stack is wider than main
                            actual_gap_cc = disp.crop_crop_gap.to_cm(
                                img_h_cm
                            )  # bottom uses height ref approx
                            if config.show_zoom_border:
                                actual_gap_cc += border_offset_cm

                            # Estimate widths
                            c_w_sum = 0
                            for ci, r_crop in enumerate(config.crop_regions):
                                c_w = 0
                                if disp.size is not None:
                                    c_w, _ = calculate_size_fit_static(
                                        100, 100, 9999, disp.size, "height"
                                    )
                                elif disp.scale is not None:
                                    th = img_h_cm * disp.scale
                                    c_w, _ = calculate_size_fit_static(
                                        100, 100, 9999, th, "height"
                                    )
                                else:
                                    sw = (
                                        img_w_cm - actual_gap_cc * (num_crops - 1)
                                    ) / num_crops
                                    c_w, _ = calculate_size_fit_static(
                                        100, 100, sw, metrics.crop_size, "fit"
                                    )
                                c_w_sum += c_w

                            if num_crops > 1:
                                c_w_sum += actual_gap_cc * (num_crops - 1)
                            if c_w_sum > item_width:
                                item_width = c_w_sum

                    row_content_width += item_width
                    valid_items += 1

                if valid_items > 1:
                    gap_val = config.gap_h.to_cm(metrics.main_width)  # approx
                    row_content_width += (valid_items - 1) * gap_val

                avail_w = config.slide_width - config.margin_left - config.margin_right
                if config.flow_align == "center":
                    cx = config.margin_left + (avail_w - row_content_width) / 2
                elif config.flow_align == "right":
                    cx = config.margin_left + (avail_w - row_content_width)

            # --- Draw Loop ---
            for c in range(config.cols):
                img_w_cm, img_h_cm = calculate_size_fit_static(
                    dummy_w_px,
                    dummy_h_px,
                    metrics.main_width,
                    metrics.main_height,
                    config.fit_mode,
                )
                this_gap_h = config.gap_h.to_cm(img_w_cm)

                if config.layout_mode == "flow":
                    main_l = cx
                else:
                    main_l = cx + (metrics.main_width - img_w_cm) / 2
                main_t = cy + (metrics.main_height - img_h_cm) / 2

                has_crops = should_apply_crop(r, c, config)
                canvas.create_rectangle(
                    tx(main_l),
                    ty(main_t),
                    tx(main_l + img_w_cm),
                    ty(main_t + img_h_cm),
                    fill="#ddf",
                    outline="blue",
                )

                content_r = main_l + img_w_cm
                if has_crops and config.crop_regions:
                    num_crops = len(config.crop_regions)
                    disp = config.crop_display
                    actual_gap_mc = config.crop_display.main_crop_gap.to_cm(
                        img_w_cm if disp.position == "right" else img_h_cm
                    )
                    actual_gap_cc = config.crop_display.crop_crop_gap.to_cm(
                        img_w_cm if disp.position == "right" else img_h_cm
                    )

                    # Compensate for border
                    if config.show_zoom_border:
                        actual_gap_mc += border_offset_cm
                        actual_gap_cc += border_offset_cm

                    for crop_idx in range(num_crops):
                        region = config.crop_regions[crop_idx]
                        dummy_cw, dummy_ch = 100, 100

                        if disp.position == "right":
                            start_x = main_l + img_w_cm + actual_gap_mc
                            c_w, c_h = 0, 0
                            if disp.scale is not None or disp.size is not None:
                                if disp.size:
                                    c_w = c_h = disp.size
                                else:
                                    c_w = img_w_cm * disp.scale
                                    c_h = c_w
                                c_l = start_x
                                c_t = main_t + crop_idx * (c_h + actual_gap_cc)
                            else:
                                single_h = (
                                    img_h_cm - actual_gap_cc * (num_crops - 1)
                                ) / num_crops
                                c_w, c_h = calculate_size_fit_static(
                                    dummy_cw,
                                    dummy_ch,
                                    metrics.crop_size,
                                    single_h,
                                    "fit",
                                )
                                c_l = start_x
                                if crop_idx == 0:
                                    c_t = main_t
                                elif num_crops > 1 and crop_idx == num_crops - 1:
                                    c_t = (main_t + img_h_cm) - c_h
                                else:
                                    slot_top = main_t + crop_idx * (
                                        single_slot_height + actual_gap_cc
                                    )
                                    c_t = slot_top + (single_slot_height - c_h) / 2

                            # Custom Alignment Override
                            this_gap_mc = (
                                region.gap if region.gap is not None else actual_gap_mc
                            )
                            if region.gap is not None and config.show_zoom_border:
                                this_gap_mc += border_offset_cm

                            c_l = main_l + img_w_cm + this_gap_mc
                            if region.align == "start":
                                c_t = main_t + region.offset
                            elif region.align == "center":
                                c_t = main_t + (img_h_cm - c_h) / 2 + region.offset
                            elif region.align == "end":
                                c_t = main_t + img_h_cm - c_h + region.offset

                            canvas.create_rectangle(
                                tx(c_l),
                                ty(c_t),
                                tx(c_l + c_w),
                                ty(c_t + c_h),
                                fill="#fdd",
                                outline="red",
                            )
                            content_r = max(content_r, c_l + c_w)

                        else:  # bottom
                            start_y = main_t + img_h_cm + actual_gap_mc
                            c_w, c_h = 0, 0
                            if disp.scale is not None or disp.size is not None:
                                if disp.size:
                                    c_w = c_h = disp.size
                                else:
                                    c_h = img_h_cm * disp.scale
                                    c_w = c_h
                                c_t = start_y
                                c_l = main_l + crop_idx * (c_w + actual_gap_cc)
                            else:
                                single_w = (
                                    img_w_cm - actual_gap_cc * (num_crops - 1)
                                ) / num_crops
                                c_w, c_h = calculate_size_fit_static(
                                    dummy_cw,
                                    dummy_ch,
                                    single_slot_width,
                                    metrics.crop_size,
                                    "fit",
                                )
                                c_t = start_y
                                if crop_idx == 0:
                                    c_l = main_l
                                elif num_crops > 1 and crop_idx == num_crops - 1:
                                    c_l = (main_l + img_w_cm) - c_w
                                else:
                                    slot_left = main_l + crop_idx * (
                                        single_slot_width + actual_gap_cc
                                    )
                                    c_l = slot_left + (single_slot_width - c_w) / 2

                            # Custom Alignment Override
                            this_gap_mc = (
                                region.gap if region.gap is not None else actual_gap_mc
                            )
                            if region.gap is not None and config.show_zoom_border:
                                this_gap_mc += border_offset_cm

                            c_t = main_t + img_h_cm + this_gap_mc
                            if region.align == "start":
                                c_l = main_l + region.offset
                            elif region.align == "center":
                                c_l = main_l + (img_w_cm - c_w) / 2 + region.offset
                            elif region.align == "end":
                                c_l = main_l + img_w_cm - c_w + region.offset

                            canvas.create_rectangle(
                                tx(c_l),
                                ty(c_t),
                                tx(c_l + c_w),
                                ty(c_t + c_h),
                                fill="#fdd",
                                outline="red",
                            )
                            content_r = max(content_r, c_l + c_w)

                if config.layout_mode == "flow":
                    cx += (content_r - main_l) + this_gap_h
                else:
                    cx += (
                        metrics.col_widths[c] + this_gap_h
                        if c < len(metrics.col_widths)
                        else metrics.main_width + this_gap_h
                    )
            cy += current_row_h + config.gap_v.to_cm(metrics.main_height)

    def save_config(self):
        f = filedialog.asksaveasfilename(
            defaultextension=".yaml", filetypes=[("YAML", "*.yaml")]
        )
        if not f:
            return
        try:
            c = self.get_current_config()
            data = {
                "slide": {"width": c.slide_width, "height": c.slide_height},
                "grid": {
                    "rows": c.rows,
                    "cols": c.cols,
                    "arrangement": c.arrangement,
                    "layout_mode": c.layout_mode,
                    "flow_align": c.flow_align,
                },
                "margin": {
                    "left": c.margin_left,
                    "top": c.margin_top,
                    "right": c.margin_right,
                    "bottom": c.margin_bottom,
                },
                "gap": {
                    "horizontal": {"value": c.gap_h.value, "mode": c.gap_h.mode},
                    "vertical": {"value": c.gap_v.value, "mode": c.gap_v.mode},
                },
                "image": {
                    "size_mode": c.size_mode,
                    "fit_mode": c.fit_mode,
                    "width": c.image_width,
                    "height": c.image_height,
                },
                "folders": c.folders,
                "output": self.output_path.get(),
                "border": {
                    "crop": {
                        "show": self.show_crop_border.get(),
                        "width": self.crop_border_w.get(),
                    },
                    "zoom": {
                        "show": self.show_zoom_border.get(),
                        "width": self.zoom_border_w.get(),
                        "shape": self.zoom_border_shape.get(),
                    },
                },
                "crop": {
                    "regions": [
                        {
                            "name": r.name,
                            "x": r.x,
                            "y": r.y,
                            "width": r.width,
                            "height": r.height,
                            "color": list(r.color),
                            "align": r.align,
                            "offset": r.offset,
                            "gap": r.gap,
                        }
                        for r in c.crop_regions
                    ],
                    "rows": c.crop_rows,
                    "cols": c.crop_cols,
                    "display": {
                        "position": c.crop_display.position,
                        "size": c.crop_display.size,
                        "scale": c.crop_display.scale,
                        "main_crop_gap": {
                            "value": c.crop_display.main_crop_gap.value,
                            "mode": c.crop_display.main_crop_gap.mode,
                        },
                        "crop_crop_gap": {
                            "value": c.crop_display.crop_crop_gap.value,
                            "mode": c.crop_display.crop_crop_gap.mode,
                        },
                        "crop_bottom_gap": {
                            "value": c.crop_display.crop_bottom_gap.value,
                            "mode": c.crop_display.crop_bottom_gap.mode,
                        },
                    },
                },
            }
            with open(f, "w", encoding="utf-8") as yf:
                yaml.dump(data, yf, allow_unicode=True)
            messagebox.showinfo("Saved", "Config saved.")
        except Exception as e:
            messagebox.showerror("Error", str(e))

    def load_config_gui(self):
        f = filedialog.askopenfilename(filetypes=[("YAML", "*.yaml")])
        if not f:
            return
        try:
            config = load_config(f)
            self.apply_config_to_gui(config)
        except Exception as e:
            messagebox.showerror("Error", str(e))

    def apply_config_to_gui(self, c: GridConfig):
        self.layout_mode.set(c.layout_mode)
        self.flow_align.set(c.flow_align)
        self.rows.set(c.rows)
        self.cols.set(c.cols)
        self.arrangement.set(c.arrangement)
        self.slide_w.set(c.slide_width)
        self.slide_h.set(c.slide_height)
        self.margin_l.set(c.margin_left)
        self.margin_t.set(c.margin_top)
        self.margin_r.set(c.margin_right)
        self.margin_b.set(c.margin_bottom)
        self.gap_h_val.set(c.gap_h.value)
        self.gap_h_mode.set(c.gap_h.mode)
        self.gap_v_val.set(c.gap_v.value)
        self.gap_v_mode.set(c.gap_v.mode)
        self.image_size_mode.set(c.size_mode)
        self.image_fit_mode.set(c.fit_mode)
        if c.image_width:
            self.image_w.set(c.image_width)
        if c.image_height:
            self.image_h.set(c.image_height)
        self.folders = c.folders
        self.folder_listbox.delete(0, tk.END)
        for p in self.folders:
            self.folder_listbox.insert(tk.END, p)
        self.crop_regions = c.crop_regions
        self.update_region_list()

        if c.crop_rows:
            self.crop_rows_filter.set(",".join(map(str, c.crop_rows)))
        else:
            self.crop_rows_filter.set("")
        if c.crop_cols:
            self.crop_cols_filter.set(",".join(map(str, c.crop_cols)))
        else:
            self.crop_cols_filter.set("")

        self.crop_pos.set(c.crop_display.position)
        self.gap_mc_val.set(c.crop_display.main_crop_gap.value)
        self.gap_mc_mode.set(c.crop_display.main_crop_gap.mode)
        self.gap_cc_val.set(c.crop_display.crop_crop_gap.value)
        self.gap_cc_mode.set(c.crop_display.crop_crop_gap.mode)
        self.gap_cb_val.set(c.crop_display.crop_bottom_gap.value)
        self.gap_cb_mode.set(c.crop_display.crop_bottom_gap.mode)

        if c.crop_display.size:
            self.crop_size_mode.set("size")
            self.crop_size_val.set(c.crop_display.size)
        elif c.crop_display.scale:
            self.crop_size_mode.set("scale")
            self.crop_scale_val.set(c.crop_display.scale)

        self.zoom_border_shape.set(c.zoom_border_shape)
        self.show_crop_border.set(c.show_crop_border)
        self.crop_border_w.set(c.crop_border_width)
        self.show_zoom_border.set(c.show_zoom_border)
        self.zoom_border_w.set(c.zoom_border_width)

    def generate(self):
        try:
            create_grid_presentation(self.get_current_config())
            messagebox.showinfo("OK", "Generated!")
        except Exception as e:
            messagebox.showerror("Error", str(e))


if __name__ == "__main__":
    root = tk.Tk()
    app = ImageGridApp(root)
    root.mainloop()
