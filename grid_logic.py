"""
Core logic for PowerPoint Image Grid Generator.

This module handles configuration loading, layout calculations, image processing,
and PowerPoint generation. Strictly follows the logic required for v24 compatibility.
"""

import os
import re
import shutil
import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import List, Optional, Tuple, Dict, Union, Any

import yaml
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

# --- Constants ---
CM_TO_EMU = 360000
PT_TO_EMU = 12700


# --- Conversion Helpers ---
def cm_to_emu(cm: float) -> int:
    """Convert centimeters to English Metric Units (EMU)."""
    return int(cm * CM_TO_EMU)


def pt_to_emu(pt: float) -> int:
    """Convert points to English Metric Units (EMU)."""
    return int(pt * PT_TO_EMU)


def pt_to_cm(pt: float) -> float:
    """Convert points to centimeters."""
    # 1 inch = 2.54 cm = 72 points
    return pt * (2.54 / 72.0)


# --- Data Classes ---
@dataclass
class CropRegion:
    """Data class representing a crop region within an image."""

    x: int
    y: int
    width: int
    height: int
    color: Tuple[int, int, int] = (255, 0, 0)
    name: str = ""
    align: str = "auto"  # 'auto', 'start', 'center', 'end'
    offset: float = 0.0
    gap: Optional[float] = None


@dataclass
class GapConfig:
    """Configuration for gaps, supporting absolute (cm) or relative (scale) modes."""

    value: float = 0.5
    mode: str = "cm"  # 'cm' or 'scale'

    def to_cm(self, ref_size: float) -> float:
        """Convert gap to centimeters based on mode and reference size."""
        if self.mode == "scale":
            return ref_size * self.value
        return self.value


@dataclass
class CropDisplayConfig:
    """Configuration for how crops are displayed relative to the main image."""

    position: str = "right"
    main_crop_gap: GapConfig = field(default_factory=lambda: GapConfig(0.15, "cm"))
    crop_crop_gap: GapConfig = field(default_factory=lambda: GapConfig(0.15, "cm"))
    crop_bottom_gap: GapConfig = field(default_factory=lambda: GapConfig(0.0, "cm"))
    size: Optional[float] = None
    scale: Optional[float] = None


@dataclass
class GridConfig:
    """Master configuration for the grid layout and slide settings."""

    slide_width: float = 33.867
    slide_height: float = 19.05
    rows: int = 2
    cols: int = 3
    margin_left: float = 1.0
    margin_top: float = 1.0
    margin_right: float = 1.0
    margin_bottom: float = 1.0

    gap_h: GapConfig = field(default_factory=lambda: GapConfig(0.5, "cm"))
    gap_v: GapConfig = field(default_factory=lambda: GapConfig(0.5, "cm"))

    layout_mode: str = "grid"  # 'grid' or 'flow'
    flow_align: str = "left"
    flow_vertical_align: str = "center"

    size_mode: str = "fit"
    fit_mode: str = "fit"
    image_width: Optional[float] = None
    image_height: Optional[float] = None
    image_scale: float = 1.0

    arrangement: str = "row"
    crop_regions: List[CropRegion] = field(default_factory=list)
    crop_rows: Optional[List[int]] = None
    crop_cols: Optional[List[int]] = None
    crop_display: CropDisplayConfig = field(default_factory=CropDisplayConfig)

    show_crop_border: bool = True
    crop_border_width: float = 1.5
    show_zoom_border: bool = True
    zoom_border_width: float = 1.5
    zoom_border_shape: str = "rectangle"

    folders: List[str] = field(default_factory=list)
    output: str = "output.pptx"


@dataclass
class LayoutMetrics:
    """Calculated metrics for layout positioning."""

    main_width: float
    main_height: float
    col_widths: List[float]
    row_heights: List[float]
    crop_size: float
    crop_main_gap: float
    crop_crop_gap: float
    crop_bottom_gap: float


# --- Parsing Helpers ---
def parse_color(color_value: Union[List, Tuple, str]) -> Tuple[int, int, int]:
    """Parse color input (hex string, list, tuple) to RGB tuple."""
    if isinstance(color_value, (list, tuple)):
        return tuple(color_value[:3])  # type: ignore
    elif isinstance(color_value, str):
        if color_value.startswith("#"):
            hex_color = color_value.lstrip("#")
            return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))
        else:
            parts = color_value.split(",")
            return tuple(int(p.strip()) for p in parts)
    return (255, 0, 0)


def parse_gap(data: Union[int, float, Dict]) -> GapConfig:
    """Parse gap configuration."""
    if isinstance(data, (int, float)):
        return GapConfig(float(data), "cm")
    elif isinstance(data, dict):
        return GapConfig(float(data.get("value", 0.5)), data.get("mode", "cm"))
    return GapConfig(0.5, "cm")


def load_config(config_path: str) -> GridConfig:
    """Load and parse YAML configuration file safely."""
    with open(config_path, "r", encoding="utf-8") as f:
        data = yaml.safe_load(f)

    if data is None:
        data = {}

    config = GridConfig()

    def get_safe_dict(source, key):
        val = source.get(key)
        return val if isinstance(val, dict) else {}

    # Slide
    slide = get_safe_dict(data, "slide")
    config.slide_width = slide.get("width", config.slide_width)
    config.slide_height = slide.get("height", config.slide_height)

    # Grid
    grid = get_safe_dict(data, "grid")
    config.rows = grid.get("rows", config.rows)
    config.cols = grid.get("cols", config.cols)
    config.arrangement = grid.get("arrangement", config.arrangement)
    config.layout_mode = grid.get("layout_mode", config.layout_mode)
    config.flow_align = grid.get("flow_align", config.flow_align)
    config.flow_vertical_align = grid.get(
        "flow_vertical_align", config.flow_vertical_align
    )

    # Margin
    margin = data.get("margin")
    if isinstance(margin, (int, float)):
        config.margin_left = config.margin_top = config.margin_right = (
            config.margin_bottom
        ) = float(margin)
    elif isinstance(margin, dict):
        config.margin_left = margin.get("left", config.margin_left)
        config.margin_top = margin.get("top", config.margin_top)
        config.margin_right = margin.get("right", config.margin_right)
        config.margin_bottom = margin.get("bottom", config.margin_bottom)

    # Gap
    gap = data.get("gap")
    if isinstance(gap, (int, float)):
        val = float(gap)
        config.gap_h = GapConfig(val, "cm")
        config.gap_v = GapConfig(val, "cm")
    elif isinstance(gap, dict):
        if "horizontal" in gap:
            config.gap_h = parse_gap(gap["horizontal"])
        if "vertical" in gap:
            config.gap_v = parse_gap(gap["vertical"])

    # Image
    img = get_safe_dict(data, "image")
    config.size_mode = img.get("size_mode", config.size_mode)
    config.fit_mode = img.get("fit_mode", config.fit_mode)
    config.image_scale = img.get("scale", config.image_scale)
    config.image_width = img.get("width")
    config.image_height = img.get("height")

    # Crop
    crop = get_safe_dict(data, "crop")
    regions_data = crop.get("regions")

    if regions_data is None:
        if "region" in crop and isinstance(crop["region"], dict):
            regions_data = [crop["region"]]
        else:
            regions_data = []

    for i, r in enumerate(regions_data):
        if not isinstance(r, dict):
            continue
        region = CropRegion(
            x=r.get("x", 0),
            y=r.get("y", 0),
            width=r.get("width", 100),
            height=r.get("height", 100),
            color=parse_color(r.get("color", "#FF0000")),
            name=r.get("name", f"crop_{i + 1}"),
            align=r.get("align", "auto"),
            offset=r.get("offset", 0.0),
            gap=r.get("gap", None),
        )
        config.crop_regions.append(region)

    config.crop_rows = crop.get("rows")
    config.crop_cols = crop.get("cols")

    if "display" in crop and isinstance(crop["display"], dict):
        disp = crop["display"]
        config.crop_display.position = disp.get("position", "right")
        config.crop_display.size = disp.get("size")
        config.crop_display.scale = disp.get("scale")

        legacy_gap = disp.get("gap")
        if legacy_gap is not None and isinstance(legacy_gap, (int, float)):
            val = float(legacy_gap)
            config.crop_display.main_crop_gap = GapConfig(val, "cm")
            config.crop_display.crop_crop_gap = GapConfig(val, "cm")

        if "main_crop_gap" in disp:
            config.crop_display.main_crop_gap = parse_gap(disp["main_crop_gap"])
        if "crop_crop_gap" in disp:
            config.crop_display.crop_crop_gap = parse_gap(disp["crop_crop_gap"])
        if "crop_bottom_gap" in disp:
            config.crop_display.crop_bottom_gap = parse_gap(disp["crop_bottom_gap"])

    # Border
    border = get_safe_dict(data, "border")
    if "crop" in border:
        cb = border["crop"]
        if isinstance(cb, dict):
            config.show_crop_border = cb.get("show", config.show_crop_border)
            config.crop_border_width = cb.get("width", config.crop_border_width)
    if "zoom" in border:
        zb = border["zoom"]
        if isinstance(zb, dict):
            config.show_zoom_border = zb.get("show", config.show_zoom_border)
            config.zoom_border_width = zb.get("width", config.zoom_border_width)
            config.zoom_border_shape = zb.get("shape", config.zoom_border_shape)

    config.folders = data.get("folders") or []
    if not isinstance(config.folders, list):
        config.folders = []

    config.output = data.get("output", config.output)

    return config


# --- Image Processing Helpers ---
def extract_number_from_filename(filename: str) -> int:
    numbers = re.findall(r"\d+", filename)
    return int(numbers[0]) if numbers else 0


def get_sorted_images(folder_path: str) -> List[str]:
    supported = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp"}
    folder = Path(folder_path)
    if not folder.exists():
        return []
    files = [
        f for f in folder.iterdir() if f.is_file() and f.suffix.lower() in supported
    ]
    files.sort(key=lambda f: extract_number_from_filename(f.stem))
    return [str(f) for f in files]


def crop_image(image_path: str, region: CropRegion, output_path: str) -> str:
    with Image.open(image_path) as img:
        img.crop(
            (region.x, region.y, region.x + region.width, region.y + region.height)
        ).save(output_path)
    return output_path


def get_image_dimensions(image_path: str) -> Tuple[int, int]:
    with Image.open(image_path) as img:
        return img.size


def should_apply_crop(row: int, col: int, config: GridConfig) -> bool:
    if not config.crop_regions:
        return False
    if config.crop_rows is not None and row not in config.crop_rows:
        return False
    if config.crop_cols is not None and col not in config.crop_cols:
        return False
    return True


def calculate_size_fit_static(
    img_w: int, img_h: int, max_width: float, max_height: float, fit_mode: str = "fit"
) -> Tuple[float, float]:
    """Calculate dimensions maintaining aspect ratio."""
    if max_height <= 0 or max_width <= 0 or img_h == 0:
        return 0, 0
    aspect = img_w / img_h
    if fit_mode == "width":
        width = max_width
        height = width / aspect
    elif fit_mode == "height":
        height = max_height
        width = height * aspect
    else:
        cell_aspect = max_width / max_height
        if aspect > cell_aspect:
            width = max_width
            height = width / aspect
        else:
            height = max_height
            width = height * aspect
    return width, height


def calculate_image_size_fit(
    image_path: str, max_width: float, max_height: float, fit_mode: str = "fit"
) -> Tuple[float, float]:
    try:
        w, h = get_image_dimensions(image_path)
    except Exception:
        return 0, 0
    return calculate_size_fit_static(w, h, max_width, max_height, fit_mode)


def _set_shape_join_style(shape, join_style="miter"):
    """
    Explicitly set the line join style of a shape using OXML.
    join_style: 'miter', 'round', 'bevel'
    """
    # FIX: Access the internal LineFormat proxy correctly
    line_format = shape.line

    # Access the underlying <a:ln> element
    ln = line_format._element
    if ln is None:
        return

    # 1. Remove existing join elements
    for tag in ["round", "bevel", "miter"]:
        element = ln.find(qn(f"a:{tag}"))
        if element is not None:
            ln.remove(element)

    # 2. Create new join element
    if join_style == "miter":
        # miter usually requires a limit (lim), default 800000
        join_elem = parse_xml(
            '<a:miter xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" lim="800000"/>'
        )
    else:
        join_elem = parse_xml(
            f'<a:{join_style} xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
        )

    # 3. Insert in correct sequence
    # Schema sequence for a:ln:
    # (noFill|solidFill|...|pattFill), (prstDash|custDash)?, (round|bevel|miter)?, (headEnd)?, (tailEnd)?, (extLst)?
    successors = ["headEnd", "tailEnd", "extLst"]
    inserted = False
    for tag in successors:
        match = ln.find(qn(f"a:{tag}"))
        if match is not None:
            match.addprevious(join_elem)
            inserted = True
            break

    if not inserted:
        ln.append(join_elem)


def add_border_shape(
    slide, left, top, width, height, color, border_width, shape_type="rectangle"
):
    mst = (
        MSO_SHAPE.ROUNDED_RECTANGLE if shape_type == "rounded" else MSO_SHAPE.RECTANGLE
    )
    shape = slide.shapes.add_shape(
        mst, cm_to_emu(left), cm_to_emu(top), cm_to_emu(width), cm_to_emu(height)
    )
    shape.fill.background()
    shape.line.color.rgb = RGBColor(*color)
    shape.line.width = pt_to_emu(border_width)
    shape.shadow.inherit = False

    # Apply join style correction
    if shape_type == "rectangle":
        _set_shape_join_style(shape, "miter")
    else:
        _set_shape_join_style(shape, "round")


def add_crop_borders_to_image(
    slide, left, top, width, height, img_path, regions, border_width
):
    try:
        ow, oh = get_image_dimensions(img_path)
    except Exception:
        return
    sx = width / ow
    sy = height / oh
    for r in regions:
        add_border_shape(
            slide,
            left + r.x * sx,
            top + r.y * sy,
            r.width * sx,
            r.height * sy,
            r.color,
            border_width,
            "rectangle",
        )


def calculate_grid_metrics(config: GridConfig) -> LayoutMetrics:
    total_w = config.slide_width - config.margin_left - config.margin_right
    total_h = config.slide_height - config.margin_top - config.margin_bottom
    est_cols = max(1, config.cols)
    est_rows = max(1, config.rows)
    est_main_w = total_w / est_cols
    est_main_h = total_h / est_rows

    gap_h = config.gap_h.to_cm(est_main_w)
    gap_v = config.gap_v.to_cm(est_main_h)

    avail_w = total_w - (gap_h * (config.cols - 1))
    avail_h = total_h - (gap_v * (config.rows - 1))

    has_crops = len(config.crop_regions) > 0
    exp_cols = set()
    exp_rows = set()
    if has_crops:
        exp_cols = (
            set(range(config.cols))
            if config.crop_cols is None
            else set(config.crop_cols)
        )
        exp_rows = (
            set(range(config.rows))
            if config.crop_rows is None
            else set(config.crop_rows)
        )

    n_exp_c = len([c for c in range(config.cols) if c in exp_cols])
    n_exp_r = len([r for r in range(config.rows) if r in exp_rows])

    disp = config.crop_display
    gap_mc = disp.main_crop_gap.to_cm(est_main_w)
    gap_cc = disp.crop_crop_gap.to_cm(est_main_w)
    gap_cb = disp.crop_bottom_gap.to_cm(
        est_main_w if disp.position == "right" else est_main_h
    )

    if config.show_zoom_border:
        bo = pt_to_cm(config.zoom_border_width)
        gap_mc += bo
        gap_cc += bo

    main_w = 0.0
    main_h = 0.0
    crop_size = 0.0
    is_fixed = (
        config.size_mode == "fixed" and config.image_width and config.image_width > 0
    )

    if is_fixed:
        main_w = config.image_width
        main_h = config.image_height
        if disp.size is not None:
            crop_size = disp.size
        elif disp.scale is not None:
            crop_size = (main_w if disp.position == "right" else main_h) * disp.scale
        else:
            crop_size = (main_w if disp.position == "right" else main_h) * 0.33
    else:
        if disp.position == "right":
            main_h = avail_h / config.rows
            ck = disp.scale if disp.scale is not None else 0.33
            cc = disp.size if disp.size is not None else 0.0
            if disp.size is not None:
                ck = 0.0

            denom = config.cols + n_exp_c * ck
            numer = avail_w - n_exp_c * (gap_mc + cc)
            main_w = numer / denom
            crop_size = main_w * ck + cc
        else:
            main_w = avail_w / config.cols
            ck = disp.scale if disp.scale is not None else 0.33
            cc = disp.size if disp.size is not None else 0.0
            if disp.size is not None:
                ck = 0.0

            denom = config.rows + n_exp_r * ck
            numer = avail_h - n_exp_r * (gap_mc + cc)
            main_h = numer / denom
            crop_size = main_h * ck + cc

    col_widths = []
    if disp.position == "right":
        extra = crop_size + gap_mc
        for c in range(config.cols):
            col_widths.append(main_w + extra if c in exp_cols else main_w)
        row_heights = [main_h] * config.rows
    else:
        extra = crop_size + gap_mc + gap_cb
        col_widths = [main_w] * config.cols
        row_heights = []
        for r in range(config.rows):
            row_heights.append(main_h + extra if r in exp_rows else main_h)

    return LayoutMetrics(
        main_w, main_h, col_widths, row_heights, crop_size, gap_mc, gap_cc, gap_cb
    )


def calculate_item_bounds(
    config, metrics, image_path, row, col, border_offset_cm=0.0, override_size=None
):
    """
    Returns (min_x, min_y, max_x, max_y) relative to Main Image Top-Left (0,0).
    """
    has_crops = should_apply_crop(row, col, config)

    if override_size:
        orig_w, orig_h = override_size
    elif image_path == "dummy":
        return 0, 0, 0, 0
    else:
        orig_w, orig_h = calculate_image_size_fit(
            image_path, metrics.main_width, metrics.main_height, config.fit_mode
        )

    min_x, min_y = 0.0, 0.0
    max_x, max_y = orig_w, orig_h
    half_border = border_offset_cm / 2.0 if config.show_zoom_border else 0.0

    if has_crops and config.crop_regions:
        num_crops = len(config.crop_regions)
        disp = config.crop_display
        ref_s = orig_w if disp.position == "right" else orig_h

        act_gap_mc = disp.main_crop_gap.to_cm(ref_s)
        act_gap_cc = disp.crop_crop_gap.to_cm(ref_s)
        act_gap_cb = disp.crop_bottom_gap.to_cm(ref_s)
        if config.show_zoom_border:
            act_gap_mc += border_offset_cm
            act_gap_cc += border_offset_cm

        if disp.position == "right":
            for i, r in enumerate(config.crop_regions):
                cw, ch = 0, 0
                if override_size:  # Dummy logic
                    if disp.size is not None:
                        cw, ch = calculate_size_fit_static(
                            100, 100, disp.size, 9999, "width"
                        )
                    elif disp.scale is not None:
                        cw, ch = calculate_size_fit_static(
                            100, 100, orig_w * disp.scale, 9999, "width"
                        )
                    else:
                        sh = (orig_h - act_gap_cc * (num_crops - 1)) / num_crops
                        cw, ch = calculate_size_fit_static(
                            100, 100, metrics.crop_size, sh, "fit"
                        )
                else:  # Real logic
                    if disp.size is not None:
                        cw, ch = calculate_image_size_fit(
                            image_path, disp.size, 9999, "width"
                        )
                    elif disp.scale is not None:
                        cw, ch = calculate_image_size_fit(
                            image_path, orig_w * disp.scale, 9999, "width"
                        )
                    else:
                        sh = (orig_h - act_gap_cc * (num_crops - 1)) / num_crops
                        cw, ch = calculate_image_size_fit(
                            image_path, metrics.crop_size, sh, "fit"
                        )

                tg = r.gap if r.gap is not None else act_gap_mc
                if r.gap is not None and config.show_zoom_border:
                    tg += border_offset_cm

                cl = orig_w + tg
                ct = 0.0
                if r.align == "start":
                    ct = 0.0 + r.offset + half_border
                elif r.align == "center":
                    ct = (orig_h - ch) / 2 + r.offset
                elif r.align == "end":
                    ct = orig_h - ch + r.offset - half_border
                else:
                    if disp.scale or disp.size:
                        ct = i * (ch + act_gap_cc)
                    else:
                        if i == 0:
                            ct = 0.0
                        elif num_crops > 1 and i == num_crops - 1:
                            ct = orig_h - ch
                        else:
                            sh = (orig_h - act_gap_cc * (num_crops - 1)) / num_crops
                            st = i * (sh + act_gap_cc)
                            ct = st + (sh - ch) / 2

                if cl < min_x:
                    min_x = cl
                if ct < min_y:
                    min_y = ct
                if (cl + cw) > max_x:
                    max_x = cl + cw
                if (ct + ch) > max_y:
                    max_y = ct + ch
        else:  # Bottom
            for i, r in enumerate(config.crop_regions):
                cw, ch = 0, 0
                if override_size:
                    if disp.size is not None:
                        cw, ch = calculate_size_fit_static(
                            100, 100, 9999, disp.size, "height"
                        )
                    elif disp.scale is not None:
                        cw, ch = calculate_size_fit_static(
                            100, 100, 9999, orig_h * disp.scale, "height"
                        )
                    else:
                        sw = (orig_w - act_gap_cc * (num_crops - 1)) / num_crops
                        cw, ch = calculate_size_fit_static(
                            100, 100, sw, metrics.crop_size, "fit"
                        )
                else:
                    if disp.size is not None:
                        cw, ch = calculate_image_size_fit(
                            image_path, 9999, disp.size, "height"
                        )
                    elif disp.scale is not None:
                        cw, ch = calculate_image_size_fit(
                            image_path, 9999, orig_h * disp.scale, "height"
                        )
                    else:
                        sw = (orig_w - act_gap_cc * (num_crops - 1)) / num_crops
                        cw, ch = calculate_image_size_fit(
                            image_path, sw, metrics.crop_size, "fit"
                        )

                tg = r.gap if r.gap is not None else act_gap_mc
                if r.gap is not None and config.show_zoom_border:
                    tg += border_offset_cm

                ct = orig_h + tg
                cl = 0.0
                if r.align == "start":
                    cl = 0.0 + r.offset + half_border
                elif r.align == "center":
                    cl = (orig_w - cw) / 2 + r.offset
                elif r.align == "end":
                    cl = orig_w - cw + r.offset - half_border
                else:
                    if disp.scale or disp.size:
                        cl = i * (cw + act_gap_cc)
                    else:
                        if i == 0:
                            cl = 0.0
                        elif num_crops > 1 and i == num_crops - 1:
                            cl = orig_w - cw
                        else:
                            sw = (orig_w - act_gap_cc * (num_crops - 1)) / num_crops
                            sl = i * (sw + act_gap_cc)
                            cl = sl + (sw - cw) / 2

                if cl < min_x:
                    min_x = cl
                if ct < min_y:
                    min_y = ct
                if (cl + cw) > max_x:
                    max_x = cl + cw
                if (ct + ch) > max_y:
                    max_y = ct + ch
            max_y += act_gap_cb

    return min_x, min_y, max_x, max_y


def calculate_flow_row_heights(config, metrics, image_grid, border_offset_cm):
    heights = []
    for r, row_imgs in enumerate(image_grid):
        if r >= config.rows:
            break
        rmax = 0.0
        for c, path in enumerate(row_imgs):
            if c >= config.cols or path is None:
                continue
            _, min_y, _, max_y = calculate_item_bounds(
                config, metrics, path, r, c, border_offset_cm
            )
            rmax = max(rmax, max_y - min_y)
        if rmax == 0.0:
            rmax = metrics.main_height
        heights.append(rmax)
    return heights


def create_grid_presentation(config: GridConfig) -> str:
    prs = Presentation()
    prs.slide_width = cm_to_emu(config.slide_width)
    prs.slide_height = cm_to_emu(config.slide_height)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    if config.arrangement == "row":
        grid = [get_sorted_images(f) for f in config.folders]
    else:
        cols = [get_sorted_images(f) for f in config.folders]
        ml = max(len(c) for c in cols) if cols else 0
        grid = [[col[i] if i < len(col) else None for col in cols] for i in range(ml)]

    temp_dir = tempfile.mkdtemp()
    num_crops = len(config.crop_regions)
    metrics = calculate_grid_metrics(config)

    bo = 0.0
    if config.show_zoom_border:
        bo = pt_to_cm(config.crop_border_width)
    half_border = bo / 2.0 if config.show_zoom_border else 0.0

    flow_heights = []
    if config.layout_mode == "flow":
        flow_heights = calculate_flow_row_heights(config, metrics, grid, bo)

    total_h = 0.0
    if config.layout_mode == "flow":
        for i, h in enumerate(flow_heights):
            total_h += h
            if i < len(flow_heights) - 1:
                total_h += config.gap_v.to_cm(metrics.main_height)

    try:
        cur_y = config.margin_top
        if config.layout_mode == "flow":
            ah = config.slide_height - config.margin_top - config.margin_bottom
            if config.flow_vertical_align == "center":
                cur_y += (ah - total_h) / 2
            elif config.flow_vertical_align == "bottom":
                cur_y = (config.margin_top + ah) - total_h

        for r, row_imgs in enumerate(grid):
            if r >= config.rows:
                break

            if config.layout_mode == "flow":
                cur_rh = (
                    flow_heights[r] if r < len(flow_heights) else metrics.main_height
                )
            else:
                cur_rh = (
                    metrics.row_heights[r]
                    if r < len(metrics.row_heights)
                    else metrics.main_height
                )

            gap_v = config.gap_v.to_cm(metrics.main_height)
            cur_x = config.margin_left

            # Flow row width calc
            if config.layout_mode == "flow":
                rw = 0.0
                cnt = 0
                for c, path in enumerate(row_imgs):
                    if c >= config.cols or path is None:
                        continue
                    mnx, _, mxx, _ = calculate_item_bounds(
                        config, metrics, path, r, c, bo
                    )
                    rw += mxx - mnx
                    cnt += 1
                if cnt > 1:
                    rw += (cnt - 1) * config.gap_h.to_cm(metrics.main_width)

                aw = config.slide_width - config.margin_left - config.margin_right
                if config.flow_align == "center":
                    cur_x += (aw - rw) / 2
                elif config.flow_align == "right":
                    cur_x += aw - rw

            for c, path in enumerate(row_imgs):
                if c >= config.cols:
                    break
                gap_h = config.gap_h.to_cm(metrics.main_width)

                if path is None:
                    if config.layout_mode == "grid":
                        w = (
                            metrics.col_widths[c]
                            if c < len(metrics.col_widths)
                            else metrics.main_width
                        )
                        cur_x += w + gap_h
                    continue

                ow, oh = calculate_image_size_fit(
                    path, metrics.main_width, metrics.main_height, config.fit_mode
                )

                # Global gaps
                disp = config.crop_display
                ref_s = ow if disp.position == "right" else oh
                g_mc = disp.main_crop_gap.to_cm(ref_s)
                g_cc = disp.crop_crop_gap.to_cm(ref_s)
                if config.show_zoom_border:
                    g_mc += bo
                    g_cc += bo

                mnx, mny, mxx, mxy = calculate_item_bounds(
                    config, metrics, path, r, c, bo
                )
                iw = mxx - mnx
                ih = mxy - mny

                if config.layout_mode == "flow":
                    dl = cur_x
                    dt = cur_y + (cur_rh - ih) / 2
                else:
                    cw = (
                        metrics.col_widths[c]
                        if c < len(metrics.col_widths)
                        else metrics.main_width
                    )
                    dl = cur_x + (cw - iw) / 2
                    dt = cur_y + (cur_rh - ih) / 2

                fml = dl - mnx
                fmt = dt - mny

                pic = slide.shapes.add_picture(
                    path,
                    cm_to_emu(fml),
                    cm_to_emu(fmt),
                    cm_to_emu(ow),
                    cm_to_emu(oh),
                )
                pic.shadow.inherit = False

                if should_apply_crop(r, c, config):
                    if config.show_crop_border:
                        add_crop_borders_to_image(
                            slide,
                            fml,
                            fmt,
                            ow,
                            oh,
                            path,
                            config.crop_regions,
                            config.crop_border_width,
                        )

                    if num_crops > 0:
                        for i, reg in enumerate(config.crop_regions):
                            cpath = os.path.join(temp_dir, f"crop_{r}_{c}_{i}.png")
                            try:
                                crop_image(path, reg, cpath)
                            except Exception:
                                continue

                            # Size logic
                            cw, ch = 0, 0
                            if disp.size is not None:
                                if disp.position == "right":
                                    cw, ch = calculate_image_size_fit(
                                        cpath, disp.size, 9999, "width"
                                    )
                                else:
                                    cw, ch = calculate_image_size_fit(
                                        cpath, 9999, disp.size, "height"
                                    )
                            elif disp.scale is not None:
                                if disp.position == "right":
                                    cw, ch = calculate_image_size_fit(
                                        cpath, ow * disp.scale, 9999, "width"
                                    )
                                else:
                                    cw, ch = calculate_image_size_fit(
                                        cpath, 9999, oh * disp.scale, "height"
                                    )
                            else:
                                if disp.position == "right":
                                    sh = (oh - g_cc * (num_crops - 1)) / num_crops
                                    cw, ch = calculate_image_size_fit(
                                        cpath, metrics.crop_size, sh, "fit"
                                    )
                                else:
                                    sw = (ow - g_cc * (num_crops - 1)) / num_crops
                                    cw, ch = calculate_image_size_fit(
                                        cpath, sw, metrics.crop_size, "fit"
                                    )

                            tg = reg.gap if reg.gap is not None else g_mc
                            if reg.gap is not None and config.show_zoom_border:
                                tg += bo

                            if disp.position == "right":
                                cl = fml + ow + tg
                                if reg.align == "start":
                                    ct = fmt + reg.offset + half_border
                                elif reg.align == "center":
                                    ct = fmt + (oh - ch) / 2 + reg.offset
                                elif reg.align == "end":
                                    ct = fmt + oh - ch + reg.offset - half_border
                                else:
                                    if disp.scale or disp.size:
                                        ct = fmt + i * (ch + g_cc)
                                    else:
                                        if i == 0:
                                            ct = fmt
                                        elif num_crops > 1 and i == num_crops - 1:
                                            ct = fmt + oh - ch
                                        else:
                                            sh = (
                                                oh - g_cc * (num_crops - 1)
                                            ) / num_crops
                                            st = fmt + i * (sh + g_cc)
                                            ct = st + (sh - ch) / 2
                            else:
                                ct = fmt + oh + tg
                                if reg.align == "start":
                                    cl = fml + reg.offset + half_border
                                elif reg.align == "center":
                                    cl = fml + (ow - cw) / 2 + reg.offset
                                elif reg.align == "end":
                                    cl = fml + ow - cw + reg.offset - half_border
                                else:
                                    if disp.scale or disp.size:
                                        cl = fml + i * (cw + g_cc)
                                    else:
                                        if i == 0:
                                            cl = fml
                                        elif num_crops > 1 and i == num_crops - 1:
                                            cl = fml + ow - cw
                                        else:
                                            sw = (
                                                ow - g_cc * (num_crops - 1)
                                            ) / num_crops
                                            sl = fml + i * (sw + g_cc)
                                            cl = sl + (sw - cw) / 2

                            cpic = slide.shapes.add_picture(
                                cpath,
                                cm_to_emu(cl),
                                cm_to_emu(ct),
                                cm_to_emu(cw),
                                cm_to_emu(ch),
                            )
                            cpic.shadow.inherit = False

                            if config.show_zoom_border:
                                add_border_shape(
                                    slide,
                                    cl,
                                    ct,
                                    cw,
                                    ch,
                                    reg.color,
                                    config.zoom_border_width,
                                    config.zoom_border_shape,
                                )

                if config.layout_mode == "flow":
                    cur_x += iw + gap_h
                else:
                    w = (
                        metrics.col_widths[c]
                        if c < len(metrics.col_widths)
                        else metrics.main_width
                    )
                    cur_x += w + gap_h

            cur_y += cur_rh + gap_v

        prs.save(config.output)
    except Exception as e:
        raise e
    finally:
        shutil.rmtree(temp_dir)
    return config.output
