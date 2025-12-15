"""
PowerPoint Image Grid Generator CLI
This script creates a PowerPoint presentation with images arranged in a grid layout.
It supports all advanced layout features found in the GUI version, including:
- Layout Modes: Grid (aligned) and Flow (packed).
- Flow Alignment: Left, Center, Right.
- Flow Vertical Alignment: Top, Center, Bottom.
- Precise Gap Control: cm or scale relative to image size.
- Image Fit Modes: Fit, Width, Height.
- Per-Crop Alignment: Start, Center, End with offsets.
- Crop-Bottom Gap support.
- Border shapes (Rectangle/Rounded).

Usage:
  python image_grid_cli.py config.yaml
  python image_grid_cli.py --init [filename]
"""

import os
import re
import sys
import shutil
import tempfile
import yaml
from pathlib import Path
from typing import Optional, List, Tuple, Dict, Union
from dataclasses import dataclass, field
from pptx import Presentation
from pptx.util import Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image


# --- Backend Logic (Synced with GUI v24) ---

CM_TO_EMU = 360000
PT_TO_EMU = 12700


def cm_to_emu(cm: float) -> int:
    return int(cm * CM_TO_EMU)


def pt_to_emu(pt: float) -> int:
    return int(pt * PT_TO_EMU)


def pt_to_cm(pt: float) -> float:
    """Convert points to centimeters."""
    # 1 inch = 2.54 cm = 72 points
    return pt * (2.54 / 72.0)


@dataclass
class CropRegion:
    x: int
    y: int
    width: int
    height: int
    color: tuple[int, int, int] = (255, 0, 0)
    name: str = ""
    # Per-crop positioning
    align: str = "auto"  # 'auto', 'start', 'center', 'end'
    offset: float = 0.0  # cm, offset from alignment anchor
    gap: Optional[float] = None  # cm, overrides global main-crop gap if set


@dataclass
class GapConfig:
    value: float = 0.5
    mode: str = "cm"  # 'cm' or 'scale'

    def to_cm(self, ref_size: float) -> float:
        if self.mode == "scale":
            return ref_size * self.value
        return self.value


@dataclass
class CropDisplayConfig:
    position: str = "right"
    main_crop_gap: GapConfig = field(default_factory=lambda: GapConfig(0.15, "cm"))
    crop_crop_gap: GapConfig = field(default_factory=lambda: GapConfig(0.15, "cm"))
    crop_bottom_gap: GapConfig = field(default_factory=lambda: GapConfig(0.0, "cm"))
    size: Optional[float] = None
    scale: Optional[float] = None


@dataclass
class GridConfig:
    slide_width: float = 33.867
    slide_height: float = 19.05
    rows: int = 2
    cols: int = 3
    margin_left: float = 1.0
    margin_top: float = 1.0
    margin_right: float = 1.0
    margin_bottom: float = 1.0

    # Gap settings
    gap_h: GapConfig = field(default_factory=lambda: GapConfig(0.5, "cm"))
    gap_v: GapConfig = field(default_factory=lambda: GapConfig(0.5, "cm"))

    # Layout Logic
    layout_mode: str = "grid"  # 'grid' (aligned) or 'flow' (compact)
    flow_align: str = "left"  # 'left', 'center', 'right' (Only for flow mode)
    flow_vertical_align: str = (
        "center"  # 'top', 'center', 'bottom' (Only for flow mode)
    )

    # Image sizing
    size_mode: str = "fit"  # 'fit' or 'fixed'
    fit_mode: str = "fit"  # 'fit', 'width', 'height'
    image_width: Optional[float] = None
    image_height: Optional[float] = None
    image_scale: float = 1.0

    arrangement: str = "row"
    crop_regions: List[CropRegion] = field(default_factory=list)
    crop_rows: Optional[List[int]] = None
    crop_cols: Optional[List[int]] = None
    crop_display: CropDisplayConfig = field(default_factory=CropDisplayConfig)
    show_crop_border: bool = True
    crop_border_width: float = 1.5
    show_zoom_border: bool = True
    zoom_border_width: float = 1.5
    zoom_border_shape: str = "rectangle"  # 'rectangle' or 'rounded'

    folders: List[str] = field(default_factory=list)
    output: str = "output.pptx"


def parse_color(color_value) -> tuple[int, int, int]:
    if isinstance(color_value, (list, tuple)):
        return tuple(color_value[:3])
    elif isinstance(color_value, str):
        if color_value.startswith("#"):
            hex_color = color_value.lstrip("#")
            return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))
        else:
            parts = color_value.split(",")
            return tuple(int(p.strip()) for p in parts)
    return (255, 0, 0)


def parse_gap(data) -> GapConfig:
    if isinstance(data, (int, float)):
        return GapConfig(float(data), "cm")
    elif isinstance(data, dict):
        return GapConfig(float(data.get("value", 0.5)), data.get("mode", "cm"))
    return GapConfig(0.5, "cm")


def load_config(config_path: str) -> GridConfig:
    with open(config_path, "r", encoding="utf-8") as f:
        data = yaml.safe_load(f)

    config = GridConfig()

    if "slide" in data:
        slide = data["slide"]
        config.slide_width = slide.get("width", config.slide_width)
        config.slide_height = slide.get("height", config.slide_height)

    if "grid" in data:
        grid = data["grid"]
        config.rows = grid.get("rows", config.rows)
        config.cols = grid.get("cols", config.cols)
        config.arrangement = grid.get("arrangement", config.arrangement)
        config.layout_mode = grid.get("layout_mode", config.layout_mode)
        config.flow_align = grid.get("flow_align", config.flow_align)
        config.flow_vertical_align = grid.get(
            "flow_vertical_align", config.flow_vertical_align
        )

    if "margin" in data:
        margin = data["margin"]
        if isinstance(margin, (int, float)):
            config.margin_left = margin
            config.margin_top = margin
            config.margin_right = margin
            config.margin_bottom = margin
        else:
            config.margin_left = margin.get("left", config.margin_left)
            config.margin_top = margin.get("top", config.margin_top)
            config.margin_right = margin.get("right", config.margin_right)
            config.margin_bottom = margin.get("bottom", config.margin_bottom)

    if "gap" in data:
        gap = data["gap"]
        if isinstance(gap, (int, float)):
            config.gap_h = GapConfig(float(gap), "cm")
            config.gap_v = GapConfig(float(gap), "cm")
        else:
            if "horizontal" in gap:
                config.gap_h = parse_gap(gap["horizontal"])
            if "vertical" in gap:
                config.gap_v = parse_gap(gap["vertical"])

    if "image" in data:
        img = data["image"]
        config.size_mode = img.get("size_mode", config.size_mode)
        config.fit_mode = img.get("fit_mode", config.fit_mode)
        config.image_scale = img.get("scale", config.image_scale)
        config.image_width = img.get("width")
        config.image_height = img.get("height")

    if "crop" in data:
        crop = data["crop"]
        if "regions" in crop:
            for i, r in enumerate(crop["regions"]):
                region = CropRegion(
                    x=r.get("x", 0),
                    y=r.get("y", 0),
                    width=r.get("width", 100),
                    height=r.get("height", 100),
                    color=parse_color(r.get("color", "#FF0000")),
                    name=r.get("name", f"crop_{i + 1}"),
                    align=r.get("align", "auto"),
                    offset=r.get("offset", 0.0),
                    gap=r.get("gap", None),
                )
                config.crop_regions.append(region)
        elif "region" in crop:
            r = crop["region"]
            region = CropRegion(
                x=r.get("x", 0),
                y=r.get("y", 0),
                width=r.get("width", 100),
                height=r.get("height", 100),
                color=parse_color(r.get("color", "#FF0000")),
                name="crop_1",
                align=r.get("align", "auto"),
                offset=r.get("offset", 0.0),
                gap=r.get("gap", None),
            )
            config.crop_regions.append(region)

        config.crop_rows = crop.get("rows")
        config.crop_cols = crop.get("cols")

        if "display" in crop:
            disp = crop["display"]
            config.crop_display.position = disp.get("position", "right")
            config.crop_display.size = disp.get("size")
            config.crop_display.scale = disp.get("scale")

            legacy_gap = disp.get("gap")
            if legacy_gap is not None:
                if isinstance(legacy_gap, (int, float)):
                    config.crop_display.main_crop_gap = GapConfig(
                        float(legacy_gap), "cm"
                    )
                    config.crop_display.crop_crop_gap = GapConfig(
                        float(legacy_gap), "cm"
                    )

            if "main_crop_gap" in disp:
                config.crop_display.main_crop_gap = parse_gap(disp["main_crop_gap"])
            if "crop_crop_gap" in disp:
                config.crop_display.crop_crop_gap = parse_gap(disp["crop_crop_gap"])
            if "crop_bottom_gap" in disp:
                config.crop_display.crop_bottom_gap = parse_gap(disp["crop_bottom_gap"])

    if "border" in data:
        border = data["border"]
        if "crop" in border:
            cb = border["crop"]
            config.show_crop_border = cb.get("show", config.show_crop_border)
            config.crop_border_width = cb.get("width", config.crop_border_width)
        if "zoom" in border:
            zb = border["zoom"]
            config.show_zoom_border = zb.get("show", config.show_zoom_border)
            config.zoom_border_width = zb.get("width", config.zoom_border_width)
            config.zoom_border_shape = zb.get("shape", config.zoom_border_shape)

    if "folders" in data:
        config.folders = data["folders"]
    config.output = data.get("output", config.output)

    return config


def extract_number_from_filename(filename: str) -> int:
    numbers = re.findall(r"\d+", filename)
    return int(numbers[0]) if numbers else 0


def get_sorted_images(folder_path: str) -> List[str]:
    supported_extensions = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp"}
    folder = Path(folder_path)
    if not folder.exists():
        return []

    image_files = [
        f
        for f in folder.iterdir()
        if f.is_file() and f.suffix.lower() in supported_extensions
    ]
    image_files.sort(key=lambda f: extract_number_from_filename(f.stem))
    return [str(f) for f in image_files]


def crop_image(image_path: str, region: CropRegion, output_path: str) -> str:
    with Image.open(image_path) as img:
        box = (region.x, region.y, region.x + region.width, region.y + region.height)
        cropped = img.crop(box)
        cropped.save(output_path)
    return output_path


def get_image_dimensions(image_path: str) -> Tuple[int, int]:
    with Image.open(image_path) as img:
        return img.size


def should_apply_crop(row: int, col: int, config: GridConfig) -> bool:
    if not config.crop_regions:
        return False
    if config.crop_rows is not None and row not in config.crop_rows:
        return False
    if config.crop_cols is not None and col not in config.crop_cols:
        return False
    return True


def calculate_size_fit_static(
    img_w: int, img_h: int, max_width: float, max_height: float, fit_mode: str = "fit"
) -> Tuple[float, float]:
    if max_height <= 0 or max_width <= 0:
        return 0, 0
    if img_h == 0:
        return 0, 0
    aspect_ratio = img_w / img_h

    if fit_mode == "width":
        width = max_width
        height = width / aspect_ratio
    elif fit_mode == "height":
        height = max_height
        width = height * aspect_ratio
    else:  # 'fit'
        cell_aspect = max_width / max_height
        if aspect_ratio > cell_aspect:
            width = max_width
            height = width / aspect_ratio
        else:
            height = max_height
            width = height * aspect_ratio

    return width, height


def calculate_image_size_fit(
    image_path: str, max_width: float, max_height: float, fit_mode: str = "fit"
) -> Tuple[float, float]:
    try:
        img_width_px, img_height_px = get_image_dimensions(image_path)
    except Exception:
        return 0, 0
    return calculate_size_fit_static(
        img_width_px, img_height_px, max_width, max_height, fit_mode
    )


def add_border_shape(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    border_color: tuple[int, int, int],
    border_width: float,
    shape_type: str = "rectangle",
):
    ms_shape_type = MSO_SHAPE.RECTANGLE
    if shape_type == "rounded":
        ms_shape_type = MSO_SHAPE.ROUNDED_RECTANGLE

    shape = slide.shapes.add_shape(
        ms_shape_type,
        cm_to_emu(left),
        cm_to_emu(top),
        cm_to_emu(width),
        cm_to_emu(height),
    )
    shape.fill.background()
    shape.line.color.rgb = RGBColor(*border_color)
    shape.line.width = pt_to_emu(border_width)
    shape.shadow.inherit = False


def add_crop_borders_to_image(
    slide,
    image_left: float,
    image_top: float,
    image_width: float,
    image_height: float,
    original_image_path: str,
    crop_regions: List[CropRegion],
    border_width: float,
):
    try:
        orig_width_px, orig_height_px = get_image_dimensions(original_image_path)
    except Exception:
        return

    scale_x = image_width / orig_width_px
    scale_y = image_height / orig_height_px

    for region in crop_regions:
        border_left = image_left + region.x * scale_x
        border_top = image_top + region.y * scale_y
        border_w = region.width * scale_x
        border_h = region.height * scale_y
        add_border_shape(
            slide,
            border_left,
            border_top,
            border_w,
            border_h,
            region.color,
            border_width,
            "rectangle",
        )  # Original image crop marks are always rects


@dataclass
class LayoutMetrics:
    main_width: float
    main_height: float
    col_widths: List[float]
    row_heights: List[float]
    crop_size: float
    crop_main_gap: float
    crop_crop_gap: float
    crop_bottom_gap: float


def calculate_grid_metrics(config: GridConfig) -> LayoutMetrics:
    total_grid_w = config.slide_width - config.margin_left - config.margin_right
    total_grid_h = config.slide_height - config.margin_top - config.margin_bottom

    est_cols = max(1, config.cols)
    est_rows = max(1, config.rows)
    est_main_w = total_grid_w / est_cols
    est_main_h = total_grid_h / est_rows

    gap_h_val = config.gap_h.to_cm(est_main_w)
    gap_v_val = config.gap_v.to_cm(est_main_h)

    avail_w_for_cells = total_grid_w - (gap_h_val * (config.cols - 1))
    avail_h_for_cells = total_grid_h - (gap_v_val * (config.rows - 1))

    has_crops = len(config.crop_regions) > 0
    expanded_cols = set()
    expanded_rows = set()
    if has_crops:
        expanded_cols = (
            set(range(config.cols))
            if config.crop_cols is None
            else set(config.crop_cols)
        )
        expanded_rows = (
            set(range(config.rows))
            if config.crop_rows is None
            else set(config.crop_rows)
        )
    num_exp_cols = len([c for c in range(config.cols) if c in expanded_cols])
    num_exp_rows = len([r for r in range(config.rows) if r in expanded_rows])

    disp = config.crop_display

    gap_mc = disp.main_crop_gap.to_cm(est_main_w)
    gap_cc = disp.crop_crop_gap.to_cm(est_main_w)
    gap_cb = disp.crop_bottom_gap.to_cm(
        est_main_w if disp.position == "right" else est_main_h
    )

    # Border offset compensation
    border_offset = 0.0
    if config.show_zoom_border:
        border_offset = pt_to_cm(config.zoom_border_width)
        gap_mc += border_offset
        gap_cc += border_offset

    main_w = 0.0
    main_h = 0.0
    crop_box_size = 0.0

    num_crops = max(1, len(config.crop_regions))

    is_fixed_size = (
        config.size_mode == "fixed"
        and config.image_width
        and config.image_width > 0
        and config.image_height
        and config.image_height > 0
    )

    if is_fixed_size:
        main_w = config.image_width
        main_h = config.image_height

        if disp.size is not None:
            crop_box_size = disp.size
        elif disp.scale is not None:
            if disp.position == "right":
                crop_box_size = main_w * disp.scale
            else:
                crop_box_size = main_h * disp.scale
        else:
            if disp.position == "right":
                crop_box_size = main_w * 0.33
            else:
                crop_box_size = main_h * 0.33

    else:
        if disp.position == "right":
            main_h = avail_h_for_cells / config.rows
            crop_k = 0.0
            crop_c = 0.0
            if disp.size is not None:
                crop_c = disp.size
            elif disp.scale is not None:
                crop_k = disp.scale
            else:
                crop_k = 0.33

            # Note: gap_mc includes border offset now
            denom = config.cols + num_exp_cols * crop_k
            numer = avail_w_for_cells - num_exp_cols * (gap_mc + crop_c)
            main_w = numer / denom
            crop_box_size = main_w * crop_k + crop_c

        else:  # bottom
            main_w = avail_w_for_cells / config.cols
            crop_k = 0.0
            crop_c = 0.0
            if disp.size is not None:
                crop_c = disp.size
            elif disp.scale is not None:
                crop_k = disp.scale
            else:
                crop_k = 0.33

            denom = config.rows + num_exp_rows * crop_k
            numer = avail_h_for_cells - num_exp_rows * (gap_mc + crop_c)
            main_h = numer / denom
            crop_box_size = main_h * crop_k + crop_c

    col_widths = []
    row_heights = []

    if disp.position == "right":
        extra = crop_box_size + gap_mc
        for c in range(config.cols):
            if c in expanded_cols:
                col_widths.append(main_w + extra)
            else:
                col_widths.append(main_w)

        row_heights = [main_h] * config.rows
    else:
        # Bottom placement expands row height
        extra = crop_box_size + gap_mc + gap_cb
        col_widths = [main_w] * config.cols
        for r in range(config.rows):
            if r in expanded_rows:
                row_heights.append(main_h + extra)
            else:
                row_heights.append(main_h)

    return LayoutMetrics(
        main_w, main_h, col_widths, row_heights, crop_box_size, gap_mc, gap_cc, gap_cb
    )


# --- Helper: Calculate Item Bounds ---
def calculate_item_bounds(
    config,
    metrics,
    image_path,
    row_idx,
    col_idx,
    border_offset_cm=0.0,
    override_size=None,
):
    """
    Returns (min_x, min_y, max_x, max_y) bounding box relative to Main Image Top-Left (0,0).
    """
    has_crops = should_apply_crop(row_idx, col_idx, config)

    if override_size:
        orig_w, orig_h = override_size
    elif image_path == "dummy":
        # Fallback if no override size is passed for dummy, use default fit
        return 0, 0, 0, 0
    else:
        orig_w, orig_h = calculate_image_size_fit(
            image_path, metrics.main_width, metrics.main_height, config.fit_mode
        )

    min_x, min_y = 0.0, 0.0
    max_x, max_y = orig_w, orig_h

    # Half-border offset for alignment logic
    # We want the "visual edge" (including border) to align with the main image edge.
    # The border is drawn centered on the shape path. So it extends border_width/2 outward.
    half_border = border_offset_cm / 2.0 if config.show_zoom_border else 0.0

    if has_crops and config.crop_regions:
        num_crops = len(config.crop_regions)
        disp = config.crop_display
        actual_gap_mc = disp.main_crop_gap.to_cm(
            orig_w if disp.position == "right" else orig_h
        )
        actual_gap_cc = disp.crop_crop_gap.to_cm(
            orig_w if disp.position == "right" else orig_h
        )
        actual_gap_cb = disp.crop_bottom_gap.to_cm(
            orig_w if disp.position == "right" else orig_h
        )

        if config.show_zoom_border:
            actual_gap_mc += border_offset_cm
            actual_gap_cc += border_offset_cm

        if disp.position == "right":
            # Base start X for crops
            start_x = orig_w + actual_gap_mc

            # We must simulate placement to find bounds
            for crop_idx, region in enumerate(config.crop_regions):
                # Calculate Crop Size
                cw, ch = 0, 0

                # Check for DUMMY context to decide if we use path or dummy assumption
                if override_size:
                    # Dummy dimensions
                    if disp.size is not None:
                        cw, ch = calculate_size_fit_static(
                            100, 100, disp.size, 9999, "width"
                        )
                    elif disp.scale is not None:
                        tw = orig_w * disp.scale
                        cw, ch = calculate_size_fit_static(100, 100, tw, 9999, "width")
                    else:
                        single_h = (
                            orig_h - actual_gap_cc * (num_crops - 1)
                        ) / num_crops
                        cw, ch = calculate_size_fit_static(
                            100, 100, metrics.crop_size, single_h, "fit"
                        )
                else:
                    # Real image path logic
                    if disp.size is not None:
                        cw, ch = calculate_image_size_fit(
                            image_path, disp.size, 9999, "width"
                        )
                    elif disp.scale is not None:
                        tw = orig_w * disp.scale
                        cw, ch = calculate_image_size_fit(image_path, tw, 9999, "width")
                    else:
                        # Fit
                        single_h = (
                            orig_h - actual_gap_cc * (num_crops - 1)
                        ) / num_crops
                        cw, ch = calculate_image_size_fit(
                            image_path, metrics.crop_size, single_h, "fit"
                        )

                # Position logic
                this_gap_mc = region.gap if region.gap is not None else actual_gap_mc
                if region.gap is not None and config.show_zoom_border:
                    this_gap_mc += border_offset_cm

                # Current crop left
                c_left = orig_w + this_gap_mc

                # Y-Alignment logic (relative to main top 0)
                c_top = 0.0
                if region.align == "start":
                    # Align outer visual edge to top
                    c_top = 0.0 + region.offset + half_border
                elif region.align == "center":
                    c_top = (orig_h - ch) / 2 + region.offset
                elif region.align == "end":
                    # Align outer visual edge to bottom
                    c_top = orig_h - ch + region.offset - half_border
                else:  # auto
                    if disp.scale is not None or disp.size is not None:
                        # Stack from top
                        c_top = crop_idx * (ch + actual_gap_cc)
                    else:
                        # Fit logic (Pin ends)
                        if crop_idx == 0:
                            c_top = 0.0
                        elif num_crops > 1 and crop_idx == num_crops - 1:
                            c_top = orig_h - ch
                        else:
                            # Recalc single_h if needed or use approx
                            single_h = (
                                orig_h - actual_gap_cc * (num_crops - 1)
                            ) / num_crops
                            slot_top = crop_idx * (single_h + actual_gap_cc)
                            c_top = slot_top + (single_h - ch) / 2

                # Update bounds
                if c_left < min_x:
                    min_x = c_left
                if c_top < min_y:
                    min_y = c_top
                if (c_left + cw) > max_x:
                    max_x = c_left + cw
                if (c_top + ch) > max_y:
                    max_y = c_top + ch

        else:  # bottom
            # Base start Y for crops
            start_y = orig_h + actual_gap_mc

            for crop_idx, region in enumerate(config.crop_regions):
                cw, ch = 0, 0
                if override_size:
                    if disp.size is not None:
                        cw, ch = calculate_size_fit_static(
                            100, 100, 9999, disp.size, "height"
                        )
                    elif disp.scale is not None:
                        th = orig_h * disp.scale
                        cw, ch = calculate_size_fit_static(100, 100, 9999, th, "height")
                    else:
                        sw = (orig_w - actual_gap_cc * (num_crops - 1)) / num_crops
                        cw, ch = calculate_size_fit_static(
                            100, 100, sw, metrics.crop_size, "fit"
                        )
                else:
                    if disp.size is not None:
                        cw, ch = calculate_image_size_fit(
                            image_path, 9999, disp.size, "height"
                        )
                    elif disp.scale is not None:
                        th = orig_h * disp.scale
                        cw, ch = calculate_image_size_fit(
                            image_path, 9999, th, "height"
                        )
                    else:
                        sw = (orig_w - actual_gap_cc * (num_crops - 1)) / num_crops
                        cw, ch = calculate_image_size_fit(
                            image_path, sw, metrics.crop_size, "fit"
                        )

                this_gap_mc = region.gap if region.gap is not None else actual_gap_mc
                if region.gap is not None and config.show_zoom_border:
                    this_gap_mc += border_offset_cm

                c_top = orig_h + this_gap_mc

                # X-Alignment
                c_left = 0.0
                if region.align == "start":
                    # Align outer visual edge to left
                    c_left = 0.0 + region.offset + half_border
                elif region.align == "center":
                    c_left = (orig_w - cw) / 2 + region.offset
                elif region.align == "end":
                    # Align outer visual edge to right
                    c_left = orig_w - cw + region.offset - half_border
                else:  # auto
                    if disp.scale is not None or disp.size is not None:
                        c_left = crop_idx * (cw + actual_gap_cc)
                    else:
                        if crop_idx == 0:
                            c_left = 0.0
                        elif num_crops > 1 and crop_idx == num_crops - 1:
                            c_left = orig_w - cw
                        else:
                            sw = (orig_w - actual_gap_cc * (num_crops - 1)) / num_crops
                            slot_left = crop_idx * (sw + actual_gap_cc)
                            c_left = slot_left + (sw - cw) / 2

                if c_left < min_x:
                    min_x = c_left
                if c_top < min_y:
                    min_y = c_top
                if (c_left + cw) > max_x:
                    max_x = c_left + cw
                if (c_top + ch) > max_y:
                    max_y = c_top + ch

            # Apply bottom gap to height?
            max_y += actual_gap_cb

    return min_x, min_y, max_x, max_y


def calculate_flow_row_heights(config, metrics, image_grid, border_offset_cm):
    """
    Simulates flow layout to calculate the max content height for each row.
    Used for vertical alignment of rows on the slide.
    """
    row_heights = []

    for row_idx, row_images in enumerate(image_grid):
        if row_idx >= config.rows:
            break

        row_max_h = 0.0

        for col_idx, image_path in enumerate(row_images):
            if col_idx >= config.cols:
                break
            if image_path is None:
                continue

            # Get TRUE bounding box height of the item
            _, _, _, max_y = calculate_item_bounds(
                config, metrics, image_path, row_idx, col_idx, border_offset_cm
            )
            # Re-calling it to get min_y as well
            min_x, min_y, max_x, max_y = calculate_item_bounds(
                config, metrics, image_path, row_idx, col_idx, border_offset_cm
            )
            item_h = max_y - min_y
            if item_h > row_max_h:
                row_max_h = item_h

        # If no items, use metric default
        if row_max_h == 0.0:
            row_max_h = metrics.main_height

        row_heights.append(row_max_h)

    return row_heights


def create_grid_presentation(config: GridConfig) -> str:
    prs = Presentation()
    prs.slide_width = cm_to_emu(config.slide_width)
    prs.slide_height = cm_to_emu(config.slide_height)
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    if config.arrangement == "row":
        image_grid = [get_sorted_images(folder) for folder in config.folders]
    else:
        columns = [get_sorted_images(folder) for folder in config.folders]
        max_len = max(len(col) for col in columns) if columns else 0
        image_grid = [
            [col[i] if i < len(col) else None for col in columns]
            for i in range(max_len)
        ]

    temp_dir = tempfile.mkdtemp()
    num_crops = len(config.crop_regions)
    metrics = calculate_grid_metrics(config)

    # Calculate Border Offset for placement
    border_offset_cm = 0.0
    if config.show_zoom_border:
        border_offset_cm = pt_to_cm(config.crop_border_width)

    # Calculate half-border for alignment logic
    half_border = border_offset_cm / 2.0 if config.show_zoom_border else 0.0

    # 1. Calculate actual row heights for Flow Mode (Vertical Packing)
    flow_row_heights = []
    if config.layout_mode == "flow":
        flow_row_heights = calculate_flow_row_heights(
            config, metrics, image_grid, border_offset_cm
        )

    # 2. Calculate Total Slide Content Height for Flow Vertical Align (Slide-level)
    total_content_height = 0.0
    if config.layout_mode == "flow":
        for i, rh in enumerate(flow_row_heights):
            total_content_height += rh
            if i < len(flow_row_heights) - 1:
                # Add gap between rows
                gap_val = config.gap_v.to_cm(metrics.main_height)
                total_content_height += gap_val
    else:
        # Grid mode height is fixed by definition
        pass

    try:
        # Determine Start Y based on Flow V-Align
        current_y = config.margin_top
        if config.layout_mode == "flow":
            avail_h = config.slide_height - config.margin_top - config.margin_bottom
            if config.flow_vertical_align == "center":
                current_y = config.margin_top + (avail_h - total_content_height) / 2
            elif config.flow_vertical_align == "bottom":
                current_y = (config.margin_top + avail_h) - total_content_height
            else:  # top
                current_y = config.margin_top

        for row_idx, row_images in enumerate(image_grid):
            if row_idx >= config.rows:
                break

            # Determine Row Height for this row
            if config.layout_mode == "flow":
                if row_idx < len(flow_row_heights):
                    current_row_height = flow_row_heights[row_idx]
                else:
                    current_row_height = metrics.main_height
            else:
                current_row_height = (
                    metrics.row_heights[row_idx]
                    if row_idx < len(metrics.row_heights)
                    else metrics.main_height
                )

            this_gap_v = config.gap_v.to_cm(metrics.main_height)
            current_x = config.margin_left

            # --- Pre-calculate row total width for alignment if in flow mode ---
            row_total_content_width = 0.0
            valid_items = 0

            if config.layout_mode == "flow":
                for col_idx, image_path in enumerate(row_images):
                    if col_idx >= config.cols:
                        break
                    if image_path is None:
                        continue

                    min_x, min_y, max_x, max_y = calculate_item_bounds(
                        config, metrics, image_path, row_idx, col_idx, border_offset_cm
                    )
                    item_w = max_x - min_x
                    # item_h = max_y - min_y # Already used for row height

                    row_total_content_width += item_w
                    valid_items += 1

                if valid_items > 1:
                    g_val = config.gap_h.to_cm(metrics.main_width)
                    row_total_content_width += (valid_items - 1) * g_val

                avail_w = config.slide_width - config.margin_left - config.margin_right
                if config.flow_align == "center":
                    current_x = (
                        config.margin_left + (avail_w - row_total_content_width) / 2
                    )
                elif config.flow_align == "right":
                    current_x = config.margin_left + (avail_w - row_total_content_width)
                else:
                    current_x = config.margin_left
            # --- End Pre-calculation ---

            for col_idx, image_path in enumerate(row_images):
                if col_idx >= config.cols:
                    break

                this_gap_h = config.gap_h.to_cm(metrics.main_width)

                if image_path is None:
                    if config.layout_mode == "grid":
                        w = (
                            metrics.col_widths[col_idx]
                            if col_idx < len(metrics.col_widths)
                            else metrics.main_width
                        )
                        current_x += w + this_gap_h
                    continue

                has_crops = should_apply_crop(row_idx, col_idx, config)
                orig_w, orig_h = calculate_image_size_fit(
                    image_path, metrics.main_width, metrics.main_height, config.fit_mode
                )

                # Global dynamic gaps
                global_gap_mc = config.crop_display.main_crop_gap.to_cm(
                    orig_w if config.crop_display.position == "right" else orig_h
                )
                global_gap_cc = config.crop_display.crop_crop_gap.to_cm(
                    orig_w if config.crop_display.position == "right" else orig_h
                )

                # Compensate for border
                if config.show_zoom_border:
                    global_gap_mc += border_offset_cm
                    global_gap_cc += border_offset_cm

                # Determine Main Image Position using Bounding Box Offset
                min_x, min_y, max_x, max_y = calculate_item_bounds(
                    config, metrics, image_path, row_idx, col_idx, border_offset_cm
                )
                item_w = max_x - min_x
                item_h = max_y - min_y

                if config.layout_mode == "flow":
                    # Item left is current_x
                    item_draw_left = current_x

                    # Align Item Vertically WITHIN Row
                    # Row height is current_row_height.
                    # We center the item in the row line.
                    item_draw_top = current_y + (current_row_height - item_h) / 2
                else:
                    cell_w = (
                        metrics.col_widths[col_idx]
                        if col_idx < len(metrics.col_widths)
                        else metrics.main_width
                    )
                    # Center in grid cell
                    item_draw_left = current_x + (cell_w - item_w) / 2
                    # Grid Row Height is in metrics or calculated.
                    item_draw_top = current_y + (current_row_height - item_h) / 2

                # Calculate absolute position for Main Image
                # Main Image is at (0-min_x, 0-min_y) relative to Item Top-Left
                final_main_left = item_draw_left - min_x
                final_main_top = item_draw_top - min_y

                pic = slide.shapes.add_picture(
                    image_path,
                    cm_to_emu(final_main_left),
                    cm_to_emu(final_main_top),
                    cm_to_emu(orig_w),
                    cm_to_emu(orig_h),
                )
                pic.shadow.inherit = False

                if has_crops and config.show_crop_border:
                    add_crop_borders_to_image(
                        slide,
                        final_main_left,
                        final_main_top,
                        orig_w,
                        orig_h,
                        image_path,
                        config.crop_regions,
                        config.crop_border_width,
                    )

                if has_crops and num_crops > 0:
                    disp = config.crop_display

                    for crop_idx, region in enumerate(config.crop_regions):
                        crop_filename = f"crop_{row_idx}_{col_idx}_{crop_idx}.png"
                        crop_path = os.path.join(temp_dir, crop_filename)
                        try:
                            crop_image(image_path, region, crop_path)
                        except Exception:
                            continue

                        # Resolve Size
                        if disp.size is not None:
                            if disp.position == "right":
                                cw, ch = calculate_image_size_fit(
                                    crop_path, disp.size, 9999, "width"
                                )
                            else:
                                cw, ch = calculate_image_size_fit(
                                    crop_path, 9999, disp.size, "height"
                                )
                        elif disp.scale is not None:
                            if disp.position == "right":
                                target_w = orig_w * disp.scale
                                cw, ch = calculate_image_size_fit(
                                    crop_path, target_w, 9999, "width"
                                )
                            else:
                                target_h = orig_h * disp.scale
                                cw, ch = calculate_image_size_fit(
                                    crop_path, 9999, target_h, "height"
                                )
                        else:
                            # Fallback fit
                            if disp.position == "right":
                                single_h = (
                                    orig_h - global_gap_cc * (num_crops - 1)
                                ) / num_crops
                                cw, ch = calculate_image_size_fit(
                                    crop_path, metrics.crop_size, single_h, "fit"
                                )
                            else:
                                single_w = (
                                    orig_w - global_gap_cc * (num_crops - 1)
                                ) / num_crops
                                cw, ch = calculate_image_size_fit(
                                    crop_path, single_w, metrics.crop_size, "fit"
                                )

                        # Resolve Gap
                        this_gap_mc = (
                            region.gap if region.gap is not None else global_gap_mc
                        )
                        if region.gap is not None and config.show_zoom_border:
                            this_gap_mc += border_offset_cm

                        # Resolve Position
                        if disp.position == "right":
                            c_left = final_main_left + orig_w + this_gap_mc
                            if region.align == "start":
                                c_top = final_main_top + region.offset + half_border
                            elif region.align == "center":
                                c_top = (
                                    final_main_top + (orig_h - ch) / 2 + region.offset
                                )
                            elif region.align == "end":
                                c_top = (
                                    final_main_top
                                    + orig_h
                                    - ch
                                    + region.offset
                                    - half_border
                                )
                            else:  # auto
                                if disp.scale is not None or disp.size is not None:
                                    # Stack logic
                                    c_top = final_main_top + crop_idx * (
                                        ch + global_gap_cc
                                    )
                                else:
                                    # Fit logic (Pin ends)
                                    if crop_idx == 0:
                                        c_top = final_main_top
                                    elif num_crops > 1 and crop_idx == num_crops - 1:
                                        c_top = (final_main_top + orig_h) - ch
                                    else:
                                        single_h = (
                                            orig_h - global_gap_cc * (num_crops - 1)
                                        ) / num_crops
                                        slot_top = final_main_top + crop_idx * (
                                            single_h + global_gap_cc
                                        )
                                        c_top = slot_top + (single_h - ch) / 2
                        else:  # bottom
                            c_top = final_main_top + orig_h + this_gap_mc
                            if region.align == "start":
                                c_left = final_main_left + region.offset + half_border
                            elif region.align == "center":
                                c_left = (
                                    final_main_left + (orig_w - cw) / 2 + region.offset
                                )
                            elif region.align == "end":
                                c_left = (
                                    final_main_left
                                    + orig_w
                                    - cw
                                    + region.offset
                                    - half_border
                                )
                            else:  # auto
                                if disp.scale is not None or disp.size is not None:
                                    c_left = final_main_left + crop_idx * (
                                        cw + global_gap_cc
                                    )
                                else:
                                    if crop_idx == 0:
                                        c_left = final_main_left
                                    elif num_crops > 1 and crop_idx == num_crops - 1:
                                        c_left = (final_main_left + orig_w) - cw
                                    else:
                                        single_w = (
                                            orig_w - global_gap_cc * (num_crops - 1)
                                        ) / num_crops
                                        slot_left = final_main_left + crop_idx * (
                                            single_w + global_gap_cc
                                        )
                                        c_left = slot_left + (single_w - cw) / 2

                        pic_crop = slide.shapes.add_picture(
                            crop_path,
                            cm_to_emu(c_left),
                            cm_to_emu(c_top),
                            cm_to_emu(cw),
                            cm_to_emu(ch),
                        )
                        pic_crop.shadow.inherit = False

                        if config.show_zoom_border:
                            add_border_shape(
                                slide,
                                c_left,
                                c_top,
                                cw,
                                ch,
                                region.color,
                                config.zoom_border_width,
                                config.zoom_border_shape,
                            )

                if config.layout_mode == "flow":
                    current_x += item_w + this_gap_h
                else:
                    w = (
                        metrics.col_widths[col_idx]
                        if col_idx < len(metrics.col_widths)
                        else metrics.main_width
                    )
                    current_x += w + this_gap_h

            # Row Advance Logic
            current_y += current_row_height + this_gap_v

        prs.save(config.output)
    except Exception as e:
        raise e
    finally:
        shutil.rmtree(temp_dir)
    return config.output


def generate_sample_config(output_path: str):
    """Generate a sample configuration file."""
    sample = """# PowerPoint Image Grid Generator - Configuration File (v24 Compatible)
# Size values are in centimeters (cm) unless otherwise specified (e.g. pt)

# Slide size settings
slide:
  width: 33.867    # 16:9 aspect ratio
  height: 19.05

# Main grid layout
grid:
  rows: 2
  cols: 3
  arrangement: row  # 'row' or 'col'
  layout_mode: flow # 'grid' (aligned) or 'flow' (compact)
  flow_align: left # 'left', 'center', 'right' (only for flow)
  flow_vertical_align: center # 'top', 'center', 'bottom' (only for flow)

# Margins around the grid (cm)
margin:
  left: 1.0
  top: 1.0
  right: 1.0
  bottom: 1.0

# Gap between images in grid
gap:
  horizontal:
    value: 0.5
    mode: cm # 'cm' or 'scale'
  vertical:
    value: 0.5
    mode: cm

# Image sizing
image:
  size_mode: fit   # 'fit' (auto calc) or 'fixed' (user defined)
  fit_mode: fit    # 'fit', 'width', 'height'
  # width: 10.0    # used if size_mode is fixed
  # height: 7.5
  scale: 1.0

# Crop settings
crop:
  # Multiple crop regions
  regions:
    - name: "Region A"
      x: 50
      y: 50
      width: 100
      height: 100
      color: "#FF0000"
      align: auto   # auto, start, center, end
      offset: 0.0   # cm
      # gap: 0.5    # overrides main_crop_gap
    - name: "Region B"
      x: 200
      y: 100
      width: 120
      height: 80
      color: "#00FF00"
  
  # Which rows/cols to crop (0-indexed), null = all
  rows: [0]
  cols: null
  
  # How to display cropped images next to the original
  display:
    position: right   # 'right' or 'bottom'
    size: null        # Absolute size (cm). If set, takes priority over scale.
    scale: 0.4        # Display size ratio relative to main image
    main_crop_gap: {value: 0.15, mode: cm}
    crop_crop_gap: {value: 0.15, mode: cm}
    crop_bottom_gap: {value: 0.0, mode: cm}

# Border settings
border:
  crop:
    show: true
    width: 1.5  # Line width in points (pt)
  zoom:
    show: true
    width: 1.5  # Line width in points (pt)
    shape: rectangle # 'rectangle' or 'rounded'

# Input folders
folders:
  - "./images/row1"
  - "./images/row2"

# Output file
output: "output.pptx"
"""
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(sample)
    print(f"Sample config created: {output_path}")


def main():
    if len(sys.argv) < 2:
        print("Usage:")
        print(
            "  python image_grid_cli.py <config.yaml>            - Generate presentation"
        )
        print(
            "  python image_grid_cli.py --init [filename]        - Create sample config"
        )
        sys.exit(1)

    if sys.argv[1] == "--init":
        output_name = sys.argv[2] if len(sys.argv) > 2 else "config.yaml"
        generate_sample_config(output_name)
    else:
        config_path = sys.argv[1]
        try:
            config = load_config(config_path)

            # Auto-determine grid size if not specified and folders exist
            if config.arrangement == "row":
                if config.rows == 0 and config.folders:
                    config.rows = len(config.folders)
                if config.cols == 0 and config.folders:
                    try:
                        max_cols = max(
                            len(get_sorted_images(f)) for f in config.folders
                        )
                        config.cols = max_cols
                    except:
                        config.cols = 3
            else:
                if config.cols == 0 and config.folders:
                    config.cols = len(config.folders)
                if config.rows == 0 and config.folders:
                    try:
                        max_rows = max(
                            len(get_sorted_images(f)) for f in config.folders
                        )
                        config.rows = max_rows
                    except:
                        config.rows = 3

            if config.rows == 0:
                config.rows = 1
            if config.cols == 0:
                config.cols = 1

            output = create_grid_presentation(config)
            print(f"Created: {output}")
        except Exception as e:
            print(f"Error: {e}")
            sys.exit(1)


if __name__ == "__main__":
    main()
