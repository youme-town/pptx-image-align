"""
PowerPoint Image Grid Generator - Core Module

This module contains shared logic for creating PowerPoint presentations
with images arranged in a grid layout. It is used by both GUI and CLI interfaces.

Features:
- Data classes for configuration
- Layout calculation algorithms
- PPTX generation functions
- Image processing utilities
"""

import os
import re
import shutil
import tempfile
from pathlib import Path
from typing import Optional, List, Tuple, Dict
from dataclasses import dataclass, field

import yaml
from PIL import Image
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_LINE_DASH_STYLE


# =============================================================================
# Constants
# =============================================================================

CM_TO_EMU = 360000
PT_TO_EMU = 12700

SUPPORTED_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp"}


# =============================================================================
# Unit Conversion Functions
# =============================================================================


def cm_to_emu(cm: float) -> int:
    """Convert centimeters to EMU (English Metric Units)."""
    return int(cm * CM_TO_EMU)


def pt_to_emu(pt: float) -> int:
    """Convert points to EMU."""
    return int(pt * PT_TO_EMU)


def pt_to_cm(pt: float) -> float:
    """Convert points to centimeters."""
    # 1 inch = 2.54 cm = 72 points
    return pt * (2.54 / 72.0)


# =============================================================================
# Data Classes
# =============================================================================


@dataclass
class CropRegion:
    """Defines a rectangular region to crop from source images."""

    x: int
    y: int
    width: int
    height: int
    color: Tuple[int, int, int] = (255, 0, 0)
    name: str = ""
    align: str = "auto"  # 'auto', 'start', 'center', 'end'
    offset: float = 0.0  # cm, offset from alignment anchor
    gap: Optional[float] = None  # cm, overrides global main-crop gap if set

    # NEW: ratio-based crop (robust across different image pixel sizes)
    mode: str = "px"  # 'px' or 'ratio'
    x_ratio: Optional[float] = None
    y_ratio: Optional[float] = None
    width_ratio: Optional[float] = None
    height_ratio: Optional[float] = None

    # 空クロップ: 枠線のみ表示し、拡大画像は表示しない
    show_zoomed: bool = True


@dataclass
class GapConfig:
    """Configuration for gap/spacing with support for absolute or relative values."""

    value: float = 0.5
    mode: str = "cm"  # 'cm' or 'scale'

    def to_cm(self, ref_size: float) -> float:
        """Convert gap value to centimeters based on mode."""
        if self.mode == "scale":
            return ref_size * self.value
        return self.value


@dataclass
class CropDisplayConfig:
    """Configuration for how cropped regions are displayed."""

    position: str = "right"  # 'right' or 'bottom'
    main_crop_gap: GapConfig = field(default_factory=lambda: GapConfig(0.15, "cm"))
    crop_crop_gap: GapConfig = field(default_factory=lambda: GapConfig(0.15, "cm"))
    crop_bottom_gap: GapConfig = field(default_factory=lambda: GapConfig(0.0, "cm"))
    size: Optional[float] = None  # Absolute size in cm
    scale: Optional[float] = None  # Scale relative to main image


@dataclass
class CropOverride:
    """Per-cell crop override (row, col) -> regions. If regions is empty, crop is disabled for that cell."""

    row: int
    col: int
    regions: List[CropRegion] = field(default_factory=list)


@dataclass
class CropPreset:
    """クロッププリセット"""

    name: str
    regions: List[CropRegion] = field(default_factory=list)
    description: str = ""
    display_position: str = "right"  # 'right' or 'bottom'


@dataclass
class LabelConfig:
    """テキストラベル/キャプションの設定"""

    enabled: bool = False
    mode: str = "filename"  # 'filename', 'number', 'custom'
    position: str = "bottom"  # 'top', 'bottom'
    font_name: str = "Arial"
    font_size: float = 10.0  # pt
    font_color: Tuple[int, int, int] = (0, 0, 0)
    font_bold: bool = False
    number_format: str = "({n})"  # 連番フォーマット
    custom_texts: List[str] = field(default_factory=list)  # カスタムモード用
    gap: float = 0.1  # cm - ラベルと画像の間隔


@dataclass
class ConnectorConfig:
    """クロップ領域と拡大画像を結ぶ連結線の設定"""

    show: bool = False
    width: float = 1.0  # pt
    color: Optional[Tuple[int, int, int]] = None  # Noneならクロップ領域の色を使用
    style: str = "straight"  # 'straight' or 'elbow'
    dash_style: str = "solid"  # 'solid', 'dash', 'dot'


@dataclass
class GridConfig:
    """Main configuration for grid-based image layout."""

    # Slide dimensions
    slide_width: float = 33.867
    slide_height: float = 19.05

    # Grid structure
    rows: int = 2
    cols: int = 3
    arrangement: str = "row"  # 'row' or 'col'

    # Margins
    margin_left: float = 1.0
    margin_top: float = 1.0
    margin_right: float = 1.0
    margin_bottom: float = 1.0

    # Gap settings
    gap_h: GapConfig = field(default_factory=lambda: GapConfig(0.5, "cm"))
    gap_v: GapConfig = field(default_factory=lambda: GapConfig(0.5, "cm"))

    # Layout mode
    layout_mode: str = "grid"  # 'grid' (aligned) or 'flow' (compact)
    flow_align: str = "left"  # 'left', 'center', 'right'
    flow_vertical_align: str = "center"  # 'top', 'center', 'bottom'

    # Image sizing
    size_mode: str = "fit"  # 'fit' or 'fixed'
    fit_mode: str = "fit"  # 'fit', 'width', 'height'
    image_width: Optional[float] = None
    image_height: Optional[float] = None
    image_scale: float = 1.0

    # Crop settings
    crop_regions: List[CropRegion] = field(default_factory=list)
    crop_rows: Optional[List[int]] = None
    crop_cols: Optional[List[int]] = None
    crop_targets: Optional[List[Tuple[int, int]]] = None  # NEW: explicit (row, col)

    # NEW: per-cell overrides
    crop_overrides: List[CropOverride] = field(default_factory=list)

    crop_display: CropDisplayConfig = field(default_factory=CropDisplayConfig)

    # Border settings
    show_crop_border: bool = True
    crop_border_width: float = 1.5
    crop_border_shape: str = "rectangle"  # 'rectangle' or 'rounded'
    show_zoom_border: bool = True
    zoom_border_width: float = 1.5
    zoom_border_shape: str = "rectangle"  # 'rectangle' or 'rounded'

    # Input/Output
    folders: List[str] = field(default_factory=list)
    images: Optional[List[str]] = None  # NEW: explicit images (row-major)
    output: str = "output.pptx"

    # Label settings
    label_config: LabelConfig = field(default_factory=LabelConfig)

    # Template settings
    template_path: Optional[str] = None
    slide_layout_index: int = 6  # Default to blank layout

    # Connector settings
    connector: ConnectorConfig = field(default_factory=ConnectorConfig)


@dataclass
class LayoutMetrics:
    """Calculated layout metrics for grid positioning."""

    main_width: float
    main_height: float
    col_widths: List[float]
    row_heights: List[float]
    crop_size: float
    crop_main_gap: float
    crop_crop_gap: float
    crop_bottom_gap: float


# =============================================================================
# Configuration Parsing
# =============================================================================


def parse_color(color_value) -> Tuple[int, int, int]:
    """Parse color from various formats (list, tuple, hex string, comma-separated)."""
    if isinstance(color_value, (list, tuple)):
        return tuple(color_value[:3])
    elif isinstance(color_value, str):
        if color_value.startswith("#"):
            hex_color = color_value.lstrip("#")
            return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))
        else:
            parts = color_value.split(",")
            return tuple(int(p.strip()) for p in parts)
    return (255, 0, 0)


def parse_gap(data) -> GapConfig:
    """Parse gap configuration from YAML data."""
    if isinstance(data, (int, float)):
        return GapConfig(float(data), "cm")
    elif isinstance(data, dict):
        return GapConfig(float(data.get("value", 0.5)), data.get("mode", "cm"))
    return GapConfig(0.5, "cm")


def _parse_crop_region_dict(r: dict, fallback_name: str) -> CropRegion:
    mode = r.get("mode", "px")
    return CropRegion(
        x=r.get("x", 0),
        y=r.get("y", 0),
        width=r.get("width", 100),
        height=r.get("height", 100),
        color=parse_color(r.get("color", "#FF0000")),
        name=r.get("name", fallback_name),
        align=r.get("align", "auto"),
        offset=r.get("offset", 0.0),
        gap=r.get("gap", None),
        mode=mode,
        x_ratio=r.get("x_ratio"),
        y_ratio=r.get("y_ratio"),
        width_ratio=r.get("width_ratio"),
        height_ratio=r.get("height_ratio"),
        show_zoomed=r.get("show_zoomed", True),
    )


def load_config(config_path: str) -> GridConfig:
    """Load configuration from a YAML file."""
    with open(config_path, "r", encoding="utf-8") as f:
        data = yaml.safe_load(f)

    config = GridConfig()

    # Slide settings
    if "slide" in data:
        slide = data["slide"]
        config.slide_width = slide.get("width", config.slide_width)
        config.slide_height = slide.get("height", config.slide_height)

    # Grid settings
    if "grid" in data:
        grid = data["grid"]
        config.rows = grid.get("rows", config.rows)
        config.cols = grid.get("cols", config.cols)
        config.arrangement = grid.get("arrangement", config.arrangement)
        config.layout_mode = grid.get("layout_mode", config.layout_mode)
        config.flow_align = grid.get("flow_align", config.flow_align)
        config.flow_vertical_align = grid.get(
            "flow_vertical_align", config.flow_vertical_align
        )

    # Margin settings
    if "margin" in data:
        margin = data["margin"]
        if isinstance(margin, (int, float)):
            config.margin_left = config.margin_top = margin
            config.margin_right = config.margin_bottom = margin
        else:
            config.margin_left = margin.get("left", config.margin_left)
            config.margin_top = margin.get("top", config.margin_top)
            config.margin_right = margin.get("right", config.margin_right)
            config.margin_bottom = margin.get("bottom", config.margin_bottom)

    # Gap settings
    if "gap" in data:
        gap = data["gap"]
        if isinstance(gap, (int, float)):
            config.gap_h = GapConfig(float(gap), "cm")
            config.gap_v = GapConfig(float(gap), "cm")
        else:
            if "horizontal" in gap:
                config.gap_h = parse_gap(gap["horizontal"])
            if "vertical" in gap:
                config.gap_v = parse_gap(gap["vertical"])

    # Image settings
    if "image" in data:
        img = data["image"]
        config.size_mode = img.get("size_mode", config.size_mode)
        config.fit_mode = img.get("fit_mode", config.fit_mode)
        config.image_scale = img.get("scale", config.image_scale)
        config.image_width = img.get("width")
        config.image_height = img.get("height")

    # Crop settings
    if "crop" in data:
        crop = data["crop"]

        # Parse crop regions (global)
        if "regions" in crop:
            for i, r in enumerate(crop["regions"]):
                config.crop_regions.append(_parse_crop_region_dict(r, f"crop_{i + 1}"))
        elif "region" in crop:
            r = crop["region"]
            config.crop_regions.append(_parse_crop_region_dict(r, "crop_1"))

        config.crop_rows = crop.get("rows")
        config.crop_cols = crop.get("cols")

        # NEW: explicit targets
        targets = crop.get("targets")
        if targets is not None:
            parsed: List[Tuple[int, int]] = []
            for t in targets:
                if isinstance(t, str) and "," in t:
                    rs, cs = t.split(",", 1)
                    parsed.append((int(rs.strip()), int(cs.strip())))
                elif isinstance(t, dict) and "row" in t and "col" in t:
                    parsed.append((int(t["row"]), int(t["col"])))
            config.crop_targets = parsed if parsed else []

        # NEW: per-cell overrides
        overrides = crop.get("overrides")
        if isinstance(overrides, list):
            parsed_overrides: List[CropOverride] = []
            for o in overrides:
                if not isinstance(o, dict):
                    continue
                if "row" not in o or "col" not in o:
                    continue
                regions_data = o.get("regions") or []
                regions: List[CropRegion] = []
                for i, rr in enumerate(regions_data):
                    if isinstance(rr, dict):
                        regions.append(
                            _parse_crop_region_dict(
                                rr, f"cell_{int(o['row'])}_{int(o['col'])}_{i + 1}"
                            )
                        )
                parsed_overrides.append(
                    CropOverride(row=int(o["row"]), col=int(o["col"]), regions=regions)
                )
            config.crop_overrides = parsed_overrides

        # Crop display settings
        if "display" in crop:
            disp = crop["display"]
            config.crop_display.position = disp.get("position", "right")
            config.crop_display.size = disp.get("size")
            config.crop_display.scale = disp.get("scale")

            # Legacy gap support
            legacy_gap = disp.get("gap")
            if legacy_gap is not None:
                if isinstance(legacy_gap, (int, float)):
                    config.crop_display.main_crop_gap = GapConfig(
                        float(legacy_gap), "cm"
                    )
                    config.crop_display.crop_crop_gap = GapConfig(
                        float(legacy_gap), "cm"
                    )

            if "main_crop_gap" in disp:
                config.crop_display.main_crop_gap = parse_gap(disp["main_crop_gap"])
            if "crop_crop_gap" in disp:
                config.crop_display.crop_crop_gap = parse_gap(disp["crop_crop_gap"])
            if "crop_bottom_gap" in disp:
                config.crop_display.crop_bottom_gap = parse_gap(disp["crop_bottom_gap"])

    # Border settings
    if "border" in data:
        border = data["border"]
        if "crop" in border:
            cb = border["crop"]
            config.show_crop_border = cb.get("show", config.show_crop_border)
            config.crop_border_width = cb.get("width", config.crop_border_width)
            config.crop_border_shape = cb.get("shape", config.crop_border_shape)
        if "zoom" in border:
            zb = border["zoom"]
            config.show_zoom_border = zb.get("show", config.show_zoom_border)
            config.zoom_border_width = zb.get("width", config.zoom_border_width)
            config.zoom_border_shape = zb.get("shape", config.zoom_border_shape)

    # Input/Output
    if "folders" in data:
        config.folders = data["folders"]
    config.images = data.get("images")  # NEW
    config.output = data.get("output", config.output)

    # Label settings
    if "label" in data:
        lbl = data["label"]
        config.label_config.enabled = lbl.get("enabled", False)
        config.label_config.mode = lbl.get("mode", "filename")
        config.label_config.position = lbl.get("position", "bottom")
        config.label_config.font_name = lbl.get("font_name", "Arial")
        config.label_config.font_size = lbl.get("font_size", 10.0)
        if "font_color" in lbl:
            config.label_config.font_color = parse_color(lbl["font_color"])
        config.label_config.font_bold = lbl.get("font_bold", False)
        config.label_config.number_format = lbl.get("number_format", "({n})")
        config.label_config.custom_texts = lbl.get("custom_texts", [])
        config.label_config.gap = lbl.get("gap", 0.1)

    # Template settings
    if "template" in data:
        tpl = data["template"]
        config.template_path = tpl.get("path")
        config.slide_layout_index = tpl.get("layout_index", 6)

    # Connector settings
    if "connector" in data:
        conn = data["connector"]
        config.connector.show = conn.get("show", False)
        config.connector.width = conn.get("width", 1.0)
        if "color" in conn and conn["color"]:
            config.connector.color = parse_color(conn["color"])
        config.connector.style = conn.get("style", "straight")
        config.connector.dash_style = conn.get("dash_style", "solid")

    return config


def save_config(config: GridConfig, output_path: str) -> None:
    """Save configuration to a YAML file."""
    data = {
        "slide": {"width": config.slide_width, "height": config.slide_height},
        "grid": {
            "rows": config.rows,
            "cols": config.cols,
            "arrangement": config.arrangement,
            "layout_mode": config.layout_mode,
            "flow_align": config.flow_align,
            "flow_vertical_align": config.flow_vertical_align,
        },
        "margin": {
            "left": config.margin_left,
            "top": config.margin_top,
            "right": config.margin_right,
            "bottom": config.margin_bottom,
        },
        "gap": {
            "horizontal": {"value": config.gap_h.value, "mode": config.gap_h.mode},
            "vertical": {"value": config.gap_v.value, "mode": config.gap_v.mode},
        },
        "image": {
            "size_mode": config.size_mode,
            "fit_mode": config.fit_mode,
            "width": config.image_width,
            "height": config.image_height,
        },
        "folders": config.folders,
        "images": config.images,  # NEW
        "output": config.output,
        "label": {
            "enabled": config.label_config.enabled,
            "mode": config.label_config.mode,
            "position": config.label_config.position,
            "font_name": config.label_config.font_name,
            "font_size": config.label_config.font_size,
            "font_color": list(config.label_config.font_color),
            "font_bold": config.label_config.font_bold,
            "number_format": config.label_config.number_format,
            "custom_texts": config.label_config.custom_texts,
            "gap": config.label_config.gap,
        },
        "template": {
            "path": config.template_path,
            "layout_index": config.slide_layout_index,
        },
        "connector": {
            "show": config.connector.show,
            "width": config.connector.width,
            "color": list(config.connector.color) if config.connector.color else None,
            "style": config.connector.style,
            "dash_style": config.connector.dash_style,
        },
        "border": {
            "crop": {
                "show": config.show_crop_border,
                "width": config.crop_border_width,
                "shape": config.crop_border_shape,
            },
            "zoom": {
                "show": config.show_zoom_border,
                "width": config.zoom_border_width,
                "shape": config.zoom_border_shape,
            },
        },
        "crop": {
            "regions": [
                {
                    "name": r.name,
                    "x": r.x,
                    "y": r.y,
                    "width": r.width,
                    "height": r.height,
                    "color": list(r.color),
                    "align": r.align,
                    "offset": r.offset,
                    "gap": r.gap,
                    "mode": r.mode,
                    "x_ratio": r.x_ratio,
                    "y_ratio": r.y_ratio,
                    "width_ratio": r.width_ratio,
                    "height_ratio": r.height_ratio,
                    "show_zoomed": r.show_zoomed,
                }
                for r in config.crop_regions
            ],
            "rows": config.crop_rows,
            "cols": config.crop_cols,
            "targets": (
                [{"row": r, "col": c} for (r, c) in config.crop_targets]
                if config.crop_targets is not None
                else None
            ),
            # NEW
            "overrides": [
                {
                    "row": o.row,
                    "col": o.col,
                    "regions": [
                        {
                            "name": r.name,
                            "x": r.x,
                            "y": r.y,
                            "width": r.width,
                            "height": r.height,
                            "color": list(r.color),
                            "align": r.align,
                            "offset": r.offset,
                            "gap": r.gap,
                            "mode": r.mode,
                            "x_ratio": r.x_ratio,
                            "y_ratio": r.y_ratio,
                            "width_ratio": r.width_ratio,
                            "height_ratio": r.height_ratio,
                            "show_zoomed": r.show_zoomed,
                        }
                        for r in (o.regions or [])
                    ],
                }
                for o in (config.crop_overrides or [])
            ],
            "display": {
                "position": config.crop_display.position,
                "size": config.crop_display.size,
                "scale": config.crop_display.scale,
                "main_crop_gap": {
                    "value": config.crop_display.main_crop_gap.value,
                    "mode": config.crop_display.main_crop_gap.mode,
                },
                "crop_crop_gap": {
                    "value": config.crop_display.crop_crop_gap.value,
                    "mode": config.crop_display.crop_crop_gap.mode,
                },
                "crop_bottom_gap": {
                    "value": config.crop_display.crop_bottom_gap.value,
                    "mode": config.crop_display.crop_bottom_gap.mode,
                },
            },
        },
    }
    with open(output_path, "w", encoding="utf-8") as f:
        yaml.dump(data, f, allow_unicode=True)


# =============================================================================
# Image Utilities
# =============================================================================


def extract_number_from_filename(filename: str) -> int:
    """Extract the first number from a filename for sorting."""
    numbers = re.findall(r"\d+", filename)
    return int(numbers[0]) if numbers else 0


def get_sorted_images(folder_path: str) -> List[str]:
    """Get sorted list of image files from a folder."""
    p = Path(folder_path)
    if not p.exists():
        return []

    # NEW: allow a single image path
    if p.is_file():
        return [str(p)] if p.suffix.lower() in SUPPORTED_IMAGE_EXTENSIONS else []

    image_files = [
        f
        for f in p.iterdir()
        if f.is_file() and f.suffix.lower() in SUPPORTED_IMAGE_EXTENSIONS
    ]
    image_files.sort(key=lambda f: extract_number_from_filename(f.stem))
    return [str(f) for f in image_files]


def build_image_grid(config: GridConfig) -> List[List[Optional[str]]]:
    """
    Build a unified rows x cols image grid.
    Priority:
      1) config.images (row-major)
      2) config.folders interpreted by config.arrangement ('row' or 'col')
    Missing cells are None.
    """
    rows = max(1, int(config.rows))
    cols = max(1, int(config.cols))
    grid: List[List[Optional[str]]] = [[None for _ in range(cols)] for _ in range(rows)]

    # 1) Explicit images (row-major)
    if config.images:
        for idx, p in enumerate(config.images):
            r = idx // cols
            c = idx % cols
            if 0 <= r < rows and 0 <= c < cols:
                grid[r][c] = p
        return grid

    # 2) Folders (row/col)
    folders = config.folders or []
    if config.arrangement == "col":
        # folders[col] -> images down rows
        for c in range(min(cols, len(folders))):
            imgs = get_sorted_images(folders[c])
            for r in range(min(rows, len(imgs))):
                grid[r][c] = imgs[r]
    else:
        # folders[row] -> images across cols
        for r in range(min(rows, len(folders))):
            imgs = get_sorted_images(folders[r])
            for c in range(min(cols, len(imgs))):
                grid[r][c] = imgs[c]

    return grid


def _clamp_int(v: int, lo: int, hi: int) -> int:
    return max(lo, min(hi, v))


def resolve_crop_box(
    region: CropRegion, img_w_px: int, img_h_px: int
) -> Tuple[int, int, int, int]:
    """
    Resolve CropRegion to pixel box (x, y, w, h).
    Supports region.mode == 'ratio'.
    """
    if region.mode == "ratio":
        xr = 0.0 if region.x_ratio is None else float(region.x_ratio)
        yr = 0.0 if region.y_ratio is None else float(region.y_ratio)
        wr = 0.0 if region.width_ratio is None else float(region.width_ratio)
        hr = 0.0 if region.height_ratio is None else float(region.height_ratio)

        x = int(round(xr * img_w_px))
        y = int(round(yr * img_h_px))
        w = int(round(wr * img_w_px))
        h = int(round(hr * img_h_px))
    else:
        x, y, w, h = region.x, region.y, region.width, region.height

    # clamp
    x = _clamp_int(x, 0, max(0, img_w_px))
    y = _clamp_int(y, 0, max(0, img_h_px))
    w = _clamp_int(w, 0, max(0, img_w_px - x))
    h = _clamp_int(h, 0, max(0, img_h_px - y))
    return x, y, w, h


def crop_image(image_path: str, region: CropRegion, output_path: str) -> str:
    """Crop an image according to a CropRegion and save to output_path."""
    with Image.open(image_path) as img:
        img_w, img_h = img.size
        x, y, w, h = resolve_crop_box(region, img_w, img_h)
        box = (x, y, x + w, y + h)
        cropped = img.crop(box)
        cropped.save(output_path)
    return output_path


def get_image_dimensions(image_path: str) -> Tuple[int, int]:
    """Get width and height of an image in pixels."""
    with Image.open(image_path) as img:
        return img.size


# =============================================================================
# Layout Calculation
# =============================================================================


def get_crop_regions_for_cell(
    row: int, col: int, config: GridConfig
) -> List[CropRegion]:
    """Return crop regions for a specific cell. Per-cell override (if present) takes priority."""
    for o in config.crop_overrides or []:
        if o.row == row and o.col == col:
            return o.regions or []
    return config.crop_regions or []


def should_apply_crop(row: int, col: int, config: GridConfig) -> bool:
    """Determine if crop regions should be applied to a specific cell."""
    # Per-cell override takes absolute priority (including disabling by empty list)
    for o in config.crop_overrides or []:
        if o.row == row and o.col == col:
            return bool(o.regions)

    if not config.crop_regions:
        return False

    # explicit targets take priority
    if config.crop_targets is not None:
        return (row, col) in set(config.crop_targets)

    if config.crop_rows is not None and row not in config.crop_rows:
        return False
    if config.crop_cols is not None and col not in config.crop_cols:
        return False
    return True


def calculate_size_fit_static(
    img_w: int, img_h: int, max_width: float, max_height: float, fit_mode: str = "fit"
) -> Tuple[float, float]:
    """Calculate image size to fit within constraints (static version)."""
    if max_height <= 0 or max_width <= 0 or img_h == 0:
        return 0, 0

    aspect_ratio = img_w / img_h

    if fit_mode == "width":
        width = max_width
        height = width / aspect_ratio
    elif fit_mode == "height":
        height = max_height
        width = height * aspect_ratio
    else:  # 'fit'
        cell_aspect = max_width / max_height
        if aspect_ratio > cell_aspect:
            width = max_width
            height = width / aspect_ratio
        else:
            height = max_height
            width = height * aspect_ratio

    return width, height


def calculate_image_size_fit(
    image_path: str, max_width: float, max_height: float, fit_mode: str = "fit"
) -> Tuple[float, float]:
    """Calculate image size to fit within constraints (from file)."""
    try:
        img_width_px, img_height_px = get_image_dimensions(image_path)
    except Exception:
        return 0, 0
    return calculate_size_fit_static(
        img_width_px, img_height_px, max_width, max_height, fit_mode
    )


def calculate_grid_metrics(config: GridConfig) -> LayoutMetrics:
    """Calculate layout metrics for grid positioning."""
    total_grid_w = config.slide_width - config.margin_left - config.margin_right
    total_grid_h = config.slide_height - config.margin_top - config.margin_bottom

    est_cols = max(1, config.cols)
    est_rows = max(1, config.rows)
    est_main_w = total_grid_w / est_cols
    est_main_h = total_grid_h / est_rows

    gap_h_val = config.gap_h.to_cm(est_main_w)
    gap_v_val = config.gap_v.to_cm(est_main_h)

    avail_w_for_cells = total_grid_w - (gap_h_val * (config.cols - 1))
    avail_h_for_cells = total_grid_h - (gap_v_val * (config.rows - 1))

    # Determine which cells have crops (global or per-cell)
    expanded_cols: set[int] = set()
    expanded_rows: set[int] = set()
    for r in range(max(1, config.rows)):
        for c in range(max(1, config.cols)):
            if should_apply_crop(r, c, config):
                expanded_cols.add(c)
                expanded_rows.add(r)

    num_exp_cols = len(expanded_cols)
    num_exp_rows = len(expanded_rows)

    disp = config.crop_display
    gap_mc = disp.main_crop_gap.to_cm(est_main_w)
    gap_cc = disp.crop_crop_gap.to_cm(est_main_w)
    gap_cb = disp.crop_bottom_gap.to_cm(
        est_main_w if disp.position == "right" else est_main_h
    )

    # Border offset compensation
    border_offset = 0.0
    if config.show_zoom_border:
        border_offset = pt_to_cm(config.zoom_border_width)
        gap_mc += border_offset
        gap_cc += border_offset

    main_w = 0.0
    main_h = 0.0
    crop_box_size = 0.0

    is_fixed_size = (
        config.size_mode == "fixed"
        and config.image_width
        and config.image_width > 0
        and config.image_height
        and config.image_height > 0
    )

    if is_fixed_size:
        main_w = config.image_width
        main_h = config.image_height

        if disp.size is not None:
            crop_box_size = disp.size
        elif disp.scale is not None:
            crop_box_size = (
                main_w * disp.scale if disp.position == "right" else main_h * disp.scale
            )
        else:
            crop_box_size = main_w * 0.33 if disp.position == "right" else main_h * 0.33
    else:
        if disp.position == "right":
            main_h = avail_h_for_cells / config.rows
            crop_k = disp.scale if disp.scale is not None else 0.33
            crop_c = disp.size if disp.size is not None else 0.0
            if disp.size is not None:
                crop_k = 0.0

            denom = config.cols + num_exp_cols * crop_k
            numer = avail_w_for_cells - num_exp_cols * (gap_mc + crop_c)
            main_w = numer / denom
            crop_box_size = main_w * crop_k + crop_c
        else:  # bottom
            main_w = avail_w_for_cells / config.cols
            crop_k = disp.scale if disp.scale is not None else 0.33
            crop_c = disp.size if disp.size is not None else 0.0
            if disp.size is not None:
                crop_k = 0.0

            denom = config.rows + num_exp_rows * crop_k
            numer = avail_h_for_cells - num_exp_rows * (gap_mc + crop_c)
            main_h = numer / denom
            crop_box_size = main_h * crop_k + crop_c

    # Calculate column widths and row heights
    col_widths = []
    row_heights = []

    if disp.position == "right":
        extra = crop_box_size + gap_mc
        for c in range(config.cols):
            col_widths.append(main_w + extra if c in expanded_cols else main_w)
        row_heights = [main_h] * config.rows
    else:  # bottom
        extra = crop_box_size + gap_mc + gap_cb
        col_widths = [main_w] * config.cols
        for r in range(config.rows):
            row_heights.append(main_h + extra if r in expanded_rows else main_h)

    return LayoutMetrics(
        main_w, main_h, col_widths, row_heights, crop_box_size, gap_mc, gap_cc, gap_cb
    )


def calculate_item_bounds(
    config: GridConfig,
    metrics: LayoutMetrics,
    image_path: str,
    row_idx: int,
    col_idx: int,
    border_offset_cm: float = 0.0,
    override_size: Optional[Tuple[float, float]] = None,
) -> Tuple[float, float, float, float]:
    """
    Calculate bounding box of an item (main image + crops) relative to main image top-left (0,0).
    Returns (min_x, min_y, max_x, max_y).
    """
    crop_regions = get_crop_regions_for_cell(row_idx, col_idx, config)
    has_crops = bool(crop_regions) and should_apply_crop(row_idx, col_idx, config)
    # show_zoomed=False の領域は拡大画像を描画しないため、bounds（auto配置）には含めない
    zoom_regions = [r for r in crop_regions if getattr(r, "show_zoomed", True)]

    if override_size:
        orig_w, orig_h = override_size
    elif image_path == "dummy":
        return 0, 0, 0, 0
    else:
        orig_w, orig_h = calculate_image_size_fit(
            image_path, metrics.main_width, metrics.main_height, config.fit_mode
        )

    min_x, min_y = 0.0, 0.0
    max_x, max_y = orig_w, orig_h

    # PowerPointの線は図形境界の外側に「半分」出るため、配置計算上は外形を膨らませる
    half_zoom_border = border_offset_cm / 2.0 if config.show_zoom_border else 0.0

    # 元画像上のクロップ枠線も同様に外側へ半分出る（端揃え時のはみ出し対策）
    if has_crops and config.show_crop_border:
        half_crop_border = pt_to_cm(config.crop_border_width) / 2.0
        min_x -= half_crop_border
        min_y -= half_crop_border
        max_x += half_crop_border
        max_y += half_crop_border

    if has_crops and zoom_regions:
        num_crops = len(zoom_regions)
        disp = config.crop_display
        actual_gap_mc = disp.main_crop_gap.to_cm(
            orig_w if disp.position == "right" else orig_h
        )
        actual_gap_cc = disp.crop_crop_gap.to_cm(
            orig_w if disp.position == "right" else orig_h
        )
        actual_gap_cb = disp.crop_bottom_gap.to_cm(
            orig_w if disp.position == "right" else orig_h
        )

        if config.show_zoom_border:
            actual_gap_mc += border_offset_cm
            actual_gap_cc += border_offset_cm

        for crop_idx, region in enumerate(zoom_regions):
            cw, ch = _calculate_crop_size(
                config,
                metrics,
                disp,
                image_path,
                orig_w,
                orig_h,
                num_crops,
                actual_gap_cc,
                override_size,
            )

            this_gap_mc = region.gap if region.gap is not None else actual_gap_mc
            if region.gap is not None and config.show_zoom_border:
                this_gap_mc += border_offset_cm

            if disp.position == "right":
                c_left = orig_w + this_gap_mc
                c_top = _calculate_crop_vertical_position(
                    region,
                    crop_idx,
                    num_crops,
                    orig_h,
                    ch,
                    actual_gap_cc,
                    half_zoom_border,
                    disp,
                )
            else:
                c_top = orig_h + this_gap_mc
                c_left = _calculate_crop_horizontal_position(
                    region,
                    crop_idx,
                    num_crops,
                    orig_w,
                    cw,
                    actual_gap_cc,
                    half_zoom_border,
                    disp,
                )

            # 拡大画像の枠線（zoom border）が外側へ半分出るのを反映
            min_x = min(min_x, c_left - half_zoom_border)
            min_y = min(min_y, c_top - half_zoom_border)
            max_x = max(max_x, c_left + cw + half_zoom_border)
            max_y = max(max_y, c_top + ch + half_zoom_border)

        if disp.position == "bottom":
            max_y += actual_gap_cb

    return min_x, min_y, max_x, max_y


def _calculate_crop_size(
    config: GridConfig,
    metrics: LayoutMetrics,
    disp: CropDisplayConfig,
    image_path: str,
    orig_w: float,
    orig_h: float,
    num_crops: int,
    actual_gap_cc: float,
    override_size: Optional[Tuple[float, float]],
) -> Tuple[float, float]:
    """Calculate the size of a crop region display."""
    if override_size:
        # Dummy/preview mode
        if disp.size is not None:
            if disp.position == "right":
                return calculate_size_fit_static(100, 100, disp.size, 9999, "width")
            else:
                return calculate_size_fit_static(100, 100, 9999, disp.size, "height")
        elif disp.scale is not None:
            if disp.position == "right":
                tw = orig_w * disp.scale
                return calculate_size_fit_static(100, 100, tw, 9999, "width")
            else:
                th = orig_h * disp.scale
                return calculate_size_fit_static(100, 100, 9999, th, "height")
        else:
            if disp.position == "right":
                single_h = (orig_h - actual_gap_cc * (num_crops - 1)) / num_crops
                return calculate_size_fit_static(
                    100, 100, metrics.crop_size, single_h, "fit"
                )
            else:
                sw = (orig_w - actual_gap_cc * (num_crops - 1)) / num_crops
                return calculate_size_fit_static(100, 100, sw, metrics.crop_size, "fit")
    else:
        # Real image mode
        if disp.size is not None:
            if disp.position == "right":
                return calculate_image_size_fit(image_path, disp.size, 9999, "width")
            else:
                return calculate_image_size_fit(image_path, 9999, disp.size, "height")
        elif disp.scale is not None:
            if disp.position == "right":
                tw = orig_w * disp.scale
                return calculate_image_size_fit(image_path, tw, 9999, "width")
            else:
                th = orig_h * disp.scale
                return calculate_image_size_fit(image_path, 9999, th, "height")
        else:
            if disp.position == "right":
                single_h = (orig_h - actual_gap_cc * (num_crops - 1)) / num_crops
                return calculate_image_size_fit(
                    image_path, metrics.crop_size, single_h, "fit"
                )
            else:
                sw = (orig_w - actual_gap_cc * (num_crops - 1)) / num_crops
                return calculate_image_size_fit(
                    image_path, sw, metrics.crop_size, "fit"
                )


def _calculate_crop_vertical_position(
    region: CropRegion,
    crop_idx: int,
    num_crops: int,
    orig_h: float,
    ch: float,
    actual_gap_cc: float,
    half_border: float,
    disp: CropDisplayConfig,
) -> float:
    """Calculate vertical position for a crop (used when position='right')."""
    if region.align == "start":
        return region.offset + half_border
    elif region.align == "center":
        return (orig_h - ch) / 2 + region.offset
    elif region.align == "end":
        return orig_h - ch + region.offset - half_border
    else:  # auto
        # 最初と最後のクロップは端に揃える（pin ends）
        if crop_idx == 0:
            return 0.0 + half_border
        elif num_crops > 1 and crop_idx == num_crops - 1:
            return orig_h - ch - half_border
        else:
            # 中間のクロップは均等配置
            if disp.scale is not None or disp.size is not None:
                # scale/size指定時: 等間隔で配置
                total_crop_h = num_crops * ch + (num_crops - 1) * actual_gap_cc
                start_y = (orig_h - total_crop_h) / 2
                return start_y + crop_idx * (ch + actual_gap_cc)
            else:
                # fit時: スロットの中央に配置
                single_h = (orig_h - actual_gap_cc * (num_crops - 1)) / num_crops
                slot_top = crop_idx * (single_h + actual_gap_cc)
                return slot_top + (single_h - ch) / 2


def _calculate_crop_horizontal_position(
    region: CropRegion,
    crop_idx: int,
    num_crops: int,
    orig_w: float,
    cw: float,
    actual_gap_cc: float,
    half_border: float,
    disp: CropDisplayConfig,
) -> float:
    """Calculate horizontal position for a crop (used when position='bottom')."""
    if region.align == "start":
        return region.offset + half_border
    elif region.align == "center":
        return (orig_w - cw) / 2 + region.offset
    elif region.align == "end":
        return orig_w - cw + region.offset - half_border
    else:  # auto
        # 最初と最後のクロップは端に揃える（pin ends）
        if crop_idx == 0:
            return 0.0 + half_border
        elif num_crops > 1 and crop_idx == num_crops - 1:
            return orig_w - cw - half_border
        else:
            # 中間のクロップは均等配置
            if disp.scale is not None or disp.size is not None:
                # scale/size指定時: 等間隔で配置
                total_crop_w = num_crops * cw + (num_crops - 1) * actual_gap_cc
                start_x = (orig_w - total_crop_w) / 2
                return start_x + crop_idx * (cw + actual_gap_cc)
            else:
                # fit時: スロットの中央に配置
                single_w = (orig_w - actual_gap_cc * (num_crops - 1)) / num_crops
                slot_left = crop_idx * (single_w + actual_gap_cc)
                return slot_left + (single_w - cw) / 2


def calculate_flow_row_heights(
    config: GridConfig,
    metrics: LayoutMetrics,
    image_grid: List[List[Optional[str]]],
    border_offset_cm: float,
) -> List[float]:
    """Calculate maximum content height for each row in flow mode."""
    row_heights = []

    for row_idx, row_images in enumerate(image_grid):
        if row_idx >= config.rows:
            break

        row_max_h = 0.0
        for col_idx, image_path in enumerate(row_images):
            if (
                col_idx >= config.cols
                or image_path is None
                or image_path == "__PLACEHOLDER__"
            ):
                continue

            min_x, min_y, max_x, max_y = calculate_item_bounds(
                config, metrics, image_path, row_idx, col_idx, border_offset_cm
            )
            item_h = max_y - min_y
            row_max_h = max(row_max_h, item_h)

        row_heights.append(row_max_h if row_max_h > 0 else metrics.main_height)

    return row_heights


def _placeholder_override_size(
    config: GridConfig, metrics: LayoutMetrics
) -> Tuple[float, float]:
    """空セル(プレースホルダー)用の仮サイズ(cm)。

    flow レイアウトで空セルを「セル1つ分」として幅・高さ計算に含めるために使う。
    """
    if (
        config.size_mode == "fixed"
        and config.image_width is not None
        and config.image_height is not None
        and config.image_width > 0
        and config.image_height > 0
    ):
        return float(config.image_width), float(config.image_height)

    return float(metrics.main_width), float(metrics.main_height)


# =============================================================================
# PPTX Generation
# =============================================================================


def add_border_shape(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    border_color: Tuple[int, int, int],
    border_width: float,
    shape_type: str = "rectangle",
) -> None:
    """Add a border shape to a slide."""
    ms_shape_type = (
        MSO_SHAPE.ROUNDED_RECTANGLE if shape_type == "rounded" else MSO_SHAPE.RECTANGLE
    )

    shape = slide.shapes.add_shape(
        ms_shape_type,
        cm_to_emu(left),
        cm_to_emu(top),
        cm_to_emu(width),
        cm_to_emu(height),
    )
    shape.fill.background()
    shape.line.color.rgb = RGBColor(*border_color)
    shape.line.width = pt_to_emu(border_width)
    shape.shadow.inherit = False


def add_crop_borders_to_image(
    slide,
    image_left: float,
    image_top: float,
    image_width: float,
    image_height: float,
    original_image_path: str,
    crop_regions: List[CropRegion],
    border_width: float,
    border_shape: str = "rectangle",
) -> None:
    """Add crop region borders overlaid on the main image."""
    try:
        orig_width_px, orig_height_px = get_image_dimensions(original_image_path)
    except Exception:
        return

    scale_x = image_width / orig_width_px
    scale_y = image_height / orig_height_px

    for region in crop_regions:
        x, y, w, h = resolve_crop_box(region, orig_width_px, orig_height_px)
        border_left = image_left + x * scale_x
        border_top = image_top + y * scale_y
        border_w = w * scale_x
        border_h = h * scale_y
        add_border_shape(
            slide,
            border_left,
            border_top,
            border_w,
            border_h,
            region.color,
            border_width,
            border_shape,
        )


def create_grid_presentation(config: GridConfig) -> str:
    """Create a PowerPoint presentation with images arranged in a grid layout."""
    # Use template if specified, otherwise create new presentation
    if config.template_path and Path(config.template_path).exists():
        prs = Presentation(config.template_path)
    else:
        prs = Presentation()
        prs.slide_width = cm_to_emu(config.slide_width)
        prs.slide_height = cm_to_emu(config.slide_height)

    # Get slide layout
    layout_index = min(config.slide_layout_index, len(prs.slide_layouts) - 1)
    slide_layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(slide_layout)

    # NEW: unified grid builder (rows x cols)
    image_grid = build_image_grid(config)

    temp_dir = tempfile.mkdtemp()
    metrics = calculate_grid_metrics(config)

    border_offset_cm = (
        pt_to_cm(config.zoom_border_width) if config.show_zoom_border else 0.0
    )
    half_border = border_offset_cm / 2.0 if config.show_zoom_border else 0.0

    # Calculate row heights for flow mode
    flow_row_heights = []
    total_content_height = 0.0
    if config.layout_mode == "flow":
        flow_row_heights = calculate_flow_row_heights(
            config, metrics, image_grid, border_offset_cm
        )
        for i, rh in enumerate(flow_row_heights):
            total_content_height += rh
            if i < len(flow_row_heights) - 1:
                total_content_height += config.gap_v.to_cm(metrics.main_height)

    try:
        # Determine start Y based on flow vertical alignment
        current_y = config.margin_top
        if config.layout_mode == "flow":
            avail_h = config.slide_height - config.margin_top - config.margin_bottom
            if config.flow_vertical_align == "center":
                current_y = config.margin_top + (avail_h - total_content_height) / 2
            elif config.flow_vertical_align == "bottom":
                current_y = (config.margin_top + avail_h) - total_content_height

        for row_idx in range(config.rows):
            # Determine row height
            if config.layout_mode == "flow":
                current_row_height = (
                    flow_row_heights[row_idx]
                    if row_idx < len(flow_row_heights)
                    else metrics.main_height
                )
            else:
                current_row_height = (
                    metrics.row_heights[row_idx]
                    if row_idx < len(metrics.row_heights)
                    else metrics.main_height
                )

            this_gap_v = config.gap_v.to_cm(metrics.main_height)
            current_x = config.margin_left

            # Pre-calculate row width for flow alignment
            if config.layout_mode == "flow":
                row_total_content_width = 0.0
                valid_items = 0

                for col_idx in range(config.cols):
                    image_path = image_grid[row_idx][col_idx]

                    # 空セルも「セル1つ分」として幅計算に含める
                    if image_path is None or image_path == "__PLACEHOLDER__":
                        ow, oh = _placeholder_override_size(config, metrics)
                        min_x, min_y, max_x, max_y = calculate_item_bounds(
                            config,
                            metrics,
                            "dummy",
                            row_idx,
                            col_idx,
                            border_offset_cm,
                            override_size=(ow, oh),
                        )
                    else:
                        min_x, min_y, max_x, max_y = calculate_item_bounds(
                            config,
                            metrics,
                            image_path,
                            row_idx,
                            col_idx,
                            border_offset_cm,
                        )

                    row_total_content_width += max_x - min_x
                    valid_items += 1

                if valid_items > 1:
                    row_total_content_width += (valid_items - 1) * config.gap_h.to_cm(
                        metrics.main_width
                    )

                avail_w = config.slide_width - config.margin_left - config.margin_right
                if config.flow_align == "center":
                    current_x = (
                        config.margin_left + (avail_w - row_total_content_width) / 2
                    )
                elif config.flow_align == "right":
                    current_x = config.margin_left + (avail_w - row_total_content_width)

            # Process each cell
            for col_idx in range(config.cols):
                image_path = image_grid[row_idx][col_idx]
                this_gap_h = config.gap_h.to_cm(metrics.main_width)

                # Skip None or placeholder cells (ただし flow では幅だけは確保する)
                if image_path is None or image_path == "__PLACEHOLDER__":
                    if config.layout_mode == "grid":
                        w = (
                            metrics.col_widths[col_idx]
                            if col_idx < len(metrics.col_widths)
                            else metrics.main_width
                        )
                        current_x += w + this_gap_h
                    else:  # flow
                        ow, oh = _placeholder_override_size(config, metrics)
                        min_x, min_y, max_x, max_y = calculate_item_bounds(
                            config,
                            metrics,
                            "dummy",
                            row_idx,
                            col_idx,
                            border_offset_cm,
                            override_size=(ow, oh),
                        )
                        item_w = max_x - min_x
                        current_x += item_w + this_gap_h
                    continue

                crop_regions = get_crop_regions_for_cell(row_idx, col_idx, config)
                has_crops = bool(crop_regions) and should_apply_crop(
                    row_idx, col_idx, config
                )

                orig_w, orig_h = calculate_image_size_fit(
                    image_path, metrics.main_width, metrics.main_height, config.fit_mode
                )

                # Calculate dynamic gaps
                global_gap_mc = config.crop_display.main_crop_gap.to_cm(
                    orig_w if config.crop_display.position == "right" else orig_h
                )
                global_gap_cc = config.crop_display.crop_crop_gap.to_cm(
                    orig_w if config.crop_display.position == "right" else orig_h
                )

                if config.show_zoom_border:
                    global_gap_mc += border_offset_cm
                    global_gap_cc += border_offset_cm

                # Calculate item bounds and position
                min_x, min_y, max_x, max_y = calculate_item_bounds(
                    config, metrics, image_path, row_idx, col_idx, border_offset_cm
                )
                item_w = max_x - min_x
                item_h = max_y - min_y

                if config.layout_mode == "flow":
                    item_draw_left = current_x
                    item_draw_top = current_y + (current_row_height - item_h) / 2
                else:
                    cell_w = (
                        metrics.col_widths[col_idx]
                        if col_idx < len(metrics.col_widths)
                        else metrics.main_width
                    )
                    item_draw_left = current_x + (cell_w - item_w) / 2
                    item_draw_top = current_y + (current_row_height - item_h) / 2

                final_main_left = item_draw_left - min_x
                final_main_top = item_draw_top - min_y

                # Add main image
                pic = slide.shapes.add_picture(
                    image_path,
                    cm_to_emu(final_main_left),
                    cm_to_emu(final_main_top),
                    cm_to_emu(orig_w),
                    cm_to_emu(orig_h),
                )
                pic.shadow.inherit = False

                if has_crops and config.show_crop_border:
                    add_crop_borders_to_image(
                        slide,
                        final_main_left,
                        final_main_top,
                        orig_w,
                        orig_h,
                        image_path,
                        crop_regions,
                        config.crop_border_width,
                        config.crop_border_shape,
                    )

                if has_crops and len(crop_regions) > 0:
                    _add_crop_images(
                        slide,
                        config,
                        metrics,
                        image_path,
                        temp_dir,
                        row_idx,
                        col_idx,
                        crop_regions,
                        final_main_left,
                        final_main_top,
                        orig_w,
                        orig_h,
                        global_gap_mc,
                        global_gap_cc,
                        border_offset_cm,
                        half_border,
                    )

                # Add text label if enabled
                if config.label_config.enabled:
                    label_text = generate_label_text(
                        config, image_path, row_idx, col_idx
                    )
                    if label_text:
                        label_h = calculate_label_height(config.label_config)
                        label_gap = config.label_config.gap

                        if config.label_config.position == "top":
                            label_top = final_main_top - label_gap - label_h
                        else:  # bottom
                            label_top = final_main_top + orig_h + label_gap

                        add_text_label(
                            slide,
                            label_text,
                            final_main_left,
                            label_top,
                            orig_w,
                            label_h,
                            config.label_config,
                        )

                # Advance X position
                if config.layout_mode == "flow":
                    current_x += item_w + this_gap_h
                else:
                    w = (
                        metrics.col_widths[col_idx]
                        if col_idx < len(metrics.col_widths)
                        else metrics.main_width
                    )
                    current_x += w + this_gap_h

            current_y += current_row_height + this_gap_v

        prs.save(config.output)
    finally:
        shutil.rmtree(temp_dir)

    return config.output


def _add_crop_images(
    slide,
    config: GridConfig,
    metrics: LayoutMetrics,
    image_path: str,
    temp_dir: str,
    row_idx: int,
    col_idx: int,
    crop_regions: List[CropRegion],
    final_main_left: float,
    final_main_top: float,
    orig_w: float,
    orig_h: float,
    global_gap_mc: float,
    global_gap_cc: float,
    border_offset_cm: float,
    half_border: float,
) -> None:
    """Add cropped images next to the main image."""
    disp = config.crop_display
    # 空クロップ（show_zoomed=False）はautoの位置計算から除外
    num_crops = sum(1 for r in crop_regions if r.show_zoomed)
    visible_crop_idx = 0

    for crop_idx, region in enumerate(crop_regions):
        # 空クロップの場合、拡大画像は描画しない
        if not region.show_zoomed:
            continue
        crop_filename = f"crop_{row_idx}_{col_idx}_{crop_idx}.png"
        crop_path = os.path.join(temp_dir, crop_filename)
        try:
            crop_image(image_path, region, crop_path)
        except Exception:
            continue

        # Calculate crop size
        cw, ch = _resolve_crop_size(
            config, metrics, disp, crop_path, orig_w, orig_h, num_crops, global_gap_cc
        )

        # Calculate gap
        this_gap_mc = region.gap if region.gap is not None else global_gap_mc
        if region.gap is not None and config.show_zoom_border:
            this_gap_mc += border_offset_cm

        # Calculate position (use visible_crop_idx for auto positioning)
        if disp.position == "right":
            c_left = final_main_left + orig_w + this_gap_mc
            c_top = _calculate_crop_absolute_vertical_position(
                region,
                visible_crop_idx,
                num_crops,
                final_main_top,
                orig_h,
                ch,
                global_gap_cc,
                half_border,
                disp,
            )
        else:  # bottom
            c_top = final_main_top + orig_h + this_gap_mc
            c_left = _calculate_crop_absolute_horizontal_position(
                region,
                visible_crop_idx,
                num_crops,
                final_main_left,
                orig_w,
                cw,
                global_gap_cc,
                half_border,
                disp,
            )

        pic_crop = slide.shapes.add_picture(
            crop_path,
            cm_to_emu(c_left),
            cm_to_emu(c_top),
            cm_to_emu(cw),
            cm_to_emu(ch),
        )
        pic_crop.shadow.inherit = False

        if config.show_zoom_border:
            add_border_shape(
                slide,
                c_left,
                c_top,
                cw,
                ch,
                region.color,
                config.zoom_border_width,
                config.zoom_border_shape,
            )

        # Draw connector line if enabled
        if config.connector.show:
            # Calculate crop border position on main image
            try:
                orig_width_px, orig_height_px = get_image_dimensions(image_path)
                x, y, w, h = resolve_crop_box(region, orig_width_px, orig_height_px)

                scale_x = orig_w / orig_width_px
                scale_y = orig_h / orig_height_px

                border_left = final_main_left + x * scale_x
                border_top = final_main_top + y * scale_y
                border_w = w * scale_x
                border_h = h * scale_y

                # Calculate connector points
                (start_x, start_y), (end_x, end_y) = calculate_connector_points(
                    border_left,
                    border_top,
                    border_w,
                    border_h,
                    c_left,
                    c_top,
                    cw,
                    ch,
                    disp.position,
                )

                # Draw connector line
                add_connector_line(
                    slide,
                    start_x,
                    start_y,
                    end_x,
                    end_y,
                    config.connector,
                    region.color,
                )
            except Exception:
                pass  # Skip connector if there's an error

        visible_crop_idx += 1


def _resolve_crop_size(
    config: GridConfig,
    metrics: LayoutMetrics,
    disp: CropDisplayConfig,
    crop_path: str,
    orig_w: float,
    orig_h: float,
    num_crops: int,
    global_gap_cc: float,
) -> Tuple[float, float]:
    """Resolve the size of a crop image."""
    if disp.size is not None:
        if disp.position == "right":
            return calculate_image_size_fit(crop_path, disp.size, 9999, "width")
        else:
            return calculate_image_size_fit(crop_path, 9999, disp.size, "height")
    elif disp.scale is not None:
        if disp.position == "right":
            return calculate_image_size_fit(
                crop_path, orig_w * disp.scale, 9999, "width"
            )
        else:
            return calculate_image_size_fit(
                crop_path, 9999, orig_h * disp.scale, "height"
            )
    else:
        if disp.position == "right":
            single_h = (orig_h - global_gap_cc * (num_crops - 1)) / num_crops
            return calculate_image_size_fit(
                crop_path, metrics.crop_size, single_h, "fit"
            )
        else:
            single_w = (orig_w - global_gap_cc * (num_crops - 1)) / num_crops
            return calculate_image_size_fit(
                crop_path, single_w, metrics.crop_size, "fit"
            )


def _calculate_crop_absolute_vertical_position(
    region: CropRegion,
    crop_idx: int,
    num_crops: int,
    final_main_top: float,
    orig_h: float,
    ch: float,
    global_gap_cc: float,
    half_border: float,
    disp: CropDisplayConfig,
) -> float:
    """Calculate absolute vertical position for a crop."""
    if region.align == "start":
        return final_main_top + region.offset + half_border
    elif region.align == "center":
        return final_main_top + (orig_h - ch) / 2 + region.offset
    elif region.align == "end":
        return final_main_top + orig_h - ch + region.offset - half_border
    else:  # auto
        # 最初と最後のクロップは端に揃える（pin ends）
        if crop_idx == 0:
            return final_main_top + half_border
        elif num_crops > 1 and crop_idx == num_crops - 1:
            return (final_main_top + orig_h) - ch - half_border
        else:
            # 中間のクロップは均等配置
            if disp.scale is not None or disp.size is not None:
                # scale/size指定時: 等間隔で配置
                total_crop_h = num_crops * ch + (num_crops - 1) * global_gap_cc
                start_y = final_main_top + (orig_h - total_crop_h) / 2
                return start_y + crop_idx * (ch + global_gap_cc)
            else:
                # fit時: スロットの中央に配置
                single_h = (orig_h - global_gap_cc * (num_crops - 1)) / num_crops
                slot_top = final_main_top + crop_idx * (single_h + global_gap_cc)
                return slot_top + (single_h - ch) / 2


def _calculate_crop_absolute_horizontal_position(
    region: CropRegion,
    crop_idx: int,
    num_crops: int,
    final_main_left: float,
    orig_w: float,
    cw: float,
    global_gap_cc: float,
    half_border: float,
    disp: CropDisplayConfig,
) -> float:
    """Calculate absolute horizontal position for a crop."""
    if region.align == "start":
        return final_main_left + region.offset + half_border
    elif region.align == "center":
        return final_main_left + (orig_w - cw) / 2 + region.offset
    elif region.align == "end":
        return final_main_left + orig_w - cw + region.offset - half_border
    else:  # auto
        # 最初と最後のクロップは端に揃える（pin ends）
        if crop_idx == 0:
            return final_main_left + half_border
        elif num_crops > 1 and crop_idx == num_crops - 1:
            return (final_main_left + orig_w) - cw - half_border
        else:
            # 中間のクロップは均等配置
            if disp.scale is not None or disp.size is not None:
                # scale/size指定時: 等間隔で配置
                total_crop_w = num_crops * cw + (num_crops - 1) * global_gap_cc
                start_x = final_main_left + (orig_w - total_crop_w) / 2
                return start_x + crop_idx * (cw + global_gap_cc)
            else:
                # fit時: スロットの中央に配置
                single_w = (orig_w - global_gap_cc * (num_crops - 1)) / num_crops
                slot_left = final_main_left + crop_idx * (single_w + global_gap_cc)
                return slot_left + (single_w - cw) / 2


# =============================================================================
# Text Label Functions
# =============================================================================


def generate_label_text(
    config: GridConfig,
    image_path: Optional[str],
    row_idx: int,
    col_idx: int,
) -> str:
    """Generate label text based on label mode."""
    label_config = config.label_config

    if label_config.mode == "filename":
        if image_path:
            return Path(image_path).stem
        return ""
    elif label_config.mode == "number":
        # Calculate sequential number (row-major order)
        n = row_idx * config.cols + col_idx + 1
        return label_config.number_format.format(n=n)
    elif label_config.mode == "custom":
        idx = row_idx * config.cols + col_idx
        if idx < len(label_config.custom_texts):
            return label_config.custom_texts[idx]
        return ""
    return ""


def add_text_label(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    label_config: LabelConfig,
) -> None:
    """Add a text label to the slide."""
    if not text:
        return

    # Create text box
    textbox = slide.shapes.add_textbox(
        cm_to_emu(left),
        cm_to_emu(top),
        cm_to_emu(width),
        cm_to_emu(height),
    )

    # Configure text frame
    tf = textbox.text_frame
    tf.word_wrap = False

    # Add text
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER

    # Configure font
    run = p.runs[0]
    run.font.name = label_config.font_name
    run.font.size = Pt(label_config.font_size)
    run.font.bold = label_config.font_bold
    run.font.color.rgb = RGBColor(*label_config.font_color)


def calculate_label_height(label_config: LabelConfig) -> float:
    """Calculate label height in cm based on font size."""
    # Approximate: font size (pt) * 1.2 / 72 inches * 2.54 cm
    return label_config.font_size * 1.5 / 72 * 2.54


# =============================================================================
# Connector Line Functions
# =============================================================================


def add_connector_line(
    slide,
    start_x: float,
    start_y: float,
    end_x: float,
    end_y: float,
    connector_config: ConnectorConfig,
    region_color: Tuple[int, int, int],
) -> None:
    """Add a connector line between crop border and zoomed image."""
    # Use region color if no specific color is set
    line_color = connector_config.color if connector_config.color else region_color

    # Create a straight connector using add_connector
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        cm_to_emu(start_x),
        cm_to_emu(start_y),
        cm_to_emu(end_x),
        cm_to_emu(end_y),
    )

    # Style the line
    connector.line.color.rgb = RGBColor(*line_color)
    connector.line.width = pt_to_emu(connector_config.width)

    # Set dash style
    if connector_config.dash_style == "dash":
        connector.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    elif connector_config.dash_style == "dot":
        connector.line.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT
    else:  # solid
        connector.line.dash_style = MSO_LINE_DASH_STYLE.SOLID


def calculate_connector_points(
    crop_border_left: float,
    crop_border_top: float,
    crop_border_width: float,
    crop_border_height: float,
    zoom_left: float,
    zoom_top: float,
    zoom_width: float,
    zoom_height: float,
    position: str,
) -> Tuple[Tuple[float, float], Tuple[float, float]]:
    """
    Calculate start and end points for connector line.

    Returns:
        Tuple of (start_point, end_point) where each point is (x, y)
    """
    if position == "right":
        # Connect from right edge of crop border to left edge of zoom image
        start_x = crop_border_left + crop_border_width
        start_y = crop_border_top + crop_border_height / 2
        end_x = zoom_left
        end_y = zoom_top + zoom_height / 2
    else:  # bottom
        # Connect from bottom edge of crop border to top edge of zoom image
        start_x = crop_border_left + crop_border_width / 2
        start_y = crop_border_top + crop_border_height
        end_x = zoom_left + zoom_width / 2
        end_y = zoom_top

    return (start_x, start_y), (end_x, end_y)


# =============================================================================
# Sample Configuration Generator
# =============================================================================

SAMPLE_CONFIG = """# PowerPoint Image Grid Generator - Configuration File
# Size values are in centimeters (cm) unless otherwise specified (e.g. pt)

# Slide size settings
slide:
  width: 33.867    # 16:9 aspect ratio
  height: 19.05

# Main grid layout
grid:
  rows: 2
  cols: 3
  arrangement: row  # 'row' or 'col'
  layout_mode: flow # 'grid' (aligned) or 'flow' (compact)
  flow_align: left # 'left', 'center', 'right' (only for flow)
  flow_vertical_align: center # 'top', 'center', 'bottom' (only for flow)

# Margins around the grid (cm)
margin:
  left: 1.0
  top: 1.0
  right: 1.0
  bottom: 1.0

# Gap between images in grid
gap:
  horizontal:
    value: 0.5
    mode: cm # 'cm' or 'scale'
  vertical:
    value: 0.5
    mode: cm

# Image sizing
image:
  size_mode: fit   # 'fit' (auto calc) or 'fixed' (user defined)
  fit_mode: fit    # 'fit', 'width', 'height'
  # width: 10.0    # used if size_mode is fixed
  # height: 7.5
  scale: 1.0

# Crop settings
crop:
  # Multiple crop regions
  regions:
    - name: "Region A"
      x: 50
      y: 50
      width: 100
      height: 100
      color: "#FF0000"
      align: auto   # auto, start, center, end
      offset: 0.0   # cm
      # gap: 0.5    # overrides main_crop_gap
    - name: "Region B"
      x: 200
      y: 100
      width: 120
      height: 80
      color: "#00FF00"
  
  # Which rows/cols to crop (0-indexed), null = all
  rows: [0]
  cols: null
  
  # How to display cropped images next to the original
  display:
    position: right   # 'right' or 'bottom'
    size: null        # Absolute size (cm). If set, takes priority over scale.
    scale: 0.4        # Display size ratio relative to main image
    main_crop_gap: {value: 0.15, mode: cm}
    crop_crop_gap: {value: 0.15, mode: cm}
    crop_bottom_gap: {value: 0.0, mode: cm}

# Border settings
border:
  crop:
    show: true
    width: 1.5  # Line width in points (pt)
    shape: rectangle # 'rectangle' or 'rounded'
  zoom:
    show: true
    width: 1.5  # Line width in points (pt)
    shape: rectangle # 'rectangle' or 'rounded'

# Input folders
folders:
  - "./images/row1"
  - "./images/row2"

# Output file
output: "output.pptx"
"""


def generate_sample_config(output_path: str) -> None:
    """Generate a sample configuration file."""
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(SAMPLE_CONFIG)


# =============================================================================
# Crop Presets
# =============================================================================

PRESETS_FILE = "crop_presets.yaml"


def get_default_presets() -> List[CropPreset]:
    """ビルトインプリセットを返す（ズーム用プリセット）"""
    return [
        CropPreset(
            name="左右ズーム",
            description="左右の2箇所を拡大表示",
            display_position="bottom",
            regions=[
                CropRegion(
                    x=0,
                    y=0,
                    width=100,
                    height=100,
                    mode="ratio",
                    x_ratio=0.05,
                    y_ratio=0.25,
                    width_ratio=0.25,
                    height_ratio=0.5,
                    color=(255, 0, 0),
                    name="左",
                ),
                CropRegion(
                    x=0,
                    y=0,
                    width=100,
                    height=100,
                    mode="ratio",
                    x_ratio=0.70,
                    y_ratio=0.25,
                    width_ratio=0.25,
                    height_ratio=0.5,
                    color=(0, 0, 255),
                    name="右",
                ),
            ],
        ),
        CropPreset(
            name="正方形2箇所ズーム（横）",
            description="左右に正方形2箇所を拡大表示",
            display_position="bottom",
            regions=[
                CropRegion(
                    x=0,
                    y=0,
                    width=100,
                    height=100,
                    mode="ratio",
                    x_ratio=0.05,
                    y_ratio=0.25,
                    width_ratio=0.25,
                    height_ratio=0.25,
                    color=(255, 0, 0),
                    name="左",
                ),
                CropRegion(
                    x=0,
                    y=0,
                    width=100,
                    height=100,
                    mode="ratio",
                    x_ratio=0.70,
                    y_ratio=0.25,
                    width_ratio=0.25,
                    height_ratio=0.25,
                    color=(0, 0, 255),
                    name="右",
                ),
            ],
        ),
        CropPreset(
            name="正方形2箇所ズーム（縦）",
            description="上下に正方形2箇所を拡大表示",
            display_position="right",
            regions=[
                CropRegion(
                    x=0,
                    y=0,
                    width=100,
                    height=100,
                    mode="ratio",
                    x_ratio=0.3,
                    y_ratio=0.05,
                    width_ratio=0.25,
                    height_ratio=0.25,
                    color=(255, 0, 0),
                    name="上",
                ),
                CropRegion(
                    x=0,
                    y=0,
                    width=100,
                    height=100,
                    mode="ratio",
                    x_ratio=0.3,
                    y_ratio=0.70,
                    width_ratio=0.25,
                    height_ratio=0.25,
                    color=(0, 0, 255),
                    name="下",
                ),
            ],
        ),
        CropPreset(
            name="上下ズーム",
            description="上下の2箇所を拡大表示",
            display_position="right",
            regions=[
                CropRegion(
                    x=0,
                    y=0,
                    width=100,
                    height=100,
                    mode="ratio",
                    x_ratio=0.25,
                    y_ratio=0.05,
                    width_ratio=0.5,
                    height_ratio=0.25,
                    color=(255, 0, 0),
                    name="上",
                ),
                CropRegion(
                    x=0,
                    y=0,
                    width=100,
                    height=100,
                    mode="ratio",
                    x_ratio=0.25,
                    y_ratio=0.70,
                    width_ratio=0.5,
                    height_ratio=0.25,
                    color=(0, 0, 255),
                    name="下",
                ),
            ],
        ),
        CropPreset(
            name="4隅ズーム",
            description="4隅を拡大表示",
            display_position="bottom",
            regions=[
                CropRegion(
                    x=0,
                    y=0,
                    width=100,
                    height=100,
                    mode="ratio",
                    x_ratio=0.0,
                    y_ratio=0.0,
                    width_ratio=0.2,
                    height_ratio=0.2,
                    color=(255, 0, 0),
                    name="左上",
                ),
                CropRegion(
                    x=0,
                    y=0,
                    width=100,
                    height=100,
                    mode="ratio",
                    x_ratio=0.8,
                    y_ratio=0.0,
                    width_ratio=0.2,
                    height_ratio=0.2,
                    color=(0, 255, 0),
                    name="右上",
                ),
                CropRegion(
                    x=0,
                    y=0,
                    width=100,
                    height=100,
                    mode="ratio",
                    x_ratio=0.0,
                    y_ratio=0.8,
                    width_ratio=0.2,
                    height_ratio=0.2,
                    color=(0, 0, 255),
                    name="左下",
                ),
                CropRegion(
                    x=0,
                    y=0,
                    width=100,
                    height=100,
                    mode="ratio",
                    x_ratio=0.8,
                    y_ratio=0.8,
                    width_ratio=0.2,
                    height_ratio=0.2,
                    color=(255, 255, 0),
                    name="右下",
                ),
            ],
        ),
        CropPreset(
            name="中央ズーム",
            description="中央部分を拡大表示",
            display_position="right",
            regions=[
                CropRegion(
                    x=0,
                    y=0,
                    width=100,
                    height=100,
                    mode="ratio",
                    x_ratio=0.35,
                    y_ratio=0.35,
                    width_ratio=0.3,
                    height_ratio=0.3,
                    color=(255, 0, 0),
                    name="中央",
                ),
            ],
        ),
        CropPreset(
            name="中央＋4隅ズーム",
            description="中央と4隅を拡大表示",
            display_position="bottom",
            regions=[
                CropRegion(
                    x=0,
                    y=0,
                    width=100,
                    height=100,
                    mode="ratio",
                    x_ratio=0.0,
                    y_ratio=0.0,
                    width_ratio=0.15,
                    height_ratio=0.15,
                    color=(255, 0, 0),
                    name="左上",
                ),
                CropRegion(
                    x=0,
                    y=0,
                    width=100,
                    height=100,
                    mode="ratio",
                    x_ratio=0.425,
                    y_ratio=0.425,
                    width_ratio=0.15,
                    height_ratio=0.15,
                    color=(0, 255, 0),
                    name="中央",
                ),
                CropRegion(
                    x=0,
                    y=0,
                    width=100,
                    height=100,
                    mode="ratio",
                    x_ratio=0.85,
                    y_ratio=0.85,
                    width_ratio=0.15,
                    height_ratio=0.15,
                    color=(0, 0, 255),
                    name="右下",
                ),
            ],
        ),
        CropPreset(
            name="横3箇所ズーム",
            description="横に3箇所を拡大表示",
            display_position="bottom",
            regions=[
                CropRegion(
                    x=0,
                    y=0,
                    width=100,
                    height=100,
                    mode="ratio",
                    x_ratio=0.05,
                    y_ratio=0.3,
                    width_ratio=0.2,
                    height_ratio=0.4,
                    color=(255, 0, 0),
                    name="左",
                ),
                CropRegion(
                    x=0,
                    y=0,
                    width=100,
                    height=100,
                    mode="ratio",
                    x_ratio=0.4,
                    y_ratio=0.3,
                    width_ratio=0.2,
                    height_ratio=0.4,
                    color=(0, 255, 0),
                    name="中",
                ),
                CropRegion(
                    x=0,
                    y=0,
                    width=100,
                    height=100,
                    mode="ratio",
                    x_ratio=0.75,
                    y_ratio=0.3,
                    width_ratio=0.2,
                    height_ratio=0.4,
                    color=(0, 0, 255),
                    name="右",
                ),
            ],
        ),
    ]


def load_crop_presets(filepath: Optional[str] = None) -> List[CropPreset]:
    """プリセットをファイルから読み込む（デフォルトプリセット + ユーザー定義）"""
    presets = get_default_presets()

    if filepath is None:
        filepath = PRESETS_FILE

    if Path(filepath).exists():
        try:
            with open(filepath, "r", encoding="utf-8") as f:
                data = yaml.safe_load(f)
            if data and "presets" in data:
                for p in data["presets"]:
                    regions = []
                    for i, r in enumerate(p.get("regions", [])):
                        regions.append(_parse_crop_region_dict(r, f"preset_{i}"))
                    presets.append(
                        CropPreset(
                            name=p.get("name", "Unknown"),
                            regions=regions,
                            description=p.get("description", ""),
                            display_position=p.get("display_position", "right"),
                        )
                    )
        except Exception:
            pass  # ファイル読み込みエラーは無視

    return presets


def save_crop_preset(preset: CropPreset, filepath: Optional[str] = None) -> None:
    """プリセットをファイルに保存"""
    if filepath is None:
        filepath = PRESETS_FILE

    # 既存のプリセットを読み込む
    presets_data: List[dict] = []
    if Path(filepath).exists():
        try:
            with open(filepath, "r", encoding="utf-8") as f:
                data = yaml.safe_load(f) or {}
            presets_data = data.get("presets", [])
        except Exception:
            pass

    # 同名のプリセットがあれば上書き
    presets_data = [p for p in presets_data if p.get("name") != preset.name]

    # 新しいプリセットを追加
    presets_data.append(
        {
            "name": preset.name,
            "description": preset.description,
            "display_position": preset.display_position,
            "regions": [
                {
                    "name": r.name,
                    "x": r.x,
                    "y": r.y,
                    "width": r.width,
                    "height": r.height,
                    "color": list(r.color),
                    "align": r.align,
                    "offset": r.offset,
                    "gap": r.gap,
                    "mode": r.mode,
                    "x_ratio": r.x_ratio,
                    "y_ratio": r.y_ratio,
                    "width_ratio": r.width_ratio,
                    "height_ratio": r.height_ratio,
                }
                for r in preset.regions
            ],
        }
    )

    with open(filepath, "w", encoding="utf-8") as f:
        yaml.dump({"presets": presets_data}, f, allow_unicode=True)
